import logging
logging.getLogger('streamlit.watcher.local_sources_watcher').disabled = True
logging.getLogger('streamlit.web.bootstrap').disabled = True

import asyncio
import base64 # Aggiunto per persistenza file_bytes
class _DummyLoop:
    def is_running(self):
        return False
_orig_get_running_loop = asyncio.get_running_loop
def _safe_get_running_loop():
    try:
        return _orig_get_running_loop()
    except RuntimeError:
        return _DummyLoop()
asyncio.get_running_loop = _safe_get_running_loop

import socketserver  # Permette SO_REUSEADDR
socketserver.TCPServer.allow_reuse_address = True
import wsgiref.simple_server  # Server OAuth2
wsgiref.simple_server.WSGIServer.allow_reuse_address = True
import os
os.environ["STREAMLIT_SERVER_FILEWATCHERTYPE"] = "none"
import streamlit as st
import os
# import openai # Rimosso
import google.generativeai as genai # Aggiunto per Gemini
from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import uuid
import json
from datetime import datetime, timezone
# import tiktoken # Rimosso
import fitz # PyMuPDF
from googleapiclient.discovery import build # Per Google Search API
import re # Aggiunto per parsing citazioni
import tempfile # Per file temporanei
# ---- Configurazione OAuth2 per Google Docs API ----
from google_auth_oauthlib.flow import InstalledAppFlow
from google.oauth2.credentials import Credentials
import io
from googleapiclient.http import MediaIoBaseUpload, MediaIoBaseDownload
from urllib.request import urlopen, Request # Aggiunto
from urllib.error import URLError, HTTPError # Aggiunto
from urllib.parse import urlparse # Aggiunto
import mimetypes # Aggiunto
import os # Già presente, ma lo useremo per os.path.basename
import xml.etree.ElementTree as ET # Aggiunto per PubMed
import urllib.parse # Aggiunto per quote nelle URL API
import requests # Per scaricare contenuto web
from bs4 import BeautifulSoup # Per parsare HTML ed estrarre testo
from fuzzywuzzy import fuzz # AGGIUNTO PER VERIFICA TITOLO FUZZY

# Import per nuovi formati
import docx
import openpyxl
from pptx import Presentation

# ---- Aggiunta import per LangChain Text Splitter ----
from langchain.text_splitter import RecursiveCharacterTextSplitter

# ---- Integrazione Ricerca Web Gratuita (DuckDuckGo) ----
# try:
#     from duckduckgo_search import DDGS
#     DUCKDUCKGO_SEARCH_AVAILABLE = True
# except ImportError:
#     DUCKDUCKGO_SEARCH_AVAILABLE = False
# # ---- Fine Integrazione Ricerca Web ----

# Funzione per iniettare lo script JS per la copia, una volta per sessione
def ensure_js_copy_script_injected():
    if "js_copy_script_injected" not in st.session_state:
        st.markdown("""
        <script>
        function copyMessageToClipboard(textToCopy, buttonId) {
            navigator.clipboard.writeText(textToCopy).then(function() {
                const button = document.getElementById(buttonId);
                if (button) {
                    const originalHTML = button.innerHTML;
                    button.innerHTML = 'Copiato! ✅';
                    button.style.backgroundColor = '#d4edda'; 
                    button.style.borderColor = '#c3e6cb';
                    button.disabled = true;
                    setTimeout(function() {
                        button.innerHTML = originalHTML;
                        button.style.backgroundColor = '';
                        button.style.borderColor = '';
                        button.disabled = false;
                    }, 2000);
                }
            }).catch(function(err) {
                console.error('Errore durante la copiatura negli appunti: ', err);
                const button = document.getElementById(buttonId);
                if (button) {
                    const originalHTML = button.innerHTML;
                    button.innerHTML = 'Errore ❌';
                    button.style.backgroundColor = '#f8d7da';
                    button.style.borderColor = '#f5c6cb';
                    button.disabled = true;
                    setTimeout(function() {
                        button.innerHTML = originalHTML;
                        button.style.backgroundColor = '';
                        button.style.borderColor = '';
                        button.disabled = false;
                    }, 2000);
                }
            });
        }
        </script>
        """, unsafe_allow_html=True)
        st.session_state.js_copy_script_injected = True

# ----- 0. Costanti e Configurazioni Globali -----
DATA_DIR = "tutor_ai_data"
CHATS_FILE = os.path.join(DATA_DIR, "chat_sessions.json")
APP_STATE_FILE = os.path.join(DATA_DIR, "app_state.json")

os.makedirs(DATA_DIR, exist_ok=True)

st.set_page_config(page_title="💬 Super Tutor AI Chat (Gemini Ed.)", layout="wide", initial_sidebar_state="expanded")

GEMINI_MODELS_INFO = {
    "gemini-1.5-flash-latest": "Veloce, multimodale, efficiente per compiti ad alto volume.",
    "gemini-1.5-pro-latest": "Modello multimodale più capace, per compiti complessi.",
    "gemini-pro": "Buon bilanciamento performance/costo per testo (legacy, preferire 1.5).",
    "gemini-2.5-flash-preview-04-17" : "Modello multimodale avanzato, per compiti complessi.",
    "gemini-2.0-flash" : "Modello multimodale avanzato, per compiti complessi.",
    "gemini-2.5-pro-exp-03-25" : "Modello multimodale più avanzato, per compiti complessi.",
}
DEFAULT_GEMINI_MODEL = "gemini-1.5-flash-latest"
DEFAULT_SIMILARITY_THRESHOLD = 0.15 # Abbassato per essere più permissivo
DEFAULT_TOP_N = 5

# ---- Configurazione OAuth2 per Google Docs API ----
CLIENT_CONFIG = {
    "installed": {
        "client_id": "76084660967-nhc46qtovlpt4n62q4eh4212gs7c52ap.apps.googleusercontent.com",
        "client_secret": "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",
        "auth_uri": "https://accounts.google.com/o/oauth2/auth",
        "token_uri": "https://oauth2.googleapis.com/token",
        "redirect_uris": ["http://localhost:8502/"]
    }
}
SCOPES = [
    "https://www.googleapis.com/auth/documents",
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/calendar",
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/cloud-platform"
]

# ----- 1. Funzioni di Caching e Utilità -----
@st.cache_resource
def load_embedding_model_cached():
    try:
        model = SentenceTransformer('all-MiniLM-L6-v2')
        return model
    except Exception as e:
        st.error(f"Errore caricamento modello embedding: {e}")
        return None

@st.cache_data
def get_chunks_and_embeddings_cached(_uploaded_file_contents_hash, raw_file_bytes, mime_type, chunk_size, overlap):
    # raw_file_bytes sono i byte originali del file
    # mime_type è il tipo MIME del file (es. "text/plain", "application/pdf", "image/png")
    
    if not raw_file_bytes:
        st.warning("File vuoto fornito.")
        return [], np.array([])

    text_content = ""
    is_text_type_for_rag = False
    current_file_name = st.session_state.get('uploaded_file_name', 'N/A') # Per i messaggi di log
    current_file_name_lower = current_file_name.lower()

    # Gestione tipi testuali e codice sorgente
    # Nota: st.session_state.uploaded_file_name potrebbe non essere ancora aggiornato qui se è un nuovo file.
    # Usiamo il nome del file passato implicitamente tramite la sidebar o quello in session_state.
    # Per l'euristica basata su estensione, è meglio fare affidamento sul nome file in session_state se disponibile e coerente.
    
    # Lista di estensioni per file di testo/codice
    text_code_extensions = [
        '.md', '.csv', '.json', '.xml', '.html', '.htm', '.rtf', # Testo/Markup
        '.py', '.js', '.java', '.c', '.cpp', '.cs', '.go', '.rb', 
        '.php', '.swift', '.kt', '.ts', '.sql', '.sh', '.bat', '.ps1' # Codice e Script
    ]
    
    office_mime_types = {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }

    if mime_type.startswith("text/") or \
       mime_type in ["application/json", "application/xml", "application/csv", 
                      "application/rtf", "application/xhtml+xml", "application/javascript",
                      "application/x-python-code", "application/x-csharp", # Esempi, non tutti i server li usano
                      # Aggiungere altri tipi MIME specifici per codice se noti
                     ] or \
       (current_file_name != 'N/A' and any(current_file_name_lower.endswith(ext) for ext in text_code_extensions)):
        is_text_type_for_rag = True
        try:
            text_content = raw_file_bytes.decode('utf-8', errors='replace')
        except Exception as e:
            st.error(f"Errore decodifica file testuale ({mime_type}, nome: {current_file_name}): {e}")
            return None, None
    
    # Gestione DOCX
    elif mime_type == office_mime_types["docx"] or current_file_name_lower.endswith(".docx"):
        is_text_type_for_rag = True
        try:
            document = docx.Document(io.BytesIO(raw_file_bytes))
            text_content = "\\n".join([para.text for para in document.paragraphs])
            if not text_content.strip():
                st.info(f"Il file DOCX '{current_file_name}' sembra non contenere testo estraibile dai paragrafi.")
        except Exception as e:
            st.error(f"Errore estrazione testo da DOCX '{current_file_name}': {e}")
            return None, None

    # Gestione XLSX
    elif mime_type == office_mime_types["xlsx"] or current_file_name_lower.endswith(".xlsx"):
        is_text_type_for_rag = True
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(raw_file_bytes))
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    if row_text: # Solo se la riga contiene testo
                        text_parts.append(" ".join(row_text)) # Unisci le celle di una riga con spazio
            text_content = "\\n".join(text_parts) # Unisci le righe con newline
            if not text_content.strip():
                st.info(f"Il file XLSX '{current_file_name}' sembra non contenere testo estraibile dalle celle.")
        except Exception as e:
            st.error(f"Errore estrazione testo da XLSX '{current_file_name}': {e}")
            return None, None
            
    # Gestione PPTX
    elif mime_type == office_mime_types["pptx"] or current_file_name_lower.endswith(".pptx"):
        is_text_type_for_rag = True
        try:
            presentation = Presentation(io.BytesIO(raw_file_bytes))
            text_parts = []
            for slide in presentation.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_texts.append(run.text)
                if slide_texts: # Solo se la slide contiene testo
                     text_parts.append(" ".join(slide_texts)) # Unisci testi da una slide
            text_content = "\\n\\n".join(text_parts) # Unisci testi da slide diverse con doppio newline
            if not text_content.strip():
                st.info(f"Il file PPTX '{current_file_name}' sembra non contenere testo estraibile dalle slide.")
        except Exception as e:
            st.error(f"Errore estrazione testo da PPTX '{current_file_name}': {e}")
            return None, None

    # Gestione PDF
    elif mime_type == "application/pdf" or current_file_name_lower.endswith(".pdf"): # aggiunto controllo estensione per sicurezza
        is_text_type_for_rag = True
        try:
            with fitz.open(stream=raw_file_bytes, filetype="pdf") as doc:
                text_content = "".join([page.get_text() + "\\n\\n" for page in doc])
            if not text_content.strip() and any(page.get_images(full=True) for page in fitz.open(stream=raw_file_bytes, filetype="pdf") if page.get_images(full=True)):
                 st.info(f"Il PDF '{current_file_name}' sembra contenere principalmente immagini o nessun testo estraibile. L'analisi del testo dalle immagini (OCR) non è attualmente attiva.")
        except Exception as e:
            st.error(f"Errore estrazione testo PDF '{current_file_name}': {e}")
            return None, None
    
    # Gestione Immagini (non per RAG testuale, ma documento è "processato")
    elif mime_type.startswith("image/"):
        st.success(f"Immagine '{current_file_name}' ({mime_type}) caricata. "
                 "Potrai fare domande dirette sull'immagine all'AI. "
                 "Non verrà usata per la ricerca contestuale basata su testo (RAG).")
        # Documento processato (caricato), ma no chunk testuali per RAG.
        # uploaded_file_bytes e uploaded_file_mime sono già in session_state.
        return [], np.array([]) 

    else:
        st.error(f"Tipo file non supportato per l'estrazione testuale automatica: {mime_type} (nome: {current_file_name}). "
                 "Se è un file di testo/codice con un MIME type inatteso, prova a rinominarlo con un'estensione standard (es. .txt, .py).")
        return None, None

    # Se abbiamo estratto testo, procedi con chunking e embedding
    if is_text_type_for_rag:
        if not text_content.strip():
            st.warning(f"Documento '{current_file_name}' ({mime_type}) risulta vuoto o non è stato possibile estrarre contenuto testuale significativo per il RAG.")
            return [], np.array([])
        
        # ---- Inizio Modifiche per Pulizia Testo e Normalizzazione Newline ----
        # 1. Normalizzazione dei newline:
        #    Sostituisce la sequenza letterale di due caratteri \ seguita da n (cioè "\\n")
        #    con un singolo carattere newline (\n).
        #    Questo è utile perché alcune routine di estrazione potrebbero inserire "\\n" come stringa.
        text_content = text_content.replace('\\\\n', '\n')

        # 2. Rimozione di spazi bianchi eccessivi e normalizzazione degli spazi:
        #    Sostituisce una o più occorrenze di caratteri di spaziatura (incluso il newline singolo,
        #    tabulazioni, spazi multipli) con un singolo spazio. 
        #    Infine, .strip() rimuove eventuali spazi bianchi iniziali o finali dal testo risultante.
        text_content = re.sub(r'\s+', ' ', text_content).strip()
        
        # Ricontrolla se text_content è diventato vuoto dopo la pulizia
        if not text_content:
            st.warning(f"Documento '{current_file_name}' ({mime_type}) risulta vuoto dopo la pulizia del testo.")
            return [], np.array([])
        # ---- Fine Modifiche per Pulizia Testo ----

        # ---- Inizio Modifiche per Chunking Avanzato con LangChain ----
        # La vecchia logica basata su words = text_content.split() e il loop manuale viene sostituita.
        
        # Configura lo splitter (es. ricorsivo basato su caratteri)
        # NOTA: chunk_size e overlap ora si riferiscono al conteggio di caratteri (definito da length_function=len),
        # non più al conteggio delle parole. Potrebbe essere necessario aggiustare i valori di default.
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size, 
            chunk_overlap=overlap,
            length_function=len, 
            add_start_index=False # Puoi impostarlo a True se vuoi gli indici per le citazioni
        )
        chunks = text_splitter.split_text(text_content)

        if not chunks:
            st.warning(f"Nessun chunk creato dal documento '{current_file_name}' ({mime_type}) dopo il tentativo di splitting avanzato.")
            return [], np.array([])
        # ---- Fine Modifiche per Chunking Avanzato con LangChain ----
        
        embedding_model = load_embedding_model_cached()
        if embedding_model is None: 
            st.error("Modello di embedding non caricato, impossibile processare il testo del documento.")
            return None, None 
        
        try:
            embeddings = embedding_model.encode(chunks)
        except Exception as e:
            st.error(f"Errore durante la creazione degli embeddings per '{current_file_name}': {e}")
            return None, None
        return chunks, embeddings
    
    # Fallback se non è is_text_type_for_rag e non è un'immagine (già gestito sopra, ma per sicurezza)
    st.error(f"Logica di processamento file incompleta per {mime_type} ({current_file_name}).")
    return None, None

def find_relevant_chunks(query_embedding, text_embeddings, text_chunks, top_n, similarity_threshold):
    if text_embeddings is None or text_embeddings.size == 0 or not text_chunks:
        return []
    similarities = cosine_similarity(query_embedding, text_embeddings)[0]
    # Ottieni gli indici dei chunk ordinati per similarità (dal più alto al più basso)
    sorted_indices = np.argsort(similarities)[::-1]
    
    relevant_chunks_data = []
    # Itera attraverso i chunk ordinati
    for i in sorted_indices:
        # Se la similarità è sopra la soglia
        if similarities[i] >= similarity_threshold:
            if i < len(text_chunks): # Controllo di sicurezza per l'indice
                 relevant_chunks_data.append({
                    "chunk": text_chunks[i], 
                    "similarity": float(similarities[i]), 
                    "original_index": int(i) # Indice originale del chunk nella lista completa
                })
        # Interrompi se abbiamo abbastanza chunk (top_n) O se la similarità scende sotto la soglia
        # Questo è importante: non prendere più di top_n chunk *anche se superano la soglia*
        if len(relevant_chunks_data) >= top_n:
            break
            
    # Se dopo aver considerato tutti i chunk sopra soglia, ne abbiamo ancora meno di top_n,
    # potremmo decidere di prendere i top_n assoluti, anche se alcuni sono sotto soglia (ma questo cambierebbe la logica)
    # Per ora, la logica è: prendi fino a top_n SE sono sopra soglia.
    return relevant_chunks_data


def estimate_gemini_tokens(text, model_name, api_key):
    if not api_key: return len(text.split())
    try:
        current_configured_key = genai.api_key if hasattr(genai, 'api_key') else None
        genai.configure(api_key=api_key)
        model_to_use = genai.GenerativeModel(model_name)
        count = model_to_use.count_tokens(text)
        if current_configured_key: genai.configure(api_key=current_configured_key)
        elif genai.api_key == api_key: genai.configure(api_key=None)
        return count.total_tokens
    except Exception: return len(text.split())

def perform_google_custom_search(query: str, num_results: int = 3) -> list[dict[str, str]] | str:
    """Esegue una ricerca Google Custom Search e restituisce i risultati formattati o un messaggio di errore."""
    api_key = st.session_state.get("google_api_key")
    cse_id = st.session_state.get("google_cse_id")

    if not api_key or not cse_id:
        st.warning("Chiave API Google (Custom Search) o ID Motore di Ricerca non configurati in st.session_state.")
        return "Configurazione API Google Custom Search mancante per eseguire la ricerca."
    
    raw_search_items = []
    try:
        service = build("customsearch", "v1", developerKey=api_key)
        res = service.cse().list(q=query, cx=cse_id, num=num_results).execute()
        raw_search_items = res.get('items', [])
    except Exception as e:
        st.error(f"Errore durante la chiamata API Google Custom Search: {e}")
        return f"Errore API Google Custom Search: {str(e)}"

    if not raw_search_items:
        return "Nessun risultato trovato dalla ricerca Google iniziale."

    # Ora, per il primo risultato (o il più rilevante), tentiamo di scaricare il contenuto completo
    # Inizializziamo una lista per i risultati che verranno formattati per l'LLM
    # Anche se scarichiamo il contenuto completo solo per uno, manteniamo la struttura a lista
    # per coerenza e per poter mostrare snippet degli altri.
    
    # Stringa per contenere il testo dettagliato dalla pagina web (se scaricato)
    detailed_content_from_page = ""
    # Lista di risultati formattati (titolo, snippet, url) per tutti gli item di Google Search
    formatted_google_results_summary = [] 

    for i, item in enumerate(raw_search_items):
        title = item.get("title", "N/A")
        snippet = item.get("snippet", "N/A")
        url = item.get("link", "#")
        formatted_google_results_summary.append({
            "title": title,
            "snippet": snippet,
            "url": url
        })

        # Tentiamo di scaricare il contenuto completo solo per il PRIMO risultato
        # e solo se l'URL sembra valido (non è un placeholder '#')
        if i == 0 and url and url != "#":
            st.write(f"DEBUG: Tentativo di fetch del contenuto da: {url}") # Log per debug
            extracted_text_from_url = fetch_and_extract_text(url)
            if extracted_text_from_url:
                is_error_message = extracted_text_from_url.startswith("Errore") or \
                                   extracted_text_from_url.startswith("Timeout")
                
                if not is_error_message and extracted_text_from_url.strip(): # Aggiunto .strip() per controllare che non sia solo whitespace
                    page_content_successfully_extracted = True
                    MAX_PAGE_TEXT_LEN = 7000 
                    if len(extracted_text_from_url) > MAX_PAGE_TEXT_LEN:
                        extracted_text_from_url = extracted_text_from_url[:MAX_PAGE_TEXT_LEN] + "... (contenuto pagina troncato)"
                    detailed_content_from_page = (f"Contenuto testuale estratto dalla pagina del primo risultato di ricerca ({url}):\\n" 
                                                f"---------------------------------------\\n"
                                                f"{extracted_text_from_url}\\n"
                                                f"---------------------------------------\\n")
                else: # L'estrazione ha restituito un messaggio di errore o testo vuoto
                    error_msg_for_user_display = extracted_text_from_url if is_error_message else 'Testo vuoto estratto'
                    if len(error_msg_for_user_display) > 250: # Tronca messaggi di errore lunghi per st.warning
                        error_msg_for_user_display = error_msg_for_user_display[:250] + "..."
                    st.warning(f"Tentativo di leggere il contenuto completo dal primo risultato ({url}) ha prodotto un errore o testo vuoto: {error_msg_for_user_display}")
                    # page_content_successfully_extracted rimane False. Non impostiamo detailed_content_from_page qui
                    # perché verrà gestito dalla logica successiva basata su page_content_successfully_extracted.
                    detailed_content_from_page = "" # Assicurati che sia vuoto se l'estrazione fallisce o produce errore
            
            # Se extracted_text_from_url era None o vuoto inizialmente.
            if not page_content_successfully_extracted and not detailed_content_from_page: # detailed_content_from_page potrebbe già essere stato gestito sopra
                st.warning(f"Non è stato possibile estrarre contenuto testuale significativo dal primo risultato di ricerca ({url}).")
                detailed_content_from_page = "" # Assicurati che sia vuoto
    
    # Ora costruiamo la stringa finale da restituire.
    # Questa stringa verrà passata all'LLM.
    final_output_parts = []

    # Includi il contenuto dettagliato SOLO se estratto con successo.
    if page_content_successfully_extracted and detailed_content_from_page.strip():
        final_output_parts.append(detailed_content_from_page)
        final_output_parts.append("\\nDi seguito un riepilogo degli altri risultati della ricerca Google (o di tutti se il primo non è stato processato in dettaglio):\\n")
    else:
        # Se l'estrazione non è riuscita o non c'era testo, l'LLM riceverà solo i riassunti.
        # Aggiungiamo una nota più neutra e meno focalizzata sul "fallimento".
        final_output_parts.append("Riepilogo dei risultati dalla ricerca Google (il contenuto completo della prima pagina potrebbe non essere stato accessibile o pienamente analizzabile, fare riferimento agli snippet forniti):\\n")

    for res_summary in formatted_google_results_summary:
        final_output_parts.append(f"Titolo: {res_summary['title']}")
        final_output_parts.append(f"Snippet: {res_summary['snippet']}")
        final_output_parts.append(f"URL: {res_summary['url']}")
        final_output_parts.append("---")
    
    return "\n".join(final_output_parts)

def fetch_and_extract_text(url: str, timeout: int = 10) -> str | None:
    """Scarica il contenuto di un URL e tenta di estrarre il testo principale."""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=timeout, allow_redirects=True, verify=False) # Aggiunto verify=False
        response.raise_for_status() # Solleva un errore per status code HTTP 4xx/5xx
        
        # Decodifica il contenuto correttamente se possibile
        response.encoding = response.apparent_encoding if response.apparent_encoding else 'utf-8'
        html_content = response.text

        soup = BeautifulSoup(html_content, 'html.parser')
        
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        
        text_parts = []
        
        # Fase 1: Trova tag <article> e <main>
        elements = soup.find_all(['article', 'main'])
        
        # Fase 2: Trova tag con attributo role="main"
        elements_with_role_main = soup.find_all(attrs={'role': 'main'})
        
        # Combina le liste di elementi, rimuovendo duplicati mantenendo l'ordine
        combined_elements = list(dict.fromkeys(elements + elements_with_role_main))

        if combined_elements:
            for element in combined_elements:
                text_parts.append(element.get_text(separator='\\n', strip=True))
        
        # Fallback solo se combined_elements non ha prodotto nulla
        if not text_parts: 
            body_text = soup.body.get_text(separator='\\n', strip=True) if soup.body else ""
            if body_text:
                 text_parts.append(body_text)

        # Fallback aggiuntivo se ANCORA non ci sono parti di testo significative
        # Questo usa stripped_strings per ottenere testo da posti più disparati
        if not any(part.strip() for part in text_parts): # Controlla se text_parts contiene testo effettivo
             temp_visible_texts = []
             # Usare stripped_strings è generalmente buono per ottenere solo testo visibile.
             all_strings = soup.stripped_strings 
             for t in all_strings:
                 # Semplifichiamo la condizione del parent: se è una stringa navigabile e non vuota, la prendiamo.
                 # L'LLM è abbastanza bravo a filtrare il rumore residuo se il contesto è sufficientemente ampio.
                 # La pulizia eccessiva qui potrebbe rimuovere testo utile.
                 if t.strip(): # Aggiungiamo qualsiasi stringa non vuota dopo lo strip
                     temp_visible_texts.append(t.strip())
             
             if temp_visible_texts:
                 # Se abbiamo trovato qualcosa con stripped_strings e text_parts era vuoto, usiamo questi.
                 # Non estendiamo se text_parts era vuoto, ma lo sostituiamo.
                 text_parts = temp_visible_texts 

        # Unisci le parti di testo. Usare \\n\\n per separare blocchi di testo distinti.
        # filter(None, text_parts) rimuove eventuali stringhe vuote prima di join.
        extracted_text = "\\n\\n".join(filter(None, text_parts))
        
        # Pulizia finale: rimuovi linee vuote multiple risultanti dal join o dal contenuto originale
        cleaned_text = "\\n".join([line for line in extracted_text.splitlines() if line.strip()])
        
        return cleaned_text if cleaned_text.strip() else None # Assicurati di restituire None se cleaned_text è vuoto o solo whitespace

    except requests.exceptions.Timeout:
        st.warning(f"Timeout durante il tentativo di scaricare {url}")
        return f"Timeout scaricando la pagina: {url}"
    except requests.exceptions.HTTPError as e:
        st.warning(f"Errore HTTP {e.response.status_code} scaricando {url}")
        return f"Errore HTTP {e.response.status_code} scaricando la pagina: {url}"
    except requests.exceptions.RequestException as e:
        st.warning(f"Errore durante il download di {url}: {e}")
        return f"Errore scaricando la pagina: {url} ({e})"
    except Exception as e:
        st.warning(f"Errore imprevisto durante l'estrazione del testo da {url}: {e}")
        return f"Errore processando la pagina: {url}"


# Definiamo il nuovo strumento per Google Search (NON PIU' USATO DIRETTAMENTE DA GEMINI)
# google_search_tool = genai.protos.Tool(
#     function_declarations=[
#         genai.protos.FunctionDeclaration(
#             name="perform_google_custom_search", 
#             description="Esegue una ricerca Google per trovare informazioni aggiornate...",
#             parameters=genai.protos.Schema(
#                 type=genai.protos.Type.OBJECT,
#                 properties={
#                     "query": genai.protos.Schema(type=genai.protos.Type.STRING, description="La query di ricerca da eseguire."),
#                     "num_results": genai.protos.Schema(type=genai.protos.Type.INTEGER, description="Numero di risultati da restituire (default 3, max 5 per concisione).")
#                 },
#                 required=["query"]
#             )
#         )
#     ]
# )

def get_answer_from_llm(
    query, 
    relevant_chunks_data, 
    api_key, 
    model_name=DEFAULT_GEMINI_MODEL, 
    chat_history=None, 
    ignore_active_contexts_for_this_query=False,
    direct_file_text_context: str | None = None, 
    image_data: bytes | None = None,             
    image_mime_type: str | None = None,
    external_web_search_results: str | None = None,
    scraped_page_content_for_query: str | None = None,
    temporary_instructions_for_this_turn: str | None = None,
    retrieved_history_context: list | None = None # NUOVO PARAMETRO PER RAG SU STORICO CHAT
):
    if not api_key: return "Errore: Chiave API Gemini non fornita.", 0, 0, {}
    try: genai.configure(api_key=api_key)
    except Exception as e: st.error(f"Errore configurazione API Gemini: {e}"); return f"Spiacente, errore configurazione AI: {e}", 0, 0, {}

    # --- Rilevamento Modalità Programmatore --- 
    is_programming_file_active = False
    active_file_language_for_markdown = ""
    # Definizioni per il rilevamento di file di codice
    # NOTA: Queste liste potrebbero essere definite globalmente se necessario altrove
    CODE_FILE_EXTENSIONS = [
        '.py', '.js', '.html', '.css', '.java', '.c', '.cpp', '.cs', '.go', '.rb', 
        '.php', '.swift', '.kt', '.ts', '.sql', '.sh', '.bat', '.ps1', 
        '.json', '.xml', '.md', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.tf', '.h'
    ]
    CODE_MIME_PREFIXES = ["text/"] # Mantenuto generico, ma l'estensione è più affidabile
    CODE_SPECIFIC_MIMES = [
        "application/javascript", "application/x-python-code", "application/x-csharp",
        "application/json", "application/xml", "application/sql", "application/x-sh",
        # Aggiungere altri tipi MIME specifici per codice se noti e affidabili
    ]
    LANGUAGE_EXTENSION_MAP = {
        '.py': 'python', '.js': 'javascript', '.html': 'html', '.css': 'css', '.java': 'java',
        '.c': 'c', '.cpp': 'cpp', '.cs': 'csharp', '.go': 'go', '.rb': 'ruby',
        '.php': 'php', '.swift': 'swift', '.kt': 'kotlin', '.ts': 'typescript',
        '.sql': 'sql', '.sh': 'bash', '.bat': 'batch', '.ps1': 'powershell',
        '.json': 'json', '.xml': 'xml', '.md': 'markdown', '.yaml': 'yaml', '.yml': 'yaml',
        '.toml': 'toml', '.h': 'c' # .h può essere C o C++ header, C è un default ragionevole
    }
 
    if not ignore_active_contexts_for_this_query and st.session_state.get("document_processed", False):
        uploaded_mime = st.session_state.get("uploaded_file_mime")
        uploaded_file_name = st.session_state.get("uploaded_file_name", "").lower()
        file_extension = os.path.splitext(uploaded_file_name)[1]
 
        if file_extension in CODE_FILE_EXTENSIONS: 
            is_programming_file_active = True
        elif uploaded_mime:
            if any(uploaded_mime.startswith(p) for p in CODE_MIME_PREFIXES) or \
               uploaded_mime in CODE_SPECIFIC_MIMES:
                # Se il MIME indica testo o un tipo di codice noto, ma l'estensione non era nella lista primaria,
                # consideralo un file di programmazione (es. un file .txt che contiene codice)
                is_programming_file_active = True 
 
        if is_programming_file_active:
            active_file_language_for_markdown = LANGUAGE_EXTENSION_MAP.get(file_extension, "")
            # Se active_file_language_for_markdown è vuoto (estensione non mappata), 
            # l'LLM potrebbe provare a indovinare o si può lasciare vuoto nel blocco Markdown.
 
    # --- 1. Istruzioni di Sistema --- (Ora dinamiche)
    system_instruction_parts = []
    if is_programming_file_active:
        # --- PROMPT PER MODALITÀ PROGRAMMATORE --- (Basato sulle linee guida dell'utente)
        system_instruction_parts = [
            "## Modalità Programmatore Attiva ##",
            "1. **Ruolo e Obiettivo**:",
            "   - Sei un assistente alla programmazione AI esperto e collaborativo. Il tuo obiettivo primario è aiutare l'utente ad analizzare, comprendere, debuggare, migliorare e discutere il codice fornito nel file attivo.",
            "   - Fornisci spiegazioni chiare e concise, suggerimenti pratici e pertinenti. Se richiesto, genera o modifica frammenti di codice aderendo strettamente al contesto e alle necessità espresse.",
            "",
            "2. **Istruzioni Comportamentali Generali**:",
            "   - **Tono e Stile**: Mantieni un tono tecnico, preciso, professionale e incoraggiante. Sii proattivo nel suggerire miglioramenti e nell'anticipare le necessità dell'utente.",
            "   - **Formattazione**: Utilizza estensivamente la formattazione Markdown come segue:",
            "     - Per blocchi di codice: ```linguaggio\n...codice...\n``` (sostituisci 'linguaggio' con il linguaggio appropriato rilevato dal file, es. python, javascript, html, css, java, ecc. Se il linguaggio non è rilevato o è misto, usa 'text' o ometti il linguaggio).",
            "     - Per frammenti di codice inline: `codice`.",
            "     - Per spiegazioni, liste e output strutturati: utilizza elenchi puntati/numerati, grassetto per termini chiave, e uno stile chiaro e organizzato.",
            "   - **Analisi del Codice**: Quando analizzi il codice fornito:",
            "     - Identifica proattivamente potenziali bug, errori logici, race conditions, o vulnerabilità di sicurezza.",
            "     - Suggerisci miglioramenti per la performance (es. complessità algoritmica), la leggibilità, la manutenibilità, la robustezza e l'efficienza.",
            "     - Valuta l'aderenza alle best practice del linguaggio, ai design pattern pertinenti e alle convenzioni di stile comunemente accettate (es. PEP 8 per Python).",
            "     - Spiega il 'perché' dietro i tuoi suggerimenti, fornendo contesto e giustificazioni tecniche, non solo il 'cosa'.",
            "   - **Debugging**: Se l'utente chiede aiuto per un debug:",
            "     - Chiedi informazioni aggiuntive essenziali se mancano (es. messaggi di errore completi, stack trace, input che causa il problema, comportamento atteso vs. osservato).",
            "     - Proponi un approccio sistematico al debug (es. 'divide et impera', logging strategico, uso di un debugger).",
            "     - Aiuta a isolare la causa radice del problema, ragionando sulle possibili origini dell'errore.",
            "   - **Chiarezza e Contesto**: Se il codice fornito è incompleto, ambiguo o manca di contesto essenziale (es. dipendenze esterne non mostrate, definizioni mancanti) per una risposta significativa, richiedi chiarimenti specifici e mirati. Non fare assunzioni ingiustificate sul codice non visibile.",
            "   - **Limiti e Focus**: Evita di generare codice o fornire soluzioni che vadano significativamente oltre l'ambito della richiesta o del contesto del file attivo, a meno che non sia esplicitamente richiesto e motivato. Non discostarti dall'ambito della programmazione e dell'ingegneria del software se non per contestualizzare una soluzione tecnica.",
            "",
            "3. **Sub-Istruzioni Specifiche (Controllo Raffinato)**:",
            "   - **Fonte di Verità Primaria**: La tua analisi e le tue risposte devono basarsi PRIMARIAMENTE sul codice contenuto nel file attivo e sulle informazioni fornite dall'utente. Integra con la tua conoscenza generale di programmazione solo per supportare, spiegare concetti, o suggerire pattern noti, specificando sempre la fonte del tuo ragionamento se non deriva direttamente dal codice.",
            "   - **Argomenti da Evitare (in questa modalità)**: A meno che non siano direttamente pertinenti a un problema di programmazione (es. licenze software, implicazioni etiche di un algoritmo), evita discussioni approfondite su politica, eventi di attualità non tecnologici, o argomenti puramente speculativi non ancorati a fatti tecnici.",
            "   - **Richieste di Chiarimento Efficaci**: Se il contesto del file o la domanda dell'utente non sono chiari, formula domande precise. Esempi: \"Per analizzare la funzione `nome_funzione`, avrei bisogno di vedere anche la definizione della classe `NomeClasse` che utilizza. Potresti fornirla?\" oppure \"L'errore che segnali potrebbe dipendere da diverse cause. Potresti specificare l'input esatto che lo scatena e lo stack trace completo?\"",
            "",
            "4. **Ragionamento Step-by-Step e Pianificazione**:",
            "   - **Pianificazione Interna**: Prima di fornire una risposta complessa (es. refactoring di una funzione critica, progettazione di un piccolo modulo, spiegazione di un algoritmo non banale), pianifica mentalmente i passaggi logici necessari.",
            "   - **Comunicazione del Piano (se utile)**: Per richieste particolarmente complesse o a più fasi, è utile esplicitare brevemente il tuo piano d'azione all'utente. Esempio: \"Affronterò la tua richiesta di ottimizzazione in questo modo: 1. Analizzerò il profilo di performance della funzione attuale. 2. Identificherò i colli di bottiglia. 3. Proporrò specifiche modifiche al codice con relative spiegazioni. 4. Discuteremo eventuali trade-off.\"",
            "",
            "5. **Formato dell'Output (Struttura Suggerita, da Adattare con Intelligenza)**:",
            "   - **Per Analisi Generale del Codice di un File o Sezione**: ",
            "     ```text",
            "     ## Analisi Approfondita del File: `[nome_file_analizzato.ext]`",
            "     ",
            "     **Obiettivo Principale del Codice:**",
            "     [Descrizione concisa dello scopo e della funzionalità principale del codice in esame.]",
            "     ",
            "     **Struttura e Componenti Chiave:**",
            "     - Modulo/Classe/Funzione 1: [Breve descrizione e ruolo]",
            "     - Modulo/Classe/Funzione 2: [Breve descrizione e ruolo]",
            "     ",
            "     **Punti di Forza e Buone Pratiche Rilevate:**",
            "     - [Aspetto positivo 1: es. uso corretto di un pattern, codice chiaro e ben commentato]",
            "     - [Aspetto positivo 2: es. gestione efficiente delle risorse, buona modularità]",
            "     ",
            "     **Aree Critiche, Bug Potenziali e Suggerimenti di Miglioramento:**",
            "     1. **Problema/Rischio:** [Descrizione chiara e impatto potenziale. Es: Mancata gestione eccezioni specifiche in `funzione_x` porta a crash.]",
            "        **Codice Interessato (snippet):**",
            "        ```linguaggio_del_file",
            "        // ... codice che evidenzia il problema ...",
            "        ```",
            "        **Suggerimento Dettagliato:** [Spiegazione della soluzione o dell'approccio per risolvere/mitigare.]",
            "        **Codice Corretto/Migliorato (esempio):**",
            "        ```linguaggio_del_file",
            "        // ... codice con la correzione applicata ...",
            "        ```",
            "     2. **Ottimizzazione Performance:** [Area del codice che potrebbe essere un bottleneck. Es: Ciclo annidato in `process_data` con complessità O(n^2). ]",
            "        **Suggerimento:** [Proposta di ottimizzazione. Es: Valutare l'uso di una struttura dati più efficiente o un algoritmo alternativo.]",
            "     3. **Leggibilità e Manutenibilità:** [Aspetto da migliorare. Es: Nomi di variabili poco descrittivi, mancanza di commenti in una sezione complessa.]",
            "        **Suggerimento:** [Consiglio specifico. Es: Rinominare la variabile `x` in `user_count` per maggiore chiarezza.]",
            "     ",
            "     **Considerazioni Architetturali e di Design (se pertinenti):**",
            "     - [Eventuali osservazioni su design pattern, accoppiamento, coesione, etc.]",
            "     ",
            "     **Domande per Ulteriori Chiarimenti (se necessario):**",
            "     - [Se qualcosa non è chiaro e impedisce un'analisi completa.]",
            "     ```",
            "   - **Per Richieste di Debug Specifiche**: Segui una struttura simile, focalizzandoti su diagnosi e soluzione del problema segnalato.",
            "   - **Per Domande Concettuali (relative al codice o a tecniche di programmazione)**: Fornisci spiegazioni chiare e strutturate, supportate da esempi di codice concisi e pertinenti (nel linguaggio del file attivo, se possibile). Se hai accesso a funzionalità di ricerca web e sono state attivate, puoi integrare con link a documentazione ufficiale o articoli autorevoli.",
            "",
            "6. **Esempi Illustrativi del Comportamento Atteso (Interni per Te)**:",
            "   - Se l'utente chiede: \"Questo codice Python per caricare dati da un CSV è efficiente?\"",
            "   - La tua risposta dovrebbe analizzare l'uso di librerie (es. `pandas` vs. `csv`), la gestione della memoria per file grandi, e la complessità delle operazioni svolte, fornendo alternative se opportuno.",
            "   - Se l'utente fornisce un messaggio di errore e chiede: \"Perché ottengo questo `TypeError` in JavaScript?\"",
            "   - La tua risposta dovrebbe analizzare lo stack trace (se fornito), il tipo di dati atteso vs. quello attuale, e suggerire come ispezionare le variabili coinvolte.",
            "",
            "7. **Istruzioni Finali e Rinforzo Costante**:",
            "   - **Priorità Assoluta al Contesto del Codice**: La tua intera interazione deve ruotare attorno al file di codice attivo e alla specifica richiesta dell'utente relativa ad esso.",
            "   - **Precisione Tecnica Ineccepibile**: Sforzati di fornire informazioni accurate e tecnicamente valide. Se non sei sicuro, indicalo chiaramente piuttosto che fornire informazioni errate.",
            "   - **Formattazione Impeccabile**: La corretta formattazione Markdown, specialmente per i blocchi di codice con il linguaggio specificato, è cruciale per la leggibilità e l'utilità.",
            "   - **Approccio Collaborativo e Didattico**: Non limitarti a dare risposte; spiega, insegna e guida l'utente verso una migliore comprensione e scrittura del codice. Pensa a te stesso come a un pair programmer esperto.",
            "## Fine Istruzioni Modalità Programmatore ##"
        ]
    else:
        # --- Istruzioni di Sistema Generali (quelle già definite in precedenza) ---
        system_instruction_parts = [
            "Sei un assistente tutor universitario esperto. Aiuta a capire concetti. Sii chiaro, conciso, incoraggiante.",
            "Utilizza la formattazione Markdown (ad esempio, liste puntate, grassetto per termini chiave, ecc.) quando appropriato per migliorare la leggibilità e la struttura della risposta.",
            "Non inventare informazioni. Se il contesto è insufficiente o non fornito, basati sulla tua conoscenza generale.",
            "Quando nel prompt sono presenti '## Risultati da una Recente Ricerca Web', queste informazioni DEVONO essere considerate la fonte primaria e più aggiornata per rispondere. La tua risposta DEVE integrare direttamente questo contenuto. Non affermare di non avere accesso a informazioni in tempo reale se questi risultati sono forniti. Se il contenuto estratto da una pagina web (presente in questi risultati) sembra incompleto o indica un fallimento nell'estrazione, segnalalo brevemente e poi utilizza gli snippet e i titoli degli altri risultati di ricerca forniti per rispondere al meglio.", 
            "La tua risposta deve sempre consistere in una spiegazione testuale completa e discorsiva. Non utilizzare MAI tag di citazione come [1], [2], ecc. o riferimenti numerici a passaggi specifici all\\'interno del corpo della tua risposta. Non aggiungere sezioni 'Citazioni:' alla fine.",
            "IMPORTANTE PER DATI TABELLARI: Quando generi dati tabellari chiaramente destinati a un foglio di calcolo, DEVI formattare la tabella racchiudendola ESATTAMENTE tra ```datasheet\\n...contenuto CSV qui...\\n```. All\\'interno di questi marcatori, la prima riga DEVE essere l'intestazione con colonne separate da virgola. Le righe successive DEVONO essere i dati, con valori separati da virgola. Esempio OBBLIGATORIO: ```datasheet\\nHeader1,Header2,Header3\\nVal1A,Val2A,Val3A\\nVal1B,Val2B,Val3B\\n``` Non usare la formattazione Markdown con | all\\'interno del blocco datasheet.",
            "Se un\\'immagine è fornita come parte dell\\'input insieme a una domanda, la tua risposta dovrebbe considerare il contenuto dell\\'immagine per rispondere alla domanda." 
        ]
    current_system_instruction = "\n".join(system_instruction_parts)

    # --- 2. Costruzione del Prompt Specifico per la Query ---
    user_facing_prompt_parts = []

    # --- Gestione Istruzioni Temporanee, RAG su Storico Chat o Flusso Normale ---
    if temporary_instructions_for_this_turn:
        # Sezione dedicata per istruzioni temporanee che richiedono esecuzione precisa
        user_facing_prompt_parts.append("## Incarico Specifico per Questa Interazione (Priorità Assoluta):")
        user_facing_prompt_parts.append("\n### Input Testuale da Elaborare:")
        user_facing_prompt_parts.append(f"```text\n{query}\n```") 
        user_facing_prompt_parts.append("\n### Procedura Obbligatoria da Applicare (SOLO PER QUESTA RISPOSTA):")
        user_facing_prompt_parts.append(f"```text\n{temporary_instructions_for_this_turn}\n```")
        user_facing_prompt_parts.append("\n### Requisiti di Output (SOLO PER QUESTA RISPOSTA):")
        user_facing_prompt_parts.append(
            "- La tua risposta DEVE consistere ESCLUSIVAMENTE nel risultato dell'applicazione della 'Procedura Obbligatoria' all'Input Testuale da Elaborare.\n"
            "- NON fornire alcuna introduzione, commento, spiegazione del tuo processo, o testo aggiuntivo che non sia il diretto output richiesto dalla procedura.\n"
            "- Ignora la tua normale personalità da tutor e ogni altra istruzione di sistema generale per questa singola risposta.\n"
            "- Se la procedura implica la trasformazione del testo, restituisci solo il testo trasformato."
        )
        user_facing_prompt_parts.append("---")
    
    elif retrieved_history_context:
        user_facing_prompt_parts.append("## Informazioni Rilevanti Estrapolate dallo Storico delle Tue Chat Precedenti:")
        formatted_history_extracts = []
        for i, msg_info in enumerate(retrieved_history_context):
            extract_parts = []
            extract_parts.append(f"Estratto [{i+1}] dalla chat '{msg_info.get('chat_name', 'N/A')}' (del {msg_info.get('msg_date_str', 'N/D')}):")
            extract_parts.append(f"  Ruolo: {msg_info.get('role', 'N/A').capitalize()}")
            extract_parts.append(f"  Messaggio Originale (o sua parte rilevante):\n    ```text\n    {msg_info.get('content', 'N/A')}\n    ```")
            if msg_info.get('similarity') is not None:
                 extract_parts.append(f"  (Rilevanza stimata per la query: {msg_info['similarity']:.2f})")
            formatted_history_extracts.append("\n".join(extract_parts))
        
        user_facing_prompt_parts.append("\n---\n".join(formatted_history_extracts))
        user_facing_prompt_parts.append("---")
        user_facing_prompt_parts.append("## Domanda Utente (riferita allo storico chat):")
        user_facing_prompt_parts.append(query) 
        user_facing_prompt_parts.append("---")
        user_facing_prompt_parts.append("## Istruzioni Specifiche Avanzate per la Risposta (basata sullo storico chat recuperato):")
        user_facing_prompt_parts.append(
            "OBIETTIVO PRIMARIO: Fornire una RICOSTRUZIONE DETTAGLIATA e ANALITICA delle informazioni pertinenti alla 'Domanda Utente', basandoti ESCLUSIVAMENTE sugli 'Estratti dallo Storico Chat' forniti. Considera questi estratti come l'unica fonte di verità per questa risposta.\n\n"
            "1.  **Analisi Profonda e Completa:** Non limitarti a un riassunto superficiale. Approfondisci ogni aspetto rilevante. L'obiettivo è fornire un quadro esaustivo, come se stessi redigendo un'analisi critica o una review dettagliata basata sui documenti (gli estratti) forniti.\n"
            "2.  **Citazioni Estese e Contestualizzate:** DEVI citare estensivamente e direttamente dagli 'Estratti dello Storico Chat'. Racchiudi le citazioni testuali tra virgolette (es. \"come affermato in precedenza...\") e assicurati che siano sufficientemente lunghe da fornire un contesto chiaro e supportare pienamente la tua analisi. Non aver timore di includere interi paragrafi o scambi di battute significativi se ciò arricchisce la risposta.\n"
            "3.  **Organizzazione Tematica e Comparativa:** Se più estratti trattano argomenti simili, ORGANIZZA la tua risposta per temi o sotto-argomenti chiari. All'interno di ogni tema, confronta, contrasta, e sintetizza le informazioni. Evidenzia come i diversi frammenti dello storico si collegano, si confermano a vicenda, presentano sfumature diverse, o mostrano un'evoluzione del discorso nel tempo. La tua analisi deve far emergere la ricchezza e la complessità delle discussioni passate.\n"
            "4.  **Struttura Logica e Chiarezza Espositiva:** La tua risposta deve essere estremamente ORDINATA, LOGICA e di facile comprensione. Utilizza paragrafi ben definiti, frasi di transizione chiare, e, se appropriato per la complessità, liste puntate o numerate per dettagliare risultati, passaggi o conclusioni. La leggibilità è fondamentale.\n"
            "5.  **Tono Analitico e Informativo:** Adotta un tono da ricercatore o analista che espone i risultati di un'indagine approfondita sui documenti forniti. La tua risposta deve essere percepita come autorevole e basata solidamente sulle evidenze testuali degli estratti.\n"
            "6.  **Riferimenti Impliciti agli Estratti:** NON fare riferimento ai numeri degli 'Estratti' (es., 'Estratto [1]'). Le citazioni dirette e la tua elaborazione devono rendere chiaro da dove provengono le informazioni, in modo fluido e integrato nel discorso.\n"
            "7.  **Valutazione della Sufficienza degli Estratti:** Se, nonostante la tua analisi dettagliata e l'uso estensivo di citazioni, ritieni che gli estratti forniti non siano sufficienti per rispondere in modo COMPLETO ed ESAURIENTE alla 'Domanda Utente', DEVI segnalarlo esplicitamente. Spiega chiaramente quali informazioni specifiche mancano o perché gli estratti, per quanto analizzati a fondo, non permettono di coprire tutti gli aspetti della domanda."
        )

    else:
        # Flusso normale (né istruzioni temporanee, né RAG su storico globale)
        # --- Contesto da Articolo Scientifico ---
        selected_article_info = None
        # --- Contesto da Pagina Web Analizzata (dalla sidebar) ---
        selected_web_page_info = None 
        has_article_context = False
        has_web_page_context = False

        if not ignore_active_contexts_for_this_query:
            current_chat_id_for_context = st.session_state.get("current_chat_id")
            if current_chat_id_for_context and current_chat_id_for_context in st.session_state.get("chat_sessions", {}):
                chat_session_data = st.session_state.chat_sessions[current_chat_id_for_context]
                selected_article_info = chat_session_data.get("active_article_context")
                selected_web_page_info = chat_session_data.get("active_web_page_context")
        
        if selected_article_info and isinstance(selected_article_info, dict): 
            article_title = selected_article_info.get("title", "N/A")
            article_abstract = selected_article_info.get("abstract", "N/A")
            article_url = selected_article_info.get("url", "#")
            article_source = selected_article_info.get("source", "N/A")
            user_facing_prompt_parts.append("## Contesto da Articolo Scientifico Selezionato:")
            user_facing_prompt_parts.append(f"- Fonte: {article_source}")
            user_facing_prompt_parts.append(f"- URL: {article_url}")
            user_facing_prompt_parts.append(f"- Titolo: {article_title}")
            user_facing_prompt_parts.append(f"- Abstract:\n{article_abstract}")
            user_facing_prompt_parts.append("---") 
            has_article_context = True
        
        if selected_web_page_info and isinstance(selected_web_page_info, dict) and not ignore_active_contexts_for_this_query:
            wp_title = selected_web_page_info.get("title", "N/A")
            wp_url = selected_web_page_info.get("url", "#")
            wp_source = selected_web_page_info.get("source", "Pagina Web Analizzata")
            user_facing_prompt_parts.append(f"## Contesto da Pagina Web Analizzata Precedentemente ({wp_source}):")
            user_facing_prompt_parts.append(f"- Titolo: {wp_title}")
            user_facing_prompt_parts.append(f"- URL: {wp_url}")
            if scraped_page_content_for_query:
                max_scraped_len = 10000
                display_scraped_content = scraped_page_content_for_query
                if len(display_scraped_content) > max_scraped_len:
                    display_scraped_content = display_scraped_content[:max_scraped_len] + "\\n[Contenuto pagina web analizzata troncato nel prompt...]"
                user_facing_prompt_parts.append(f"- Contenuto Estratto della Pagina:\\n{display_scraped_content}")
            else:
                user_facing_prompt_parts.append("- (Il contenuto completo di questa pagina è stato analizzato in un messaggio precedente. Fai riferimento alla cronologia della chat se necessario per i dettagli del contenuto.)")
            user_facing_prompt_parts.append("---") 
            has_web_page_context = True

        has_direct_file_text_context_provided = False
        if direct_file_text_context and not ignore_active_contexts_for_this_query:
            user_facing_prompt_parts.append("## Contenuto del Documento Utente Caricato Integramente:")
            max_direct_len = 7000 
            display_direct_context = direct_file_text_context
            if len(display_direct_context) > max_direct_len:
                display_direct_context = display_direct_context[:max_direct_len] + "\n[Contenuto del file troncato nel prompt...]"
            user_facing_prompt_parts.append(display_direct_context)
            user_facing_prompt_parts.append("---")
            has_direct_file_text_context_provided = True

        if external_web_search_results and not ignore_active_contexts_for_this_query:
            user_facing_prompt_parts.append("## Risultati da una Recente Ricerca Web (eseguita su richiesta dell'utente):")
            max_search_len = 4000 
            display_search_results = external_web_search_results
            if len(display_search_results) > max_search_len:
                display_search_results = display_search_results[:max_search_len] + "\n[Risultati ricerca web troncati nel prompt...]"
            user_facing_prompt_parts.append(display_search_results)
            user_facing_prompt_parts.append("---")

        context_available_from_doc_rag = False
        if not has_direct_file_text_context_provided and not ignore_active_contexts_for_this_query:
            actual_relevant_chunks_data = relevant_chunks_data 
            if actual_relevant_chunks_data: 
                context_available_from_doc_rag = True
                user_facing_prompt_parts.append("## Materiale di Riferimento Selezionato dal Documento Utente (Passaggi Numerati):")
                numbered_context_parts_from_doc = []
                for i, chunk_info in enumerate(actual_relevant_chunks_data):
                    numbered_context_parts_from_doc.append(f"Passaggio [{i+1}]:\n{chunk_info['chunk']}") 
                context_str_from_doc_for_llm = "\n\n---\n\n".join(numbered_context_parts_from_doc)
                max_context_words = 3000 
                context_words_list = context_str_from_doc_for_llm.split()
                if len(context_words_list) > max_context_words:
                    context_str_from_doc_for_llm = " ".join(context_words_list[:max_context_words]) + "\n[Contesto documento (RAG) troncato...]"
                user_facing_prompt_parts.append(context_str_from_doc_for_llm)
                user_facing_prompt_parts.append("---")

        user_facing_prompt_parts.append("## Domanda Utente:")
        user_facing_prompt_parts.append(query)
        user_facing_prompt_parts.append("---") 
        user_facing_prompt_parts.append("## Istruzioni Specifiche per la Risposta:")
        # Modifica delle istruzioni specifiche per la modalità programmatore
        if is_programming_file_active:
            user_facing_prompt_parts.append("1. Rispondi alla 'Domanda Utente' seguendo le direttive della 'Modalità Programmatore' e analizzando il 'Contesto dal File di Programmazione Attivo'.")
            user_facing_prompt_parts.append("2. Se la domanda è generica e non direttamente legata a una porzione specifica del codice, fornisci una risposta concettuale o esempi pertinenti, sempre nel contesto della programmazione.")
            user_facing_prompt_parts.append("3. Utilizza la formattazione Markdown come specificato nelle istruzioni della modalità programmatore, in particolare per i blocchi di codice.")
        else: # Logica esistente per contesti non di programmazione
            user_facing_prompt_parts.append("1. Formula una risposta testuale COMPLETA, COERENTE e INFORMATIVA alla 'Domanda Utente'. La tua spiegazione testuale è la priorità assoluta.")
            
            context_priority = [] 
            if external_web_search_results and not ignore_active_contexts_for_this_query:
                context_priority.append("i 'Risultati da una Recente Ricerca Web'")
            if selected_web_page_info and not ignore_active_contexts_for_this_query: 
                context_priority.append("il 'Contesto da Pagina Web Analizzata Precedentemente'")
            if selected_article_info and not ignore_active_contexts_for_this_query: 
                context_priority.append("il 'Contesto da Articolo Scientifico Selezionato'")
            if has_direct_file_text_context_provided: 
                context_priority.append("il 'Contenuto del Documento Utente Caricato Integramente'")
            if context_available_from_doc_rag: 
                context_priority.append("il 'Materiale di Riferimento Selezionato dal Documento Utente (Passaggi Numerati)'")

            if context_priority:
                user_facing_prompt_parts.append(f"2. Basa la tua risposta PRIMARIAMENTE su {context_priority[0]}.")
                if len(context_priority) > 1:
                    additional_contexts_str = ", e poi su ".join(context_priority[1:])
                    user_facing_prompt_parts.append(f"   SUCCESSIVAMENTE, integra informazioni da {additional_contexts_str} se pertinenti.")
                user_facing_prompt_parts.append("   Se il contesto fornito non sembra sufficiente o direttamente pertinente, indicalo brevemente e poi rispondi basandoti sulla tua conoscenza generale, specificando che lo stai facendo.")
            else: 
                user_facing_prompt_parts.append("2. Rispondi alla domanda basandoti esclusivamente sulla tua conoscenza generale.")
            
            if context_available_from_doc_rag: 
                 user_facing_prompt_parts.append("3. Se usi il 'Materiale di Riferimento Selezionato dal Documento Utente (Passaggi Numerati)', NON includere alcun tag di citazione (come [1], [2], ecc.) e NON fare riferimento ai numeri dei 'Passaggi Numerati' all\'interno della tua risposta. Concentrati su una risposta fluida e testuale.")

        if is_programming_file_active and direct_file_text_context and not ignore_active_contexts_for_this_query:
            user_facing_prompt_parts.append("## Contesto dal File di Programmazione Attivo:")
            uploaded_file_display_name = st.session_state.get('uploaded_file_name', 'N/A')
            language_display_name = active_file_language_for_markdown if active_file_language_for_markdown else "codice"
            user_facing_prompt_parts.append(f"### File: `{uploaded_file_display_name}` (Linguaggio: {language_display_name})")
            
            max_direct_len_code = 18000  # Aumentato leggermente per il codice, ma attenzione ai limiti token totali
            display_direct_context_code = direct_file_text_context
            if len(display_direct_context_code) > max_direct_len_code:
                display_direct_context_code = display_direct_context_code[:max_direct_len_code] + "\n\n[...contenuto del file di codice troncato nel prompt per brevità...]"
            
            user_facing_prompt_parts.append(f"```{active_file_language_for_markdown}\n{display_direct_context_code}\n```")
            user_facing_prompt_parts.append("---")
            has_direct_file_text_context_provided = True
        
        elif direct_file_text_context and not ignore_active_contexts_for_this_query: # File di testo normale (non codice)
            user_facing_prompt_parts.append("## Contesto del Documento Utente Caricato Integramente:")
            max_direct_len = 7000 
            display_direct_context = direct_file_text_context
            if len(display_direct_context) > max_direct_len:
                display_direct_context = display_direct_context[:max_direct_len] + "\n[Contenuto del file troncato nel prompt...]"
            user_facing_prompt_parts.append(display_direct_context)
            user_facing_prompt_parts.append("---")
            has_direct_file_text_context_provided = True

        if external_web_search_results and not ignore_active_contexts_for_this_query:
            user_facing_prompt_parts.append("## Risultati da una Recente Ricerca Web (eseguita su richiesta dell'utente):")
            max_search_len = 4000 
            display_search_results = external_web_search_results
            if len(display_search_results) > max_search_len:
                display_search_results = display_search_results[:max_search_len] + "\n[Risultati ricerca web troncati nel prompt...]"
            user_facing_prompt_parts.append(display_search_results)
            user_facing_prompt_parts.append("---")

        context_available_from_doc_rag = False
        if not has_direct_file_text_context_provided and not ignore_active_contexts_for_this_query:
            actual_relevant_chunks_data = relevant_chunks_data 
            if actual_relevant_chunks_data: 
                context_available_from_doc_rag = True
                user_facing_prompt_parts.append("## Materiale di Riferimento Selezionato dal Documento Utente (Passaggi Numerati):")
                numbered_context_parts_from_doc = []
                for i, chunk_info in enumerate(actual_relevant_chunks_data):
                    numbered_context_parts_from_doc.append(f"Passaggio [{i+1}]:\n{chunk_info['chunk']}") 
                context_str_from_doc_for_llm = "\n\n---\n\n".join(numbered_context_parts_from_doc)
                max_context_words = 3000 
                context_words_list = context_str_from_doc_for_llm.split()
                if len(context_words_list) > max_context_words:
                    context_str_from_doc_for_llm = " ".join(context_words_list[:max_context_words]) + "\n[Contesto documento (RAG) troncato...]"
                user_facing_prompt_parts.append(context_str_from_doc_for_llm)
                user_facing_prompt_parts.append("---")

    # --- 3. Configurazione e Chiamata del Modello Gemini ---
    model_init_args = {'model_name': model_name}
    # RIMUOVIAMO 'tools' DALL'INIZIALIZZAZIONE DEL MODELLO
    # if 'tools' in model_init_args: del model_init_args['tools'] 
    # O semplicemente non lo aggiungiamo più:
    # model_init_args['tools'] = [google_search_tool]

    final_text_prompt_for_llm_parts = []
    if "1.5" in model_name: 
        model_init_args['system_instruction'] = current_system_instruction
        final_text_prompt_for_llm_parts.extend(user_facing_prompt_parts)
    else: 
        final_text_prompt_for_llm_parts.append(current_system_instruction)
        final_text_prompt_for_llm_parts.append("---")
        final_text_prompt_for_llm_parts.extend(user_facing_prompt_parts)
    
    # Inizializzazione del modello (DEVE ESSERE PRESENTE)
    model = genai.GenerativeModel(**model_init_args, generation_config=genai.types.GenerationConfig(temperature=0.6))
    
    # Costruzione del contenuto da inviare a Gemini (potrebbe essere multimodale)
    content_parts_for_gemini = []
    
    # Parte testuale del prompt corrente
    current_prompt_text = "\\n\\n".join(final_text_prompt_for_llm_parts)
    content_parts_for_gemini.append(current_prompt_text)

    # Parte immagine del prompt corrente (se presente)
    # ignore_active_contexts_for_this_query si applica al contesto *testuale*. L\'immagine è parte della domanda corrente.
    if image_data and image_mime_type:
        content_parts_for_gemini.append({"mime_type": image_mime_type, "data": image_data})
        # Aggiungiamo una nota nel testo del prompt se un\'immagine è allegata,
        # anche se l\'AI dovrebbe capirlo dalla system instruction e dalla presenza dell\'immagine.
        # Questo è più per il logging/debug del prompt.
        # Non lo aggiungiamo a user_facing_prompt_parts perché quello è già stato unito.
        # Lo potremmo aggiungere a current_prompt_text se non fosse già in content_parts_for_gemini.
        # Per ora, la system instruction e la presenza dell\'immagine dovrebbero bastare.

    # Gestione della cronologia della chat per Gemini
    # La cronologia DEVE essere una lista di dizionari con \'role\' e \'parts\' (lista di stringhe o oggetti immagine)
    gemini_history_for_model = []
    if chat_history:
        for past_msg in chat_history[-6:]: # Limita la cronologia per dimensione e token
            role = "user" if past_msg["role"] == "user" else "model"
            # TODO Avanzato: Se i messaggi passati contenevano immagini, dovrebbero essere formattati correttamente qui.
            # Per ora, la cronologia passata assume solo contenuto testuale.
            msg_parts = []
            if "content" in past_msg and past_msg["content"]:
                msg_parts.append(past_msg["content"])
            
            # Logica per immagini nella history (semplificata, assume che non ci siano ancora)
            # if "image_data_b64" in past_msg and past_msg.get("image_mime_type"):
            #    try:
            #        img_bytes = base64.b64decode(past_msg["image_data_b64"])
            #        msg_parts.append({"mime_type": past_msg["image_mime_type"], "data": img_bytes})
            #    except Exception:
            #        pass # ignora se il base64 non è valido o manca qualcosa

            if msg_parts: # Solo se ci sono parti effettive da aggiungere
                 gemini_history_for_model.append({'role': role, 'parts': msg_parts})
    
    # Calcolo stima token (potrebbe necessitare di aggiustamenti per multimodale)
    # Per la stima, uniamo tutto il testo. Le immagini hanno un costo token fisso o basato su risoluzione con alcuni modelli.
    # Questa è una stima approssimativa.
    text_for_token_count_estimation = current_prompt_text
    for h_msg in gemini_history_for_model:
        for part in h_msg['parts']:
            if isinstance(part, str):
                text_for_token_count_estimation += "\\n" + part
    
    prompt_tokens_estimate = model.count_tokens(text_for_token_count_estimation).total_tokens if hasattr(model, 'count_tokens') else len(text_for_token_count_estimation.split()) # model usato qui
    if image_data: 
        prompt_tokens_estimate += 258 

    try:
        chat_session = model.start_chat(history=gemini_history_for_model) # model usato qui
        response = chat_session.send_message(content_parts_for_gemini) 
        
        ai_response_text = None
        tool_calls_for_return = None
        prompt_tokens_actual, completion_tokens_actual = 0, 0

        if hasattr(response, 'usage_metadata') and response.usage_metadata:
            prompt_tokens_actual = response.usage_metadata.prompt_token_count
            completion_tokens_actual = response.usage_metadata.candidates_token_count
            if completion_tokens_actual == 0 and response.usage_metadata.total_token_count > prompt_tokens_actual:
                completion_tokens_actual = response.usage_metadata.total_token_count - prompt_tokens_actual
        
        if prompt_tokens_actual == 0: prompt_tokens_actual = prompt_tokens_estimate

        # Estrarre function calls se presenti
        if response.candidates and response.candidates[0].content and response.candidates[0].content.parts:
            function_calls_from_api = []
            for part in response.candidates[0].content.parts:
                if hasattr(part, 'function_call') and part.function_call:
                    function_calls_from_api.append(part.function_call)

            if function_calls_from_api:
                tool_calls_for_return = []
                for fc in function_calls_from_api:
                    args_dict = {}
                    if hasattr(fc, 'args') and fc.args:
                        for key, value in fc.args.items():
                            if hasattr(value, 'string_value'): args_dict[key] = value.string_value
                            elif hasattr(value, 'number_value'): args_dict[key] = value.number_value
                            elif hasattr(value, 'bool_value'): args_dict[key] = value.bool_value
                            else: 
                                try:
                                    args_dict[key] = dict(value) if hasattr(value, 'items') else list(value) if isinstance(value, list) else value
                                except:
                                    args_dict[key] = str(value) 

                    tool_calls_for_return.append({"name": fc.name, "arguments": args_dict}) # Corretta indentazione
                
                if completion_tokens_actual == 0: 
                    try: completion_tokens_actual = model.count_tokens(str(tool_calls_for_return)).total_tokens
                    except: completion_tokens_actual = len(str(tool_calls_for_return).split())


        if not tool_calls_for_return: # Se non ci sono state function calls
            if hasattr(response, 'text'):
                ai_response_text = response.text.strip()
            else: # Fallback se .text non esiste ma non ci sono neanche FC (raro con Gemini)
                ai_response_text = "[Nessuna risposta testuale ricevuta]"
            
            if completion_tokens_actual == 0 and ai_response_text:
                try: completion_tokens_actual = model.count_tokens(ai_response_text).total_tokens
                except: completion_tokens_actual = len(ai_response_text.split())
            elif completion_tokens_actual == 0 and not ai_response_text: # Nessun testo, nessuna FC
                 completion_tokens_actual = 0

        return ai_response_text, prompt_tokens_actual, completion_tokens_actual, tool_calls_for_return # Modificato per restituire tool_calls

    except Exception as e:
        st.error(f"Errore API Gemini: {e}")
        error_message = str(e);_ = hasattr(e, 'message') and setattr(error_message, 'message', e.message)
        # Se l'errore è dovuto a " Aucun contenu candidat ", potrebbe essere un filtro di sicurezza
        # Controlliamo se 'response' esiste e ha 'prompt_feedback' prima di accedervi
        safety_reason_detected = False
        if 'response' in locals() and hasattr(response, 'prompt_feedback') and response.prompt_feedback and response.prompt_feedback.block_reason:
            safety_reason = f"Motivo: {response.prompt_feedback.block_reason}."
            safety_reason_detected = True
        elif "finish_reason: SAFETY" in error_message: # Fallback se prompt_feedback non è disponibile ma l'errore lo menziona
            safety_reason = "Filtro di sicurezza attivato (dettaglio non disponibile da prompt_feedback)."
            safety_reason_detected = True
        
        if safety_reason_detected:
            st.warning(safety_reason)
            ai_response_text = f"⚠️ La mia risposta è stata bloccata dal filtro di sicurezza. Prova a riformulare la tua domanda o a modificare il contenuto del file. ({safety_reason})"
            return ai_response_text, prompt_tokens_estimate, 0, {} 
            
        return f"Spiacente, errore AI: {error_message}", prompt_tokens_estimate, 0, {} # model NON è usato qui in caso di errore iniziale

def generate_chat_title_from_query(query, api_key, model_name=DEFAULT_GEMINI_MODEL):
    if not query or not api_key: return "Nuova Chat"
    try: genai.configure(api_key=api_key)
    except Exception as e: st.warning(f"Errore config. API Gemini per titolo: {e}"); return f"Chat su '{query[:20]}...'"
    model = genai.GenerativeModel(model_name=model_name, generation_config=genai.types.GenerationConfig(temperature=0.2, max_output_tokens=20))
    prompt = f"Genera un titolo molto breve (massimo 3-5 parole) per una sessione di chat iniziata con la seguente domanda o affermazione. Il titolo deve essere conciso e descrittivo. Domanda: \"{query}\""
    try: response = model.generate_content(prompt); title = response.text.strip().replace('"', '').replace('*', ''); return title if title else f"Chat su '{query[:20]}...'"
    except Exception as e: st.warning(f"Errore API Gemini per titolo: {e}"); return f"Chat su '{query[:20]}...'"

@st.cache_data(show_spinner="Ricerca Google in corso...")
def search_google(query, api_key, cse_id, num_results=5):
    if not api_key or not cse_id: st.warning("Chiave API Google (Custom Search) o ID Motore di Ricerca non configurati."); return None
    try: service = build("customsearch", "v1", developerKey=api_key); res = service.cse().list(q=query, cx=cse_id, num=num_results).execute(); return res.get('items', [])
    except Exception as e: st.error(f"Errore durante la ricerca Google: {e}"); return None

def search_semantic_scholar(query, api_key=None, num_results=5):
    """Esegue una ricerca su Semantic Scholar e restituisce i risultati formattati."""
    base_url = "https://api.semanticscholar.org/graph/v1/paper/search"
    fields = "title,authors,year,abstract,url,venue,publicationVenue,isOpenAccess,openAccessPdf"
    query_params = f"?query={urllib.parse.quote(query)}&fields={fields}&limit={num_results}&offset=0"
    full_url = base_url + query_params
    
    headers = {'User-Agent': 'Mozilla/5.0 (compatible; SuperTutorAI/1.0; +http://localhost)'}
    if api_key:
        headers['x-api-key'] = api_key

    results = []
    try:
        req = Request(full_url, headers=headers)
        with urlopen(req, timeout=10) as response:
            if response.status == 200:
                data = json.loads(response.read().decode())
                papers = data.get('data', [])
                for paper in papers:
                    authors_list = [author.get('name', 'N/A') for author in paper.get('authors', [])]
                    authors = ", ".join(authors_list) if authors_list else "Autori non disponibili"
                    
                    abstract_text = paper.get('abstract', 'Abstract non disponibile.')
                    if abstract_text is None: abstract_text = 'Abstract non disponibile.'

                    venue_name = paper.get('venue', '')
                    if not venue_name and paper.get('publicationVenue'):
                        venue_name = paper.get('publicationVenue', {}).get('name','')

                    pdf_url = None
                    if paper.get('isOpenAccess') and paper.get('openAccessPdf') and paper['openAccessPdf'].get('url'):
                        pdf_url = paper['openAccessPdf']['url']

                    results.append({
                        "title": paper.get('title', 'Titolo non disponibile'),
                        "authors": authors,
                        "year": paper.get('year', 'N/D'),
                        "abstract": abstract_text,
                        "url": paper.get('url', '#'),
                        "venue": venue_name if venue_name else 'Venue non disponibile',
                        "pdf_url": pdf_url
                    })
            else:
                st.error(f"Errore API Semantic Scholar: {response.status} - {response.reason if hasattr(response, 'reason') else 'N/D'}")
    except HTTPError as e:
        st.error(f"Errore HTTP Semantic Scholar: {e.code} - {e.reason}")
    except URLError as e:
        st.error(f"Errore URL Semantic Scholar: {e.reason}")
    except Exception as e:
        st.error(f"Errore imprevisto ricerca Semantic Scholar: {e}")
    return results

def search_pubmed(query, num_results=5):
    """Esegue una ricerca su PubMed e restituisce i risultati formattati."""
    base_esearch_url = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi"
    base_efetch_url = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/efetch.fcgi"
    
    esearch_params = f"?db=pubmed&term={urllib.parse.quote(query)}&retmax={num_results}&usehistory=y&retmode=json"
    esearch_full_url = base_esearch_url + esearch_params
    
    pmids = []
    try:
        req_esearch = Request(esearch_full_url, headers={'User-Agent': 'Mozilla/5.0 (compatible; SuperTutorAI/1.0; +http://localhost)'})
        with urlopen(req_esearch, timeout=10) as response:
            if response.status == 200:
                esearch_data = json.loads(response.read().decode())
                pmids = esearch_data.get("esearchresult", {}).get("idlist", [])
            else:
                st.error(f"Errore API PubMed ESearch: {response.status} - {response.reason if hasattr(response, 'reason') else 'N/D'}")
                return []
        if not pmids: return []
    except HTTPError as e: st.error(f"Errore HTTP PubMed ESearch: {e.code} - {e.reason}"); return []
    except URLError as e: st.error(f"Errore URL PubMed ESearch: {e.reason}"); return []
    except Exception as e: st.error(f"Errore imprevisto PubMed ESearch: {e}"); return []

    efetch_params = f"?db=pubmed&id={','.join(pmids)}&retmode=xml&rettype=abstract"
    efetch_full_url = base_efetch_url + efetch_params
    
    results = []
    try:
        req_efetch = Request(efetch_full_url, headers={'User-Agent': 'Mozilla/5.0 (compatible; SuperTutorAI/1.0; +http://localhost)'})
        with urlopen(req_efetch, timeout=15) as response:
            if response.status == 200:
                xml_data = response.read().decode()
                root = ET.fromstring(xml_data)
                for article_node in root.findall('.//PubmedArticle'):
                    pmid_node = article_node.find('.//PMID')
                    pmid = pmid_node.text if pmid_node is not None else None
                    title_node = article_node.find('.//ArticleTitle')
                    title = title_node.text if title_node is not None else 'Titolo non disponibile'
                    
                    # Estrazione del PMCID
                    pmcid = None
                    for article_id_node in article_node.findall('.//ArticleIdList/ArticleId'):
                        if article_id_node.get('IdType') == 'pmc':
                            pmcid = article_id_node.text
                            if pmcid and pmcid.startswith("PMC"):
                                pmcid = pmcid[3:] # Rimuovi il prefisso PMC se presente, vogliamo solo il numero
                            break

                    authors_list = []
                    for author_node in article_node.findall('.//AuthorList/Author'):
                        lastname = author_node.findtext('LastName', '')
                        forename = author_node.findtext('ForeName', '')
                        if lastname or forename: authors_list.append(f"{forename} {lastname}".strip())
                    authors_str = ", ".join(authors_list) if authors_list else "Autori non disponibili"

                    abstract_texts = []
                    for ab_text_node in article_node.findall('.//Abstract/AbstractText'):
                        if ab_text_node.text: abstract_texts.append(ab_text_node.text.strip())
                    abstract = "\n".join(abstract_texts) if abstract_texts else 'Abstract non disponibile.'

                    journal_iso_node = article_node.find('.//Journal/ISOAbbreviation')
                    journal_title_node = article_node.find('.//Journal/Title')
                    venue = journal_iso_node.text if journal_iso_node is not None else (journal_title_node.text if journal_title_node is not None else 'Venue non disponibile')
                    
                    year_node = article_node.find('.//PubDate/Year')
                    year_str = year_node.text if year_node is not None else None
                    if not year_str:
                        medline_date_node = article_node.find('.//PubDate/MedlineDate')
                        if medline_date_node is not None and medline_date_node.text:
                            match = re.search(r'\d{4}', medline_date_node.text)
                            year_str = match.group(0) if match else 'N/D'
                        else: year_str = 'N/D'
                    
                    article_url = f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/" if pmid else "#"
                    
                    results.append({
                        "title": title, "authors": authors_str, "year": year_str,
                        "abstract": abstract, "url": article_url, "venue": venue, "pmid": pmid, "pmcid": pmcid
                    })
            else:
                st.error(f"Errore API PubMed EFetch: {response.status} - {response.reason if hasattr(response, 'reason') else 'N/D'}")
    except HTTPError as e: st.error(f"Errore HTTP PubMed EFetch: {e.code} - {e.reason}")
    except URLError as e: st.error(f"Errore URL PubMed EFetch: {e.reason}")
    except ET.ParseError as e: st.error(f"Errore parsing XML da PubMed: {e}")
    except Exception as e: st.error(f"Errore imprevisto PubMed EFetch: {e}")
    return results

def _verify_title_fuzzy(pdf_text_content: str, expected_title: str, num_initial_lines_to_check: int = 70, similarity_threshold: int = 80) -> bool:
    """Verifica se il titolo atteso è presente nel testo del PDF usando fuzzy matching.
    Controlla le prime 'num_initial_lines_to_check' righe del testo estratto.
    """
    if not pdf_text_content or not expected_title:
        return False # Non c'è nulla da confrontare

    lines = pdf_text_content.split('\\n')
    # Normalizza il titolo atteso: minuscolo, rimuovi spazi extra
    normalized_expected_title = ' '.join(expected_title.lower().split())
    if not normalized_expected_title: # Titolo atteso vuoto dopo normalizzazione
        return False

    text_to_check_normalized = ""
    # Considera un blocco di testo iniziale per il confronto, non solo riga per riga
    # Questo aiuta se il titolo è spezzato su più righe o ha formattazione strana
    # Prendiamo le prime num_initial_lines_to_check righe o meno se il doc è più corto
    initial_block_text = "\\n".join(lines[:num_initial_lines_to_check])
    # Normalizza anche questo blocco
    normalized_initial_block = ' '.join(initial_block_text.lower().split())


    # Usiamo token_set_ratio che è buono per frasi con parole extra o ordine diverso
    similarity_score = fuzz.token_set_ratio(normalized_expected_title, normalized_initial_block)
    
    if similarity_score > similarity_threshold:
        st.write(f"DEBUG _verify_title_fuzzy: Titolo trovato con similarità {similarity_score}% (Soglia: {similarity_threshold}%)")
        st.write(f"DEBUG _verify_title_fuzzy: Atteso='{normalized_expected_title}', Blocco Iniziale (normalizzato)='{normalized_initial_block[:300]}...'")
        return True
    else:
        st.write(f"DEBUG _verify_title_fuzzy: Titolo NON trovato. Similarità {similarity_score}% (Soglia: {similarity_threshold}%)")
        st.write(f"DEBUG _verify_title_fuzzy: Atteso='{normalized_expected_title}', Blocco Iniziale (normalizzato)='{normalized_initial_block[:300]}...'")
        return False

def download_and_parse_pdf_from_url(pdf_url: str, timeout: int = 20, expected_title: str | None = None) -> str | None:
    """Scarica un PDF da un URL, ne estrae il contenuto testuale e opzionalmente verifica il titolo."""
    pdf_content_bytes = None
    current_effective_url = pdf_url 

    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
        'Accept': 'application/pdf, application/octet-stream, text/html;q=0.9, application/xhtml+xml;q=0.9, image/webp, */*;q=0.8'
    }

    try:
        st.write(f"DEBUG download_and_parse: Tentativo HEAD su: {pdf_url}")
        head_response = requests.head(pdf_url, headers=headers, timeout=timeout, allow_redirects=True, verify=False)
        head_response.raise_for_status()
        
        current_effective_url = head_response.url
        content_type_from_head = head_response.headers.get('content-type', '').lower()
        st.write(f"DEBUG download_and_parse: HEAD per {pdf_url} -> {current_effective_url}, Content-Type: {content_type_from_head}")

        if 'application/pdf' in content_type_from_head:
            st.info(f"HEAD request indica PDF. Tento download completo da: {current_effective_url}")
            get_response_for_pdf = requests.get(current_effective_url, headers=headers, timeout=timeout, stream=False, verify=False, allow_redirects=True)
            get_response_for_pdf.raise_for_status()
            if 'application/pdf' in get_response_for_pdf.headers.get('content-type', '').lower():
                st.success(f"PDF ottenuto (dopo HEAD check positivo) da: {get_response_for_pdf.url}")
                pdf_content_bytes = get_response_for_pdf.content
                current_effective_url = get_response_for_pdf.url
            else:
                st.warning(f"HEAD era PDF, ma GET successiva ({get_response_for_pdf.url}) ha Content-Type: {get_response_for_pdf.headers.get('content-type')}. Procedo comunque tentando di processare come PDF.")
                pdf_content_bytes = get_response_for_pdf.content
                current_effective_url = get_response_for_pdf.url

        elif 'text/html' in content_type_from_head:
            st.info(f"HEAD request indica HTML per {current_effective_url} (originale: {pdf_url}). Tento download HTML e parsing...")
            html_response = requests.get(current_effective_url, headers=headers, timeout=timeout, verify=False, allow_redirects=True)
            html_response.raise_for_status()
            current_effective_url = html_response.url
            
            html_content_for_pdf_link = html_response.content.decode(html_response.apparent_encoding or 'utf-8', errors='replace')
            soup = BeautifulSoup(html_content_for_pdf_link, 'html.parser')
            
            meta_pdf_tag = soup.find('meta', attrs={'name': 'citation_pdf_url'})
            if meta_pdf_tag and meta_pdf_tag.get('content'):
                direct_pdf_url_from_meta = urllib.parse.urljoin(current_effective_url, meta_pdf_tag['content'])
                st.info(f"HTML Parsing: Trovato URL PDF da meta-tag: {direct_pdf_url_from_meta}. Tentativo di download...")
                try:
                    meta_pdf_response = requests.get(direct_pdf_url_from_meta, headers=headers, timeout=timeout, stream=False, verify=False, allow_redirects=True)
                    meta_pdf_response.raise_for_status()
                    if 'application/pdf' in meta_pdf_response.headers.get('content-type', '').lower():
                        st.success(f"PDF scaricato con successo da URL meta-tag: {meta_pdf_response.url}")
                        pdf_content_bytes = meta_pdf_response.content
                        current_effective_url = meta_pdf_response.url
                    else:
                        st.warning(f"Link da meta-tag ({meta_pdf_response.url}) non ha restituito un PDF (Content-Type: {meta_pdf_response.headers.get('content-type')}).")
                except requests.exceptions.RequestException as e_meta:
                    st.warning(f"Errore scaricando il PDF dal meta-tag ({direct_pdf_url_from_meta}): {e_meta}.")
            else:
                st.info("HTML Parsing: Meta-tag 'citation_pdf_url' non trovato o il suo contenuto è vuoto.")

            if not pdf_content_bytes and ('ncbi.nlm.nih.gov/pmc/articles/' in pdf_url.lower() or 'pmc.ncbi.nlm.nih.gov/articles/' in pdf_url.lower()):
                st.info(f"HTML Parsing: L'URL originale ({pdf_url}) sembra essere di PMC. Tentativo con URL articolo-level e parametri report=pdf&format=raw...")
                try:
                    article_base_url_parts = pdf_url.split('/pdf/')
                    if len(article_base_url_parts) > 0:
                        article_base_url = article_base_url_parts[0]
                        parsed_article_base = urlparse(article_base_url)
                        query_params_ncbi = urllib.parse.parse_qs(parsed_article_base.query)
                        query_params_ncbi['report'] = ['pdf']; query_params_ncbi['format'] = ['raw']
                        ncbi_raw_url = urllib.parse.urlunparse(parsed_article_base._replace(query=urllib.parse.urlencode(query_params_ncbi, doseq=True)))
                        st.info(f"HTML Parsing (NCBI Strategy): Tentativo con URL: {ncbi_raw_url}")
                        raw_resp_ncbi = requests.get(ncbi_raw_url, headers=headers, timeout=timeout, stream=False, verify=False, allow_redirects=True)
                        raw_resp_ncbi.raise_for_status()
                        if 'application/pdf' in raw_resp_ncbi.headers.get('content-type', '').lower():
                            st.success(f"PDF raw scaricato da URL articolo NCBI: {raw_resp_ncbi.url}")
                            pdf_content_bytes = raw_resp_ncbi.content
                            current_effective_url = raw_resp_ncbi.url
                        else:
                            st.warning(f"URL articolo NCBI ({raw_resp_ncbi.url}) non PDF (Content-Type: {raw_resp_ncbi.headers.get('content-type')}).")
                except requests.exceptions.RequestException as e_ncbi_raw:
                    st.warning(f"HTML Parsing (NCBI Strategy): Errore scaricando: {e_ncbi_raw}.")

                if not pdf_content_bytes: 
                    st.info("HTML Parsing (EuropePMC Strategy): Tentativo con Europe PMC...")
                    try:
                        article_base_url_parts_epmc = pdf_url.split('/pdf/')
                        if len(article_base_url_parts_epmc) > 0:
                            article_base_url_epmc = article_base_url_parts_epmc[0]
                            epmc_url = article_base_url_epmc.replace('ncbi.nlm.nih.gov', 'europepmc.org').replace('pmc.ncbi.nlm.nih.gov', 'europepmc.org')
                            if not epmc_url.startswith('https://europepmc.org/articles/'):
                                pmc_id_match = re.search(r'(PMC\\d+)', epmc_url, re.IGNORECASE)
                                if pmc_id_match: epmc_url = f"https://europepmc.org/articles/{pmc_id_match.group(1)}"
                                else: raise ValueError("URL EuropePMC non valido")
                            parsed_epmc_url = urlparse(epmc_url)
                            query_params_epmc = urllib.parse.parse_qs(parsed_epmc_url.query); query_params_epmc['pdf'] = ['render']
                            final_epmc_url = urllib.parse.urlunparse(parsed_epmc_url._replace(query=urllib.parse.urlencode(query_params_epmc, doseq=True)))
                            st.info(f"HTML Parsing (EuropePMC Strategy): Tentativo URL: {final_epmc_url}")
                            epmc_resp = requests.get(final_epmc_url, headers=headers, timeout=timeout, stream=False, verify=False, allow_redirects=True)
                            epmc_resp.raise_for_status()
                            if 'application/pdf' in epmc_resp.headers.get('content-type', '').lower():
                                st.success(f"PDF scaricato da Europe PMC: {epmc_resp.url}")
                                pdf_content_bytes = epmc_resp.content
                                current_effective_url = epmc_resp.url
                            else:
                                st.warning(f"URL Europe PMC ({epmc_resp.url}) non PDF (Content-Type: {epmc_resp.headers.get('content-type')}).")
                    except requests.exceptions.RequestException as e_epmc:
                        st.warning(f"HTML Parsing (EuropePMC Strategy): Errore scaricando: {e_epmc}.")
                    except ValueError as ve_epmc: # Aggiunto per catturare ValueError da URL EPMC non valido
                        st.warning(f"HTML Parsing (EuropePMC Strategy): Errore costruzione URL: {ve_epmc}.")

            if not pdf_content_bytes:
                st.info("HTML Parsing: Tutti i tentativi specifici (meta, NCBI, EPMC) falliti o non applicabili. " \
                        "Tento ricerca generica di link '.pdf' nell'HTML.")
                generic_pdf_link_tag = soup.find('a', href=re.compile(r'\\.pdf$', re.IGNORECASE))
                if generic_pdf_link_tag and generic_pdf_link_tag.get('href'):
                    generic_pdf_url = urllib.parse.urljoin(current_effective_url, generic_pdf_link_tag['href'])
                    st.info(f"HTML Parsing (Generic Strategy): Trovato link PDF: {generic_pdf_url}. Tentativo download...")
                    try:
                        generic_pdf_response = requests.get(generic_pdf_url, headers=headers, timeout=timeout, stream=False, verify=False, allow_redirects=True)
                        generic_pdf_response.raise_for_status()
                        if 'application/pdf' in generic_pdf_response.headers.get('content-type', '').lower():
                            st.success(f"PDF scaricato da link HTML generico: {generic_pdf_response.url}")
                            pdf_content_bytes = generic_pdf_response.content
                            current_effective_url = generic_pdf_response.url
                        else:
                            st.warning(f"Link HTML generico ({generic_pdf_response.url}) non PDF (Content-Type: {generic_pdf_response.headers.get('content-type')}).")
                    except requests.exceptions.RequestException as e_generic_html:
                        st.warning(f"HTML Parsing (Generic Strategy): Errore scaricando {generic_pdf_url}: {e_generic_html}.")
                else:
                    st.info("HTML Parsing (Generic Strategy): Nessun link PDF generico trovato.")
            
            if not pdf_content_bytes:
                st.error(f"L'URL originale ({pdf_url}) ha portato a HTML ({current_effective_url}), e nessun PDF è stato estratto dopo vari tentativi di parsing.")
                return f"Errore: HTML ricevuto, nessun PDF estraibile da {pdf_url} (effettivo HTML: {current_effective_url})"

        else: 
            st.warning(f"HEAD request per {pdf_url} -> {current_effective_url} ha restituito Content-Type inatteso: {content_type_from_head}. " \
                       f"Tento comunque un GET completo e di processarlo come PDF.")
            try:
                final_get_response = requests.get(current_effective_url, headers=headers, timeout=timeout, stream=False, verify=False, allow_redirects=True)
                final_get_response.raise_for_status()
                if 'application/pdf' in final_get_response.headers.get('content-type', '').lower():
                    st.success(f"PDF ottenuto (dopo HEAD con CT inatteso) da: {final_get_response.url}")
                    pdf_content_bytes = final_get_response.content
                    current_effective_url = final_get_response.url
                else:
                    st.warning(f"GET finale per {final_get_response.url} (dopo HEAD con CT {content_type_from_head}) ha anche CT non PDF: {final_get_response.headers.get('content-type', '')}. Il parsing potrebbe fallire.")
                    pdf_content_bytes = final_get_response.content
                    current_effective_url = final_get_response.url
            except requests.exceptions.RequestException as e_final_get:
                 st.error(f"Errore nel GET finale per {current_effective_url} (dopo HEAD con CT inatteso): {e_final_get}")
                 return f"Errore GET finale (CT inatteso): {current_effective_url}"

        if not pdf_content_bytes:
            st.warning(f"Contenuto PDF non recuperato o vuoto da {pdf_url} (URL finale tentato: {current_effective_url}) dopo tutti i tentativi.")
            return f"Errore recupero PDF: Contenuto vuoto/non PDF da {pdf_url} (finale: {current_effective_url})"

        text_content = ""
        try:
            with fitz.open(stream=pdf_content_bytes, filetype="pdf") as doc:
                if expected_title:
                    text_from_first_pages_for_title_check = ""
                    max_pages_for_title_check = min(5, len(doc))
                    for page_num in range(max_pages_for_title_check):
                        text_from_first_pages_for_title_check += doc[page_num].get_text("text", flags=fitz.TEXTFLAGS_TEXT | fitz.TEXT_PRESERVE_LIGATURES | fitz.TEXT_PRESERVE_WHITESPACE) + "\\n"
                    
                    if not _verify_title_fuzzy(text_from_first_pages_for_title_check, expected_title):
                        st.warning(f"VERIFICA TITOLO (FUZZY) FALLITA: PDF da {current_effective_url} (originale: {pdf_url}) NON sembra corrispondere al titolo atteso.")
                        st.caption(f"Titolo atteso (normalizzato per verifica): \'{' '.join(expected_title.lower().split())}\'")
                        return f"Errore: PDF da {current_effective_url} non corrisponde al titolo atteso \'{expected_title}\' (verifica fuzzy)."
                    else:
                        st.success(f"VERIFICA TITOLO (FUZZY) SUPERATA: PDF da {current_effective_url} sembra corrispondere al titolo atteso.")
                text_content = "".join([page.get_text() + "\\n\\n" for page in doc])
        except Exception as fitz_error:
            error_type_name = type(fitz_error).__name__
            st.warning(f"Errore PyMuPDF ({error_type_name}) durante parsing PDF da {current_effective_url} (originale: {pdf_url}). Dettaglio: {fitz_error}")
            return f"Errore parsing PDF ({error_type_name} da fitz): {current_effective_url} (originale: {pdf_url})"
            
        cleaned_text = "\\n".join([line for line in text_content.splitlines() if line.strip()])
        return cleaned_text if cleaned_text.strip() else None

    except requests.exceptions.Timeout as e_timeout:
        st.warning(f"Timeout ({timeout}s) durante il tentativo di scaricare/processare {pdf_url}. Dettaglio: {e_timeout}")
        return f"Timeout scaricando/processando PDF: {pdf_url}"
    except requests.exceptions.HTTPError as e_http:
        st.warning(f"Errore HTTP {e_http.response.status_code} per {pdf_url} (URL effettivo tentato: {current_effective_url}). Dettaglio: {e_http}")
        return f"Errore HTTP {e_http.response.status_code} scaricando/processando PDF: {pdf_url} (effettivo: {current_effective_url})"
    except requests.exceptions.RequestException as e_req:
        st.warning(f"Errore Richiesta generico per {pdf_url} (URL effettivo tentato: {current_effective_url}). Dettaglio: {e_req}")
        return f"Errore Richiesta scaricando/processando PDF: {pdf_url} ({e_req})"
    except Exception as e_general: 
        error_type_name = type(e_general).__name__
        st.warning(f"Errore imprevisto ({error_type_name}) durante download/parsing PDF ({pdf_url}, effettivo tentato: {current_effective_url}). Dettaglio: {e_general}")
        return f"Errore imprevisto ({error_type_name}) processando PDF: {pdf_url}"

def find_pmc_pdf_url(pmid: str, pmcid: str | None = None, timeout: int = 15) -> str | None:
    """Prova a trovare l'URL diretto di un PDF su PMC dato PMID e/o PMCID.
    Integra la ricerca di link PDF diretti sulla pagina dell'articolo PMC."""
    st.write(f"DEBUG find_pmc_pdf_url: Inizio ricerca PDF per PMID {pmid}, PMCID {pmcid}")
    session = requests.Session()
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
        'Accept-Language': 'en-US,en;q=0.5',
        'Connection': 'keep-alive',
    }
    session.headers.update(headers)

    # Tentativo 0: Accesso diretto alla pagina articolo PMC (se PMCID fornito) e ricerca link PDF specifico
    if pmcid:
        numeric_pmcid_val = pmcid.replace("PMC", "").strip()
        if numeric_pmcid_val.isdigit():
            article_url_pmc_page = f"https://www.ncbi.nlm.nih.gov/pmc/articles/PMC{numeric_pmcid_val}/"
            st.write(f"DEBUG find_pmc_pdf_url (PMCID:{numeric_pmcid_val}): Tentativo 0 - Scraping pagina {article_url_pmc_page}")
            try:
                response_page = session.get(article_url_pmc_page, headers=headers, timeout=timeout)
                response_page.raise_for_status()
                soup = BeautifulSoup(response_page.content, 'html.parser')
                
                # Cerca link con "PDF" nel testo e /pdf/ nell'href, come suggerito
                # Questa regex cerca link che hanno /pmc/articles/PMCxxxxxx/pdf/ nel loro href
                # e che contengono la parola "PDF" (case-insensitive) nel testo del link.
                pdf_link_tag = soup.find('a', href=re.compile(r'/pmc/articles/PMC\\d+/pdf/', re.IGNORECASE), string=re.compile('PDF', re.IGNORECASE))
                
                if pdf_link_tag and pdf_link_tag.get('href'):
                    pdf_href_found = pdf_link_tag['href']
                    # Costruisci l'URL assoluto. urllib.parse.urljoin è robusto.
                    # Assicurati che la base sia corretta, specialmente se l'href fosse già assoluto (urljoin lo gestisce)
                    absolute_pdf_url = urllib.parse.urljoin(response_page.url, pdf_href_found) 
                    st.success(f"DEBUG find_pmc_pdf_url (PMCID:{numeric_pmcid_val}): Trovato URL PDF da scraping pagina (metodo 0): {absolute_pdf_url}")
                    # Potremmo fare un HEAD check qui per confermare 'application/pdf', ma lo lasciamo a download_and_parse_pdf_from_url
                    return absolute_pdf_url
                else:
                    st.write(f"DEBUG find_pmc_pdf_url (PMCID:{numeric_pmcid_val}): Nessun link PDF specifico trovato con regex su {article_url_pmc_page}")

            except requests.exceptions.RequestException as e_req:
                st.warning(f"DEBUG find_pmc_pdf_url (PMCID:{numeric_pmcid_val}): Errore richiesta HTTP durante scraping pagina {article_url_pmc_page}: {e_req}")
            except Exception as e_parse: # Cattura altri errori (es. BeautifulSoup)
                st.warning(f"DEBUG find_pmc_pdf_url (PMCID:{numeric_pmcid_val}): Errore parsing pagina {article_url_pmc_page}: {e_parse}")
        else:
            st.write(f"DEBUG find_pmc_pdf_url: PMCID '{pmcid}' non valido per Tentativo 0.")


    # Tentativo 1: URL diretto del PDF su PMC usando PMCID (se pmcid ha già il prefisso PMC, va bene)
    if pmcid:
        # Assicuriamoci che pmcid sia solo il numero per costruire l'URL, anche se dovrebbe già esserlo dalla modifica precedente
        numeric_pmcid = pmcid.replace("PMC", "").strip()
        if numeric_pmcid.isdigit(): # Controllo base
            # Formato 1: /pmc/articles/PMC<pmcid>/pdf/
            pdf_url_try1 = f"https://www.ncbi.nlm.nih.gov/pmc/articles/PMC{numeric_pmcid}/pdf/"
            try:
                response = session.head(pdf_url_try1, timeout=timeout, allow_redirects=True, verify=False)
                if response.status_code == 200 and 'application/pdf' in response.headers.get('content-type', '').lower():
                    st.write(f"DEBUG find_pmc_pdf_url: Trovato PDF diretto (metodo 1): {pdf_url_try1}")
                    return pdf_url_try1
            except requests.exceptions.RequestException:
                pass # Continua con il prossimo tentativo

    # Tentativo 2: Pagina articolo PMC (usando PMCID) e cerca link PDF
    if pmcid:
        numeric_pmcid = pmcid.replace("PMC", "").strip()
        if numeric_pmcid.isdigit():
            article_pmc_url = f"https://www.ncbi.nlm.nih.gov/pmc/articles/PMC{numeric_pmcid}/"
            try:
                response = session.get(article_pmc_url, timeout=timeout, verify=False)
                response.raise_for_status()
                soup = BeautifulSoup(response.content, 'html.parser')
                # Cerca link con href che termina in .pdf o attributi specifici
                # Esempio: <a href="...article.pdf" data-format="pdf">PDF</a>
                # Esempio: link nella sezione "Formats" o "Download links"
                pdf_links = soup.find_all('a', href=True)
                for link in pdf_links:
                    href = link['href']
                    # Logica di matching più specifica per i link PDF su PMC
                    is_pdf_link_by_text = any(kw.lower() in link.get_text().lower() for kw in ['pdf', 'full text'])
                    is_pdf_link_by_attr = link.get('data-format') == 'pdf' or link.get('title', '').lower().count('pdf') > 0
                    
                    if href.lower().endswith('.pdf') and (is_pdf_link_by_text or is_pdf_link_by_attr):
                        # Assicura che l'URL sia assoluto
                        pdf_url = urllib.parse.urljoin(article_pmc_url, href)
                        # Verifica che l'URL sia ancora sotto ncbi.nlm.nih.gov per sicurezza
                        if "ncbi.nlm.nih.gov" in urllib.parse.urlparse(pdf_url).netloc:
                            st.write(f"DEBUG find_pmc_pdf_url: Trovato PDF da pagina PMC (metodo 2, PMCID {numeric_pmcid}): {pdf_url}")
                            return pdf_url
            except requests.exceptions.RequestException as e:
                st.warning(f"DEBUG find_pmc_pdf_url: Errore accesso pagina PMC {article_pmc_url}: {e}")
            except Exception as e_parse:
                 st.warning(f"DEBUG find_pmc_pdf_url: Errore parsing pagina PMC {article_pmc_url}: {e_parse}")

    # Tentativo 3: Pagina articolo PubMed (usando PMID) e cerca link a PMC o PDF
    if pmid:
        article_pubmed_url = f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/"
        try:
            response = session.get(article_pubmed_url, timeout=timeout, verify=False)
            response.raise_for_status()
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Cerca link a PMC prima, poi PDF generici
            all_links = soup.find_all('a', href=True)
            pmc_article_link_found = None
            
            for link in all_links:
                href = link['href']
                # Cerca link a un articolo PMC completo
                pmc_match = re.search(r'/pmc/articles/PMC(\d+)/?', href)
                if pmc_match:
                    # Trovato un link a un articolo PMC. Estrai il PMCID e richiama.
                    new_pmcid = pmc_match.group(1)
                    st.write(f"DEBUG find_pmc_pdf_url: Trovato link a PMC (PMCID {new_pmcid}) da pagina PubMed. Ritento con questo.")
                    # Chiamata ricorsiva con il nuovo PMCID (solo se diverso da quello già tentato)
                    if new_pmcid and new_pmcid != pmcid: # Evita loop se pmcid iniziale era già da qui
                        return find_pmc_pdf_url(pmid, new_pmcid, timeout=timeout) # Passa anche pmid per contesto
                    pmc_article_link_found = True # Segna che abbiamo trovato un link PMC, anche se non ha dato un PDF diretto
                
                # Se non abbiamo trovato un link PMC e questo link sembra un PDF
                if not pmc_article_link_found and href.lower().endswith('.pdf'):
                    link_text_lower = link.get_text().lower()
                    # Cerca keyword come "full text", "pdf" nel testo del link o nei suoi dintorni
                    if 'full text' in link_text_lower or 'pdf' in link_text_lower or 'download' in link_text_lower:
                        pdf_url = urllib.parse.urljoin(article_pubmed_url, href)
                        # Verifica che l'URL sia ancora ncbi o un dominio fidato se necessario
                        # Per ora, lo accettiamo se finisce in .pdf e ha testo rilevante
                        st.write(f"DEBUG find_pmc_pdf_url: Trovato PDF generico da pagina PubMed (metodo 3, PMID {pmid}): {pdf_url}")
                        return pdf_url
                        
        except requests.exceptions.RequestException as e:
            st.warning(f"DEBUG find_pmc_pdf_url: Errore accesso pagina PubMed {article_pubmed_url}: {e}")
        except Exception as e_parse:
            st.warning(f"DEBUG find_pmc_pdf_url: Errore parsing pagina PubMed {article_pubmed_url}: {e_parse}")
            
    st.write(f"DEBUG find_pmc_pdf_url: Nessun PDF diretto trovato per PMID {pmid} / PMCID {pmcid}.")
    return None

# NUOVA FUNZIONE PER RICERCA SEMANTICA SU INTERO STORICO CHAT
@st.cache_data(show_spinner="Ricerca nell'intero storico chat in corso...", persist="disk")
def search_entire_chat_history_semantically(_chat_sessions_hash, query_text, _embedding_model, top_n=10, similarity_threshold=0.2):
    """Cerca semanticamente in tutti i messaggi di tutte le chat."""
    if not query_text or not _chat_sessions_hash: # _chat_sessions_hash per invalidare cache se le chat cambiano
        return []

    all_messages_data = []
    for chat_id, chat_data in st.session_state.chat_sessions.items():
        chat_name = chat_data.get("name", f"Chat ID: {chat_id[:8]}")
        for msg_idx, msg in enumerate(chat_data.get("messages", [])):
            content = msg.get("content")
            role = msg.get("role")
            timestamp = msg.get("timestamp")
            if content and role: # Includi solo messaggi con contenuto e ruolo
                # Crea una data leggibile, se possibile
                msg_date_str = "N/D"
                if timestamp:
                    try: msg_date_str = datetime.fromisoformat(timestamp.replace("Z","+00:00")).strftime('%d/%m/%Y %H:%M')
                    except: pass
                
                all_messages_data.append({
                    "chat_id": chat_id,
                    "chat_name": chat_name,
                    "msg_idx": msg_idx,
                    "role": role,
                    "content": content,
                    "timestamp": timestamp,
                    "msg_date_str": msg_date_str,
                    "unique_id": f"{chat_id}_{msg_idx}" # Per possibile uso futuro, es. dedup
                })

    if not all_messages_data:
        return []

    message_contents = [msg["content"] for msg in all_messages_data]
    
    try:
        message_embeddings = _embedding_model.encode(message_contents)
        query_embedding = _embedding_model.encode([query_text])
    except Exception as e:
        st.error(f"Errore durante la creazione degli embeddings per la ricerca nello storico: {e}")
        return []

    if message_embeddings is None or query_embedding is None or message_embeddings.size == 0:
        return []

    similarities = cosine_similarity(query_embedding, message_embeddings)[0]
    
    relevant_messages_info = []
    sorted_indices = np.argsort(similarities)[::-1]

    for i in sorted_indices:
        if similarities[i] >= similarity_threshold and len(relevant_messages_info) < top_n:
            msg_data = all_messages_data[i].copy() # Copia per non modificare l'originale
            msg_data["similarity"] = float(similarities[i])
            # Per l'LLM, potremmo voler troncare il contenuto del messaggio se troppo lungo
            # msg_data["content"] = msg_data["content"][:1000] + ("..." if len(msg_data["content"]) > 1000 else "")
            relevant_messages_info.append(msg_data)
        elif len(relevant_messages_info) >= top_n:
            break # Raggiunto il numero massimo di risultati
            
    return relevant_messages_info
def save_json_data(data, file_path):
    temp_file_path = ""; dir_name = os.path.dirname(file_path)
    try:
        # Assicurati che la directory esista, specialmente per il primo salvataggio di APP_STATE_FILE
        os.makedirs(dir_name, exist_ok=True) 
        with tempfile.NamedTemporaryFile(mode='w', encoding='utf-8', delete=False, dir=dir_name, prefix=os.path.basename(file_path) + '~') as tmp_f:
            temp_file_path = tmp_f.name; json.dump(data, tmp_f, ensure_ascii=False, indent=4)
        os.replace(temp_file_path, file_path)
    except Exception as e:
        st.error(f"Errore durante il salvataggio sicuro di {os.path.basename(file_path)}: {e}")
        if temp_file_path and os.path.exists(temp_file_path):
            try: os.remove(temp_file_path)
            except Exception as re: st.warning(f"Impossibile rimuovere il file temporaneo {temp_file_path} dopo errore: {re}")
    finally: # Pulizia finale del file temporaneo se esiste e la sostituzione non è avvenuta.
        if temp_file_path and os.path.exists(temp_file_path):
            # Verifica se il file originale esiste e se il temporaneo è più recente (significa che replace potrebbe non essere avvenuto)
            # O se il file originale non esiste affatto (significa che il replace non è avvenuto)
            original_exists = os.path.exists(file_path)
            if not original_exists or (original_exists and os.path.getmtime(temp_file_path) > os.path.getmtime(file_path)):
                try: os.remove(temp_file_path)
                except Exception: pass


def load_json_data(file_path, default_data_type=dict):
    try:
        if os.path.exists(file_path):
            with open(file_path, 'r', encoding='utf-8') as f: content = f.read()
            if not content.strip(): return default_data_type() if callable(default_data_type) else default_data_type
            return json.loads(content)
    except json.JSONDecodeError as e: print(f"DEBUG: Errore decodifica JSON in {os.path.basename(file_path)}: {e}"); st.warning(f"Errore decodifica JSON {os.path.basename(file_path)}. Dettagli: {e}. Verranno usati i dati di default.")
    except Exception as e: st.warning(f"Errore caricamento {os.path.basename(file_path)}: {e}. Verranno usati i dati di default.")
    return default_data_type() if callable(default_data_type) else default_data_type

if "app_initialized" not in st.session_state:
    load_dotenv()
    app_state_loaded = load_json_data(APP_STATE_FILE, default_data_type=dict)
    
    # Inizializza tutti i valori di session_state con i default o quelli caricati
    st.session_state.gemini_api_key = app_state_loaded.get("gemini_api_key", os.getenv("GEMINI_API_KEY") or "")
    st.session_state.selected_gemini_model = app_state_loaded.get("selected_gemini_model", DEFAULT_GEMINI_MODEL)
    st.session_state.document_processed = app_state_loaded.get("document_processed", False)
    st.session_state.uploaded_file_name = app_state_loaded.get("uploaded_file_name", None)
    st.session_state.uploaded_file_content_hash = app_state_loaded.get("uploaded_file_content_hash", None) # Verrà ricalcolato se i parametri di chunking cambiano
    st.session_state.current_chat_id = app_state_loaded.get("current_chat_id", None)
    st.session_state.token_usage_log = app_state_loaded.get("token_usage_log", [])
    
    # Impostazioni di chunking e similarità
    st.session_state.chunk_size = app_state_loaded.get("chunk_size", 700)
    st.session_state.overlap = app_state_loaded.get("overlap", 150)
    st.session_state.top_n = app_state_loaded.get("top_n", DEFAULT_TOP_N)
    st.session_state.similarity_threshold = app_state_loaded.get("similarity_threshold", DEFAULT_SIMILARITY_THRESHOLD)

    # Persistenza file caricato
    st.session_state.uploaded_file_mime = app_state_loaded.get("uploaded_file_mime", None)
    uploaded_file_bytes_b64 = app_state_loaded.get("uploaded_file_bytes_b64", None)
    if uploaded_file_bytes_b64:
        try:
            st.session_state.uploaded_file_bytes = base64.b64decode(uploaded_file_bytes_b64)
        except Exception as e:
            st.warning(f"Errore nel decodificare i byte del file caricato dallo stato: {e}")
            st.session_state.uploaded_file_bytes = None
    else:
        st.session_state.uploaded_file_bytes = None

    # Caricamento chunks e embeddings (devono essere coerenti con document_processed)
    if st.session_state.document_processed:
        text_chunks_tuple = app_state_loaded.get("text_chunks_tuple", None)
        chunk_embeddings_list = app_state_loaded.get("chunk_embeddings_list", None)
        if text_chunks_tuple and chunk_embeddings_list:
            st.session_state.text_chunks = list(text_chunks_tuple)
            st.session_state.chunk_embeddings = np.array(chunk_embeddings_list)
            if not st.session_state.text_chunks or not isinstance(st.session_state.chunk_embeddings, np.ndarray) or st.session_state.chunk_embeddings.size == 0:
                st.session_state.document_processed = False # Incoerenza trovata
        else:
            st.session_state.document_processed = False # Dati mancanti
    else: # Se non è processato, assicurati che siano None
        st.session_state.text_chunks = None
        st.session_state.chunk_embeddings = None

    st.session_state.chat_sessions = load_json_data(CHATS_FILE, default_data_type=dict)
    migrated_chats_flag = False; dt_now = datetime.now(timezone.utc).isoformat()
    for chat_id, chat_data in st.session_state.chat_sessions.items():
        updated = False
        for f, v in [("updated_at", chat_data.get("created_at", dt_now)), ("created_at", dt_now)]:
            if f not in chat_data: chat_data[f] = v; updated = True
        for tf in ["tokens_prompt", "tokens_completion", "tokens_total"]:
            if tf not in chat_data: chat_data[tf] = 0; updated = True
        if updated: migrated_chats_flag = True
    if migrated_chats_flag: save_json_data(st.session_state.chat_sessions, CHATS_FILE); st.info("Dati chat aggiornati.")

    st.session_state.google_api_key = app_state_loaded.get("google_api_key", os.getenv("GOOGLE_API_KEY") or "")
    st.session_state.google_cse_id = app_state_loaded.get("google_cse_id", os.getenv("GOOGLE_CSE_ID") or "")
    st.session_state.google_creds = app_state_loaded.get("google_creds", None)
    st.session_state.app_initialized = True

    # Aggiunta per Semantic Scholar API Key
    st.session_state.semantic_scholar_api_key = app_state_loaded.get("semantic_scholar_api_key", "")

    if 'action_in_progress' not in st.session_state:
        st.session_state.action_in_progress = False
    
    # Assicuriamoci che ogni chat abbia il campo 'active_article_context' e 'active_web_page_context'
    for chat_id in st.session_state.chat_sessions:
        if 'active_article_context' not in st.session_state.chat_sessions[chat_id]:
            st.session_state.chat_sessions[chat_id]['active_article_context'] = None
        if 'active_web_page_context' not in st.session_state.chat_sessions[chat_id]:
            st.session_state.chat_sessions[chat_id]['active_web_page_context'] = None

def persist_app_data():
    # Calcola l'hash del contenuto del documento basato ANCHE sui parametri di chunking
    # così se cambiano, l'hash per la cache di get_chunks_and_embeddings_cached sarà diverso
    doc_content_for_hash = None
    # Utilizzo di .get() per accedere in sicurezza agli attributi di session_state
    if st.session_state.get("document_processed", False) and st.session_state.get("uploaded_file_name"):
        # Idealmente avremmo i bytes originali qui per l'hash,
        # ma per semplicità usiamo il nome e i parametri di chunking
        # Questo non è l'hash usato per get_chunks_and_embeddings_cached, ma per APP_STATE_FILE
        doc_content_for_hash = (
            st.session_state.get("uploaded_file_name"), 
            st.session_state.get("chunk_size", 700), # Default come nell'inizializzazione
            st.session_state.get("overlap", 150)    # Default come nell'inizializzazione
        )
    
    app_state_to_save = {
        "gemini_api_key": st.session_state.get("gemini_api_key", os.getenv("GEMINI_API_KEY") or ""),
        "selected_gemini_model": st.session_state.get("selected_gemini_model", DEFAULT_GEMINI_MODEL),
        "document_processed": st.session_state.get("document_processed", False),
        "uploaded_file_name": st.session_state.get("uploaded_file_name"), # None è un default accettabile
        "uploaded_file_content_hash": hash(doc_content_for_hash) if doc_content_for_hash else None, # Aggiornato per riflettere parametri
        "current_chat_id": st.session_state.get("current_chat_id"),
        "token_usage_log": st.session_state.get("token_usage_log", []),
        
        # Gestione sicura per text_chunks e chunk_embeddings
        "text_chunks_tuple": tuple(st.session_state.get("text_chunks", [])) if st.session_state.get("text_chunks") is not None else None,
        "chunk_embeddings_list": st.session_state.get("chunk_embeddings").tolist() if isinstance(st.session_state.get("chunk_embeddings"), np.ndarray) else None,
        
        "chunk_size": st.session_state.get("chunk_size", 700),
        "overlap": st.session_state.get("overlap", 150),
        "top_n": st.session_state.get("top_n", DEFAULT_TOP_N),
        "similarity_threshold": st.session_state.get("similarity_threshold", DEFAULT_SIMILARITY_THRESHOLD),
        
        "google_api_key": st.session_state.get("google_api_key", os.getenv("GOOGLE_API_KEY") or ""),
        "google_cse_id": st.session_state.get("google_cse_id", os.getenv("GOOGLE_CSE_ID") or ""),
        "google_creds": st.session_state.get("google_creds"),
        "uploaded_file_mime": st.session_state.get("uploaded_file_mime"),
        "uploaded_file_bytes_b64": base64.b64encode(st.session_state.get("uploaded_file_bytes")).decode('utf-8') if st.session_state.get("uploaded_file_bytes") else None,
        "semantic_scholar_api_key": st.session_state.get("semantic_scholar_api_key", "")
    }
    save_json_data(app_state_to_save, APP_STATE_FILE)
    save_json_data(st.session_state.get("chat_sessions", {}), CHATS_FILE) # Le chat ora contengono active_article_context

def get_styled_citation_label(citation_numbers_list):
    tags = []
    for num_str in citation_numbers_list:
        try:
            num = int(num_str)
            if 1 <= num <= 20: # Numeri cerchiati Unicode da U+2460 (①) a U+2473 (⑳)
                tags.append(chr(0x245F + num))
            else:
                tags.append(f"({num})") # Fallback per numeri > 20
        except ValueError:
            tags.append(f"[{num_str}]") # Se non è un numero, non dovrebbe succedere con il regex attuale ma è un fallback
    return "".join(tags)

def display_response_with_citations(text_content, cited_chunks_map):
    ensure_js_copy_script_injected() # Assicura che lo script sia presente

    # Se text_content è None, impostalo a stringa vuota per evitare errori con json.dumps e split
    current_text_content = text_content if text_content is not None else ""

    button_id = f"copy_btn_{uuid.uuid4().hex}"
    text_to_copy_js_escaped = json.dumps(current_text_content) # Usa current_text_content

    button_html = f"""
    <div style="display: flex; justify-content: flex-end; margin-bottom: 0px; margin-top: -5px; position: relative; z-index: 10;">
        <button 
            id="{button_id}" 
            onclick='copyMessageToClipboard({text_to_copy_js_escaped}, "{button_id}")' 
            title="Copia il testo del messaggio negli appunti"
            style="
                padding: 3px 8px; 
                font-size: 0.75em; 
                background-color: #f0f0f0; 
                border: 1px solid #ccc; 
                border-radius: 4px; 
                cursor: pointer;
                line-height: 1.2;
            "
            onmouseover="this.style.backgroundColor='#e0e0e0';"
            onmouseout="this.style.backgroundColor='#f0f0f0';"
        >
            📄 Copia
        </button>
    </div>
    """
    # Visualizza prima il pulsante, poi il contenuto formattato.
    # NOTA: Questo markdown è separato da quello del contenuto del messaggio.
    # Il pulsante apparirà sopra il blocco di testo del messaggio.
    st.markdown(button_html, unsafe_allow_html=True)

    # Pattern per trovare [1], [1,2], [ 1 , 2 ], [1], ecc. Le quadre sono letterali.
    # Assicuriamoci che 're' sia importato nel contesto di questo file.
    citation_pattern = re.compile(r'\\[\\s*(\\d+(?:\\s*,\\s*\\d+)*)\\s*\\]')
    parts = re.split(citation_pattern, current_text_content) # Usa current_text_content
    text_buffer = ""
    for part in parts:
        if not part: continue
        is_citation_tag = re.fullmatch(citation_pattern, part)
        if is_citation_tag:
            if text_buffer.strip(): 
                st.markdown(text_buffer, unsafe_allow_html=True)
            text_buffer = ""
            citation_numbers_str_list = [x.strip() for x in part.strip("[]").split(',')]
            
            popover_inner_content = ""
            valid_citations_for_popover = []
            for num_str_idx in citation_numbers_str_list:
                try:
                    idx_1_based = int(num_str_idx) 
                    if idx_1_based in cited_chunks_map:
                        chunk_text = cited_chunks_map[idx_1_based]
                        max_len = 700 # Lunghezza massima anteprima chunk nel popover
                        preview = chunk_text[:max_len] + ("..." if len(chunk_text) > max_len else "")
                        popover_inner_content += f"**Passaggio Citato [{idx_1_based}]**:\n```text\n{preview}\n```\n\n---\n"
                        valid_citations_for_popover.append(num_str_idx)
                    else:
                        popover_inner_content += f"*Riferimento [{idx_1_based}] non trovato nei chunk forniti.*\n---\n"
                except ValueError:
                    # Questo non dovrebbe accadere se il regex matcha correttamente solo numeri
                    popover_inner_content += f"*Riferimento non valido: {num_str_idx}.*\n---\n"

            if valid_citations_for_popover:
                label = get_styled_citation_label(valid_citations_for_popover)
                content_for_popover = popover_inner_content.strip().removesuffix("---").strip()
                if content_for_popover: # Assicurati che ci sia contenuto prima di creare il popover
                    with st.popover(label):
                        st.markdown(content_for_popover, unsafe_allow_html=True)
                else: # Se per qualche motivo non c'è contenuto (es. chunk non trovato e nessun messaggio di errore), mostra il tag originale
                    text_buffer += part 
            else: # Nessuna citazione valida nel tag, trattalo come testo normale
                text_buffer += part
    if text_buffer.strip(): 
        st.markdown(text_buffer, unsafe_allow_html=True)

# Helper per creare client di qualunque Google API usando le stesse credenziali
def get_google_service(api_name, version, silent_on_no_creds=False):
    creds_info_original = st.session_state.get("google_creds", None)
    if not creds_info_original:
        if not silent_on_no_creds:
            st.error("Servizi Google non autorizzati. Rieseguire l'autorizzazione dalla sidebar.")
        return None
    try:
        creds_info_processed = creds_info_original.copy()

        # Controllo preventivo per i campi necessari al refresh
        required_fields_for_refresh = ['refresh_token', 'token_uri', 'client_id', 'client_secret']
        missing_fields = [field for field in required_fields_for_refresh if not creds_info_processed.get(field)]
        if missing_fields:
            if not silent_on_no_creds:
                error_message = (
                    f"Credenziali Google incomplete o corrotte (campi mancanti: {', '.join(missing_fields)}). "
                    "L'autorizzazione esistente verrà resettata. "
                    "Per favore, riesegui l'autorizzazione dei Servizi Google al prossimo caricamento della pagina."
                )
                st.error(error_message)
                if 'google_creds' in st.session_state: del st.session_state['google_creds']
                persist_app_data()
                st.rerun() # Forza un rerun immediato per aggiornare la UI della sidebar
            return None

        if "expiry" in creds_info_processed and isinstance(creds_info_processed["expiry"], str):
            try:
                creds_info_processed["expiry"] = datetime.fromisoformat(creds_info_processed["expiry"].replace("Z", "+00:00"))
            except ValueError:
                # Se la conversione fallisce, potrebbe essere un formato non atteso o None
                # In tal caso, potremmo volerlo rimuovere o impostare a None per evitare errori nel costruttore Credentials
                st.warning(f"Attenzione: impossibile convertire 'expiry' ({creds_info_processed['expiry']}) in datetime. Verrà ignorato.")
                del creds_info_processed["expiry"] # O creds_info_processed["expiry"] = None
        
        creds = Credentials(**creds_info_processed)
        return build(api_name, version, credentials=creds)
    except Exception as e:
        if not silent_on_no_creds:
            st.error(f"Errore creazione servizio Google {api_name} v{version}: {e}")
        return None

# Aggiungo funzioni helper per Drive, Sheets e Calendar

def upload_file_to_drive(file_bytes, file_name, mime_type):
    drive_service = get_google_service('drive', 'v3')
    if not drive_service:
        return None
    file_metadata = {'name': file_name}
    media = MediaIoBaseUpload(io.BytesIO(file_bytes), mimetype=mime_type)
    file = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
    return file.get('id')

def create_google_sheet(title):
    sheets_service = get_google_service('sheets', 'v4')
    if not sheets_service:
        return None
    spreadsheet_body = {'properties': {'title': title}}
    spreadsheet = sheets_service.spreadsheets().create(body=spreadsheet_body, fields='spreadsheetId').execute()
    return spreadsheet.get('spreadsheetId')

def list_calendar_events(calendar_id='primary', max_results=10):
    calendar_service = get_google_service('calendar', 'v3')
    if not calendar_service:
        return None
    now = datetime.utcnow().isoformat() + 'Z'
    events_result = calendar_service.events().list(calendarId=calendar_id, timeMin=now, maxResults=max_results, singleEvents=True, orderBy='startTime').execute()
    return events_result.get('items', [])

def add_assistant_response_and_rerun_v2(assistant_content, current_chat_data, prompt_tokens=0, completion_tokens=0, citation_map=None, chunks_info=None, mark_action_completed=False):
    """Aggiunge una risposta dell'assistente, aggiorna i token (se forniti), persiste e fa rerun."""
    ai_msg_data = {
        "role": "assistant", 
        "content": assistant_content,
        "timestamp": datetime.now(timezone.utc).isoformat() # Aggiungi timestamp UTC ISO per AI
    }
    if citation_map: # Attualmente citation_map è sempre vuoto, ma lo manteniamo per coerenza
        ai_msg_data["cited_chunks_map"] = citation_map
    if chunks_info:
        ai_msg_data["chunks_info"] = chunks_info
    
    current_chat_data.setdefault("messages", []).append(ai_msg_data)
    
    # Aggiorna i token solo se sono stati forniti valori positivi (tipicamente da una chiamata LLM)
    if prompt_tokens > 0 or completion_tokens > 0:
        current_chat_data.update({
            "tokens_prompt": current_chat_data.get("tokens_prompt", 0) + prompt_tokens,
            "tokens_completion": current_chat_data.get("tokens_completion", 0) + completion_tokens,
            "tokens_total": current_chat_data.get("tokens_total", 0) + prompt_tokens + completion_tokens
        })
        # Assicurati che current_chat_id e selected_gemini_model siano accessibili o passati se necessario
        # Qui presumiamo che siano in st.session_state come al solito
        if st.session_state.get("current_chat_id") and st.session_state.get("selected_gemini_model"):
            st.session_state.setdefault("token_usage_log", []).append({
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "chat_id": st.session_state.current_chat_id,
                "model": st.session_state.selected_gemini_model,
                "prompt_tokens": prompt_tokens,
                "completion_tokens": completion_tokens,
                "total_tokens": prompt_tokens + completion_tokens
            })

    current_chat_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    
    if mark_action_completed:
        st.session_state.action_in_progress = False # Sblocca per la prossima azione
        
    persist_app_data() # Salva lo stato PRIMA del rerun
    st.rerun()

# ----- Funzioni Helper Aggiuntive -----

def parse_markdown_table_from_text(md_text):
    """Estrae una tabella da testo formattato come Markdown semplice o da un blocco ```datasheet ... ```."""
    if not md_text: return None

    # Cerca prima i blocchi ```datasheet ... ```
    datasheet_match = re.search(r"```datasheet\n(.*?)\n```", md_text, re.DOTALL)
    if datasheet_match:
        table_content = datasheet_match.group(1).strip()
        parsed_table = []
        lines = table_content.split('\n')
        if not lines: return None

        # Tenta di determinare il delimitatore (virgola o tabulazione)
        # Controlla la prima riga di dati (o l'header se è l'unica)
        sample_line_for_delimiter = lines[0] if len(lines) == 1 else lines[1] if len(lines) > 1 else lines[0]
        
        if ',' in sample_line_for_delimiter:
            delimiter = ','
        elif '\t' in sample_line_for_delimiter:
            delimiter = '\t'
        else:
            # Se non ci sono virgole o tab, prova a splittare per spazi multipli (euristica)
            # o assumi che sia una tabella Markdown con pipe già gestita dal fallback sotto
            # Per ora, se non è CSV o TSV chiaro, consideriamo ogni linea come una singola colonna
            # o affidiamoci al parser Markdown sottostante se questo blocco non è ben formattato.
            # Questa parte potrebbe essere migliorata con parsing più intelligente di spazi/allineamenti.
            delimiter = None 

        header = [h.strip() for h in lines[0].split(delimiter if delimiter else ',')] # Default a virgola se non trovato
        if not header: return None # Richiede almeno un header
        parsed_table.append(header)

        for line in lines[1:]:
            if line.strip(): # Salta righe vuote
                row_data = [cell.strip() for cell in line.split(delimiter if delimiter else ',')] # Default a virgola
                # Assicura che ogni riga abbia lo stesso numero di colonne dell'header, riempiendo con stringhe vuote se necessario
                if len(row_data) < len(header):
                    row_data.extend([''] * (len(header) - len(row_data)))
                elif len(row_data) > len(header):
                    row_data = row_data[:len(header)] # Tronca se più lunga
                parsed_table.append(row_data)
        
        return parsed_table if len(parsed_table) > 0 else None # Richiede almeno l'header

    # Fallback: parsing di tabelle Markdown con | come prima
    table = []
    lines = md_text.strip().split('\n')
    header_found = False
    skipped_separator = False

    for line in lines:
        line = line.strip()
        if line.startswith("|") and line.endswith("|"):
            potential_separator_chars = line[1:-1].replace(" ", "") 
            is_separator_line = all(c in '-|:' for c in potential_separator_chars) and '-' in potential_separator_chars

            if header_found and is_separator_line and not skipped_separator:
                skipped_separator = True
                continue
            
            row = [cell.strip() for cell in line[1:-1].split("|")]
            if row: 
                if not header_found and not is_separator_line:
                    header_found = True
                table.append(row)
        elif table and not line.startswith("|"): # Se abbiamo iniziato una tabella con | e la linea non è più formattata così, interrompi
                        break
                
    return table if header_found and len(table) > (1 if skipped_separator else 0) else None

def write_data_to_google_sheet(spreadsheet_id, data_values, start_cell="A1"):
    """Scrive dati (lista di liste) in un Google Sheet specificato."""
    sheets_service = get_google_service('sheets', 'v4')
    if not sheets_service or not data_values or not isinstance(data_values, list):
        st.warning(f"Servizio Sheets non disponibile o dati non validi per la scrittura. Dati: {type(data_values)}")
        return False
    try:
        spreadsheet_metadata = sheets_service.spreadsheets().get(spreadsheetId=spreadsheet_id).execute()
        sheets = spreadsheet_metadata.get('sheets', '')
        if not sheets:
            st.error("Impossibile trovare fogli nello spreadsheet per la scrittura.")
            return False
        # Usa il titolo del primo foglio per costruire il range
        sheet_title = sheets[0].get('properties', {}).get('title', 'Foglio1') 
        
        range_name_for_update = f"{sheet_title}!{start_cell}"
        
        body = {
            'values': data_values
        }
        result = sheets_service.spreadsheets().values().update(
            spreadsheetId=spreadsheet_id,
            range=range_name_for_update,
            valueInputOption="USER_ENTERED",
            body=body
        ).execute()
        return result.get('updatedCells', 0) > 0
    except Exception as e:
        st.error(f"Errore durante la scrittura dei dati su Google Sheet: {e}")
        return False

with st.sidebar:
    st.title("📚 Super Tutor AI (Gemini Ed.)")
    st.caption("Il tuo assistente di studio potenziato con Gemini.")
    sidebar_update_time_placeholder = st.empty()
    with st.expander("🔑 Configurazione API e Modello", expanded=True):
        st.session_state.gemini_api_key = st.text_input("Google Gemini API Key", type="password", value=st.session_state.gemini_api_key, on_change=persist_app_data)
        
        model_opts = list(GEMINI_MODELS_INFO.keys())
        try: cur_idx = model_opts.index(st.session_state.selected_gemini_model)
        except ValueError: cur_idx = model_opts.index(DEFAULT_GEMINI_MODEL) if DEFAULT_GEMINI_MODEL in model_opts else 0
        sel_model = st.selectbox("Modello LLM Gemini:", model_opts, index=cur_idx, format_func=lambda x: f"{x} ({GEMINI_MODELS_INFO.get(x, 'Info non disp.')})", on_change=persist_app_data)
        if sel_model != st.session_state.selected_gemini_model: st.session_state.selected_gemini_model = sel_model # on_change gestisce persist

        st.markdown("---"); st.subheader("Configurazione Ricerca Google")
        st.session_state.google_api_key = st.text_input("Google API Key (Custom Search)", type="password", value=st.session_state.google_api_key, on_change=persist_app_data)
        st.session_state.google_cse_id = st.text_input("Google Search Engine ID (CX)", value=st.session_state.google_cse_id, on_change=persist_app_data)

    with st.expander("🔑 Configurazione API Scientifiche", expanded=False):
        st.session_state.semantic_scholar_api_key = st.text_input(
            "Semantic Scholar API Key (Opzionale)", 
            type="password", 
            value=st.session_state.semantic_scholar_api_key, 
            on_change=persist_app_data,
            help="Inserisci la tua API key di Semantic Scholar per limiti di richieste più alti (aiuta con errore 429)."
        )

    with st.expander("📄 Documento Allegato", expanded=True):
        # Chiave unica per file_uploader per permettere il reset quando le impostazioni di chunking cambiano
        # Non è il modo ideale, idealmente si invalida la cache di get_chunks_and_embeddings_cached
        uploader_key = f"file_uploader_{st.session_state.get('uploaded_file_content_hash', 'new')}"
        
        accepted_file_types = [
            "txt", "md", "rtf", "html", "htm", "xml", "json", "csv", 
            "py", "js", "java", "c", "cpp", "cs", "go", "rb", "php", "swift", "kt", "ts", "sql", "sh", "bat", "ps1", # Codice e Script
            "pdf", 
            "png", "jpg", "jpeg", "webp", "heic", "heif", # Immagini
            "docx", "xlsx", "pptx" # Nuovi formati Office
        ]
        up_file = st.file_uploader(
            label="Carica file (testo, codice, PDF, immagini)",
            type=accepted_file_types, 
            key=uploader_key
        )
        
        if up_file:
            file_bytes = up_file.getvalue()
            file_type = up_file.type # Questo è il MIME type
            # Aggiorna subito uploaded_file_name in session_state così è disponibile in get_chunks_and_embeddings_cached
            st.session_state.uploaded_file_name = up_file.name 
            st.session_state.uploaded_file_bytes = file_bytes
            st.session_state.uploaded_file_mime = file_type
            
            current_doc_processing_hash = hash((file_bytes, st.session_state.chunk_size, st.session_state.overlap))

            # Confronta con un hash che traccia lo stato processato del documento nella UI
            # Se l'hash dei parametri di processo (current_doc_processing_hash) è diverso da quello salvato
            # OPPURE se il nome del file è diverso, allora riprocessa.
            if not st.session_state.get("document_processed", False) or \
               st.session_state.uploaded_file_content_hash != current_doc_processing_hash or \
               st.session_state.uploaded_file_name != up_file.name: # Controllo ridondante sul nome ma sicuro
                
                with st.spinner(f"Elaborazione di '{up_file.name}'..."):
                    # Passiamo file_bytes e file_type (MIME)
                    # uploaded_file_name è già in session_state per essere usato da get_chunks_and_embeddings_cached
                    chunks, embeds = get_chunks_and_embeddings_cached(
                        current_doc_processing_hash, 
                        file_bytes, # Passiamo i byte grezzi
                        file_type,  # Passiamo il MIME type
                        st.session_state.chunk_size, 
                        st.session_state.overlap
                    )
                    if chunks is not None and embeds is not None: # Accetta anche chunks/embeds vuoti (es. per immagini)
                        st.session_state.update({
                            "text_chunks": chunks, 
                            "chunk_embeddings": embeds, 
                            "document_processed": True, # Il file è stato caricato/processato (anche se no RAG testuale per immagini)
                            # "uploaded_file_name": up_file.name, # Già settato prima della chiamata
                            "uploaded_file_content_hash": current_doc_processing_hash
                        })
                        # Il messaggio di successo/info ora è gestito DENTRO get_chunks_and_embeddings_cached
                        # per differenziare tra estrazione testuale e caricamento immagine.
                    else: 
                        st.error(f"Errore durante l'elaborazione di '{up_file.name}'. Il file potrebbe non essere supportato o corrotto.")
                        st.session_state.document_processed = False # Fallito il processamento
                    persist_app_data() # Salva lo stato aggiornato
        elif st.session_state.document_processed and st.session_state.uploaded_file_name:
            st.caption(f"Doc: {st.session_state.uploaded_file_name}")
            if st.button("Rimuovi documento", use_container_width=True):
                st.session_state.update({"document_processed": False, "text_chunks": None, "chunk_embeddings": None, 
                                         "uploaded_file_name": None, "uploaded_file_content_hash": None,
                                         "uploaded_file_bytes": None, "uploaded_file_mime": None}) # Aggiunto reset byte e mime
                st.info("Documento rimosso."); persist_app_data(); st.rerun()
    
    with st.expander("⚙️ Parametri di Analisi Documento", expanded=False):
        # Callback per invalidare l'hash del documento quando queste impostazioni cambiano
        def settings_changed_callback():
            st.session_state.uploaded_file_content_hash = None # Invalida l'hash per forzare il riprocessamento
            persist_app_data() 

        st.slider("Dimensione Chunk (parole)", 100, 2000, st.session_state.get("chunk_size", 700), 50, key="chunk_size", on_change=settings_changed_callback)
        st.slider("Sovrapposizione Chunk (parole)", 0, st.session_state.chunk_size // 2 if st.session_state.chunk_size > 0 else 0 , st.session_state.get("overlap", 150), 10, key="overlap", on_change=settings_changed_callback)
        st.slider("Max Chunk Rilevanti (Top N)", 1, 10, st.session_state.get("top_n", DEFAULT_TOP_N), 1, key="top_n", on_change=persist_app_data) # Non richiede riprocessamento documento
        st.slider("Soglia Similarità Minima", 0.0, 1.0, st.session_state.get("similarity_threshold", DEFAULT_SIMILARITY_THRESHOLD), 0.01, key="similarity_threshold", on_change=persist_app_data) # Non richiede riprocessamento

    st.markdown("---"); st.subheader("📜 Storico Chat")
    search_q_chat = st.text_input("Cerca chat:", key="chat_search_q_v2", placeholder="Cerca nei messaggi...") # Nuova chiave e placeholder
    search_q_chat_low = search_q_chat.lower().strip() if search_q_chat else ""

    if st.button("➕ Nuova Chat", use_container_width=True, type="primary"):
        chat_id = str(uuid.uuid4()); ts = datetime.now(timezone.utc).isoformat()
        st.session_state.chat_sessions[chat_id] = {"name": "Nuova Chat...", "messages": [], "created_at": ts, "updated_at": ts, "tokens_prompt":0, "tokens_completion":0, "tokens_total":0, "active_article_context": None, "active_web_page_context": None}
        st.session_state.current_chat_id = chat_id
        if search_q_chat: st.session_state.chat_search_q_v2 = "" # Resetta la ricerca
        persist_app_data(); st.rerun()

    if search_q_chat_low:
        found_messages_previews = []
        
        # Ordina le chat per data più recente prima di cercarvi dentro
        sorted_chat_ids_for_search = sorted(
            st.session_state.chat_sessions.keys(),
            key=lambda x: st.session_state.chat_sessions[x].get("updated_at", "1970"),
            reverse=True
        )

        for chat_id in sorted_chat_ids_for_search:
            chat_data = st.session_state.chat_sessions[chat_id]
            chat_name = chat_data.get("name", "Chat Sconosciuta")
            for msg_idx, msg in enumerate(chat_data.get("messages", [])):
                content = msg.get("content", "").lower()
                original_content = msg.get("content", "") # Contenuto originale per l'anteprima
                
                if search_q_chat_low in content and original_content.strip(): # Assicurati che ci sia contenuto originale
                    start_index = content.find(search_q_chat_low)
                    end_index = start_index + len(search_q_chat_low)
                    
                    preview_context_chars = 35 # Caratteri prima e dopo
                    
                    prev_start = max(0, start_index - preview_context_chars)
                    prev_end = min(len(original_content), end_index + preview_context_chars)
                    
                    preview_parts = []
                    if prev_start > 0: preview_parts.append("...")
                    preview_parts.append(original_content[prev_start:start_index])
                    preview_parts.append(f"**{original_content[start_index:end_index]}**") # Evidenzia
                    preview_parts.append(original_content[end_index:prev_end])
                    if prev_end < len(original_content): preview_parts.append("...")
                    
                    preview_text = "".join(preview_parts)
                    
                    found_messages_previews.append({
                        "chat_id": chat_id,
                        "msg_idx": msg_idx,
                        "preview_text": preview_text,
                        "chat_name": chat_name,
                        "msg_timestamp": msg.get("timestamp")
                    })
        
        if found_messages_previews:
            st.markdown("---") 
            st.caption(f"Risultati ({len(found_messages_previews)}):")
            container_height = min(len(found_messages_previews) * 90, 360) # Altezza per item leggermente aumentata, max 360px
            with st.container(height=container_height): # Richiede Streamlit 1.29+
                for res_idx, res in enumerate(found_messages_previews[:15]): # Limita a N risultati mostrati (es. 15)
                    btn_key = f"search_res_btn_{res['chat_id']}_{res['msg_idx']}_{res_idx}" 
                    
                    ts_str = ""
                    if res['msg_timestamp']:
                        try: ts_str = datetime.fromisoformat(res['msg_timestamp'].replace("Z","+00:00")).strftime('%d/%m %H:%M')
                        except: ts_str = "Data N/D"
                    
                    button_label = f"Chat: '{res['chat_name']}' ({ts_str})\n" \
                                   f"_{res['preview_text']}_"

                    if st.button(button_label, key=btn_key, use_container_width=True):
                        st.session_state.current_chat_id = res['chat_id']
                        st.session_state.chat_search_q_v2 = "" # Pulisce la ricerca dopo il click
                        persist_app_data() 
                        st.rerun()
        elif search_q_chat_low: 
            st.markdown("---")
            st.caption("Nessun messaggio trovato per la tua ricerca.")
        # Non mostrare l'elenco completo delle chat se c'è una ricerca attiva.
    
    else: # search_q_chat_low è vuoto, mostra l'elenco normale delle chat
        st.markdown("---") 
        all_chat_ids = list(st.session_state.chat_sessions.keys())
        if not all_chat_ids:
            st.caption("Nessuna chat ancora.")
        else:
            sorted_ids_display = sorted(all_chat_ids, key=lambda x: st.session_state.chat_sessions[x].get("updated_at", "1970"), reverse=True)
            for sid_display in sorted_ids_display:
                sdata_display = st.session_state.chat_sessions.get(sid_display)
                if not sdata_display: continue
                sname_display = sdata_display.get("name", "Chat")
                ts_iso_display = sdata_display.get("updated_at")
                try: ts_str_display = datetime.fromisoformat(ts_iso_display.replace("Z","+00:00")).strftime('%d/%m %H:%M') if ts_iso_display else "N/D"
                except: ts_str_display = "N/D"

                cols_chat_list = st.columns([0.85, 0.15])
                with cols_chat_list[0]:
                    btn_key_select = f"sel_chat_sidebar_{sid_display}"
                    if st.button(f"{sname_display} ({ts_str_display})", key=btn_key_select, use_container_width=True, type="primary" if st.session_state.current_chat_id == sid_display else "secondary"):
                        if st.session_state.current_chat_id != sid_display:
                            st.session_state.current_chat_id = sid_display
                            persist_app_data()
                            st.rerun()
                with cols_chat_list[1]:
                    btn_key_delete = f"del_chat_sidebar_{sid_display}"
                    if st.button("🗑️", key=btn_key_delete, help="Elimina chat"):
                        if sid_display in st.session_state.chat_sessions: del st.session_state.chat_sessions[sid_display]
                        if st.session_state.current_chat_id == sid_display:
                            remaining_ids = list(st.session_state.chat_sessions.keys())
                            st.session_state.current_chat_id = sorted(remaining_ids, key=lambda x: st.session_state.chat_sessions[x].get("updated_at","1970"), reverse=True)[0] if remaining_ids else None
                        persist_app_data(); st.rerun()
    st.markdown("---"); st.subheader("🔍 Ricerca Google")
    g_query = st.text_input("Cerca su Google:", key="g_search_in")
    if st.button("Cerca", key="g_search_btn", use_container_width=True) and g_query:
        if st.session_state.google_api_key and st.session_state.google_cse_id:
            st.session_state.last_google_search_results = search_google(g_query, st.session_state.google_api_key, st.session_state.google_cse_id)
            st.session_state.last_google_query = g_query
        else: st.warning("Configura API Google Search e ID Motore.")
    if "last_google_search_results" in st.session_state and st.session_state.last_google_search_results is not None:
        st.markdown(f"##### Risultati per: *\"{st.session_state.get('last_google_query', '')}\"*")
        if not st.session_state.last_google_search_results: st.caption("Nessun risultato Google.")
        for i, item in enumerate(st.session_state.last_google_search_results):
            with st.expander(f"**{i+1}. {item.get('title', 'N/A')}**", expanded=i<2):
                st.markdown(f"[{item.get('link')}]({item.get('link')})\n\n{item.get('snippet', 'N/D')}")
                
                # NUOVO: Bottone "Analizza con AI" per risultati Google
                button_key_gs = f"analyze_gs_{item.get('link', '')}_{i}" # Assicurati che la chiave sia univoca
                if st.button("🔬 Analizza con AI", key=button_key_gs, use_container_width=True, help="Estrai contenuto, riassumi con AI e imposta come contesto chat."):
                    if not st.session_state.get('action_in_progress', False): # Controllo più esplicito
                        st.session_state.action_in_progress = True
                        url_to_analyze = item.get('link')
                        item_title = item.get('title', 'Pagina Web')
                        p_tok_analysis, c_tok_analysis = 0, 0 # Inizializza token per questa azione

                        if url_to_analyze and st.session_state.current_chat_id and st.session_state.current_chat_id in st.session_state.chat_sessions:
                            current_chat_data_for_analysis = st.session_state.chat_sessions[st.session_state.current_chat_id]
                            analysis_output_message = ""
                            
                            with st.spinner(f"Estrazione contenuto da: {url_to_analyze[:70]}..."):
                                extracted_text = fetch_and_extract_text(url_to_analyze)

                            if extracted_text and extracted_text.strip():
                                prompt_per_analisi = (
                                    f"Il seguente testo è stato estratto dalla pagina web intitolata '{item_title}' ({url_to_analyze}). Per favore:\n"
                                    "1. Fornisci un riassunto conciso (circa 100-150 parole) del contenuto principale.\n"
                                    "2. Elenca e spiega brevemente i 3-5 punti chiave o argomenti trattati nel testo.\n"
                                    "3. Se il testo sembra essere un articolo di notizie o un post di un blog, indica la data di pubblicazione se riesci a identificarla chiaramente nel testo fornito.\n"
                                    "Rispondi in modo chiaro e strutturato, usando Markdown per la formattazione.\n\n"
                                    "--- INIZIO TESTO ESTRATTO ---\n"
                                    f"{extracted_text[:15000]}"  # Limite generoso per l'analisi iniziale
                                    "\n--- FINE TESTO ESTRATTO ---"
                                )
                                
                                user_action_msg_content = f"Richiesta analisi AI per la pagina: '{item_title}' ({url_to_analyze})"
                                user_action_msg = {
                                    "role": "user", 
                                    "content": user_action_msg_content,
                                    "timestamp": datetime.now(timezone.utc).isoformat()
                                }
                                current_chat_data_for_analysis.setdefault("messages",[]).append(user_action_msg)
                                
                                # Prepara la history per l'LLM (escludendo l'ultimo messaggio utente appena aggiunto se non si vuole che si auto-risponda)
                                # In questo caso, l'LLM deve rispondere al prompt_per_analisi, non a user_action_msg_content
                                chat_history_for_llm_analysis = current_chat_data_for_analysis.get("messages", [])[:-1] 

                                with st.spinner(f"L'AI sta analizzando '{item_title}'..."):
                                    ai_summary_text, p_tok_analysis, c_tok_analysis, _ = get_answer_from_llm(
                                        query=prompt_per_analisi, # Questa è la query effettiva per l'LLM
                                        relevant_chunks_data=[],
                                        api_key=st.session_state.gemini_api_key,
                                        model_name=st.session_state.selected_gemini_model,
                                        chat_history=chat_history_for_llm_analysis, 
                                        ignore_active_contexts_for_this_query=True, 
                                        direct_file_text_context=None, 
                                        image_data=None, image_mime_type=None,
                                        external_web_search_results=None,
                                        scraped_page_content_for_query=None # Non passiamo il testo qui perché è già nel prompt_per_analisi
                                    )
                                
                                analysis_output_message = ai_summary_text
                                current_chat_data_for_analysis['active_web_page_context'] = {
                                    "url": url_to_analyze,
                                    "title": item_title,
                                    "source": "Risultato Google Analizzato",
                                    "analyzed_at": datetime.now(timezone.utc).isoformat(),
                                }
                                current_chat_data_for_analysis.pop('active_article_context', None)
                                st.success(f"'{item_title}' analizzato e impostato come contesto per la chat corrente!")

                            elif extracted_text is None or not extracted_text.strip():
                                analysis_output_message = f"❌ Impossibile estrarre contenuto testuale significativo da {url_to_analyze} per l'analisi."
                                st.error(analysis_output_message)
                            else: 
                                analysis_output_message = f"⚠️ Errore durante il tentativo di estrazione da {url_to_analyze}: {extracted_text[:200]}..."
                                st.warning(analysis_output_message)
                            
                            add_assistant_response_and_rerun_v2(
                                assistant_content=analysis_output_message, 
                                current_chat_data=current_chat_data_for_analysis,
                                prompt_tokens=p_tok_analysis, 
                                completion_tokens=c_tok_analysis,
                                mark_action_completed=True
                            )
                        else:
                            st.warning("Nessuna chat attiva o URL non valido per l'analisi.")
                            st.session_state.action_in_progress = False
                    else:
                        st.caption("Azione precedente in corso...")

    # ----- Autorizzazione Servizi Google (Docs, Sheets, Drive, Calendar, Vision) -----
    st.markdown("---")
    st.subheader("🔑 Autorizza Servizi Google")
    if not st.session_state.get("google_creds"):
        if st.button("Autorizza Servizi Google"):
            flow = InstalledAppFlow.from_client_config(CLIENT_CONFIG, SCOPES)
            creds = flow.run_local_server(port=8502, prompt='consent', access_type='offline')
            
            st.session_state.google_creds = {
                "token": creds.token,
                "refresh_token": creds.refresh_token, 
                "token_uri": creds.token_uri,         
                "client_id": creds.client_id,       
                "client_secret": creds.client_secret, 
                "scopes": creds.scopes,
                "expiry": creds.expiry.isoformat() if creds.expiry else None 
            }
            persist_app_data()
            st.success("Autenticazione Servizi Google completata! L'app si ricaricherà.")
            st.rerun()
    else:
        st.success("Servizi Google già autorizzati")
        # Bottone per la revoca manuale dell'autorizzazione
        if st.button("Revoca Autorizzazione Google", key="revoke_google_auth"):
            if 'google_creds' in st.session_state: del st.session_state['google_creds']
            persist_app_data()
            st.info("Autorizzazione Google revocata. L'app si ricaricherà.")
            st.rerun()

    st.sidebar.metric(label="Token Usati (Chat Attuale)", value=f"{st.session_state.chat_sessions[st.session_state.current_chat_id].get('tokens_total',0) if st.session_state.get('current_chat_id') and st.session_state.current_chat_id in st.session_state.get('chat_sessions', {}) else 0:,}")

    # ----- Ricerca Scientifica Specializzata -----
    st.markdown("---"); st.subheader("🔬 Ricerca Scientifica")

    # Sottosezione Semantic Scholar
    with st.expander("🎓 Semantic Scholar Search", expanded=False):
        ss_query = st.text_input("Cerca su Semantic Scholar:", key="ss_search_in", help="Inserisci termini di ricerca per articoli scientifici.")
        if st.button("Cerca Semantic Scholar", key="ss_search_btn", use_container_width=True):
            if ss_query.strip():
                with st.spinner("Ricerca Semantic Scholar in corso..."):
                    st.session_state.last_semantic_scholar_results = search_semantic_scholar(ss_query, api_key=st.session_state.semantic_scholar_api_key)
                st.session_state.last_semantic_scholar_query = ss_query
                st.rerun()
            else:
                st.caption("Inserisci un termine di ricerca.")
        st.caption("Nota: Semantic Scholar potrebbe limitare il n. di richieste (errore 429). Se accade, attendi prima di riprovare.") # AVVISO AGGIUNTO
        
        if "last_semantic_scholar_results" in st.session_state and st.session_state.last_semantic_scholar_results is not None:
            query_display = st.session_state.get('last_semantic_scholar_query', '')
            st.markdown(f"##### Risultati per: *\"{query_display}\"*")
            if not st.session_state.last_semantic_scholar_results:
                st.caption(f"Nessun risultato da Semantic Scholar per '{query_display}'.")
            else:
                for i, item in enumerate(st.session_state.last_semantic_scholar_results):
                    exp_title = f"**{i+1}. {item.get('title', 'N/A')}** ({item.get('year', 'N/D')})"
                    # Rimuoviamo l'st.expander annidato qui. Mostriamo il titolo come markdown.
                    st.markdown(f"--- \n {exp_title}") # Aggiungo un separatore e mostro il titolo
                    
                    # I dettagli ora vengono mostrati direttamente, non dentro un expander
                    st.markdown(f"**Autori:** {item.get('authors', 'N/D')}")
                    st.markdown(f"**Venue:** {item.get('venue', 'N/D')}")
                    
                    abstract_text = item.get('abstract', 'N/D')
                    max_abstract_len = 300
                    display_abstract = (abstract_text[:max_abstract_len] + '...' if len(abstract_text) > max_abstract_len else abstract_text)
                    if display_abstract.strip(): 
                        st.markdown(f"**Abstract:** {display_abstract}")
                    else: 
                        st.caption("Abstract non disponibile.")
                    
                    ss_url = item.get('url', '#')
                    pdf_url = item.get('pdf_url')
                    
                    links_md_parts = [f"[Pagina Semantic Scholar]({ss_url})"]
                    if pdf_url: 
                        links_md_parts.append(f"[PDF Open Access]({pdf_url})")
                    st.markdown(" | ".join(links_md_parts))

                        # Bottone Analizza con AI per Semantic Scholar
                    button_key_ss = f"analyze_ss_{item.get('url', '')}_{i}"  # Chiave univoca
                    if st.button("🔬 Analizza con AI", key=button_key_ss, use_container_width=True, help="Analizza abstract e tenta di scaricare e analizzare il PDF completo, se un URL Open Access è disponibile."):
                            if not st.session_state.get('action_in_progress', False):
                                st.session_state.action_in_progress = True
                                p_tok_ss, c_tok_ss = 0, 0  # Token per questa azione specifica
                                analysis_output_parts_ss = [] # Parti del messaggio di output

                                if st.session_state.current_chat_id and st.session_state.current_chat_id in st.session_state.chat_sessions:
                                    current_chat_data_ss = st.session_state.chat_sessions[st.session_state.current_chat_id]
                                    
                                    item_title_ss = item.get('title', 'N/A')
                                    item_abstract_ss = item.get('abstract', 'N/A')
                                    item_url_ss = item.get('url', '#') # URL della pagina Semantic Scholar
                                    item_pdf_url_ss = item.get('pdf_url') # URL diretto al PDF, se disponibile

                                    user_action_msg_content_ss = f"Richiesta analisi AI per l'articolo Semantic Scholar: '{item_title_ss}' (URL: {item_url_ss})"
                                    user_action_msg_ss = {"role": "user", "content": user_action_msg_content_ss, "timestamp": datetime.now(timezone.utc).isoformat()}
                                    current_chat_data_ss.setdefault("messages",[]).append(user_action_msg_ss)
                                    
                                    extracted_pdf_text_ss = None
                                    pdf_download_error_detail_ss = ""
                                    attempted_pdf_url_info_ss = "Nessun URL PDF diretto fornito da Semantic Scholar."

                                    if item_pdf_url_ss:
                                        attempted_pdf_url_info_ss = f"URL PDF da Semantic Scholar: {item_pdf_url_ss}"
                                        with st.spinner(f"Download e parsing PDF da: {item_pdf_url_ss}..."):
                                            extracted_pdf_text_ss = download_and_parse_pdf_from_url(item_pdf_url_ss, expected_title=item_title_ss)
                                        
                                        if extracted_pdf_text_ss and not extracted_pdf_text_ss.startswith("Errore") and not extracted_pdf_text_ss.startswith("Timeout"):
                                            st.success(f"PDF per '{item_title_ss}' scaricato e analizzato con successo da: {item_pdf_url_ss}")
                                        else:
                                            pdf_download_error_detail_ss = extracted_pdf_text_ss if extracted_pdf_text_ss else "Nessun testo estratto o errore sconosciuto."
                                            st.warning(f"Non è stato possibile scaricare/analizzare il PDF per '{item_title_ss}' da {item_pdf_url_ss}. L'analisi AI si baserà sull'abstract. (Dettaglio: {pdf_download_error_detail_ss})")
                                            extracted_pdf_text_ss = None # Assicura che sia None se fallisce
                                    else:
                                        st.info(f"Nessun URL PDF Open Access diretto fornito da Semantic Scholar per '{item_title_ss}'. L'analisi si baserà sull'abstract.")
                                        pdf_download_error_detail_ss = "Nessun URL PDF disponibile."

                                    # Costruzione del prompt per l'LLM
                                    prompt_per_analisi_ss_parts = [
                                        f"Sto analizzando l'articolo scientifico proveniente da Semantic Scholar intitolato: '{item_title_ss}'.",
                                        f"URL Pagina Semantic Scholar: {item_url_ss}",
                                    ]
                                    abstract_section_ss = "Abstract fornito:\n"
                                    if item_abstract_ss and item_abstract_ss != 'Abstract non disponibile.':
                                        abstract_section_ss += str(item_abstract_ss)
                                    else:
                                        abstract_section_ss += "N/A"
                                    prompt_per_analisi_ss_parts.append(abstract_section_ss)

                                    if extracted_pdf_text_ss:
                                        prompt_per_analisi_ss_parts.append(f"\\nIn aggiunta, è stato possibile estrarre il seguente testo dal PDF completo dell'articolo (URL tentato: {attempted_pdf_url_info_ss}):\\n---\\n{extracted_pdf_text_ss[:10000]}\\n---")
                                    else:
                                        prompt_per_analisi_ss_parts.append(f"\\nNon è stato possibile analizzare il PDF completo dell'articolo. {attempted_pdf_url_info_ss}. Dettaglio errore PDF: {pdf_download_error_detail_ss}")
                                    
                                    prompt_per_analisi_ss_parts.extend([
                                        "\\nPer favore, esegui le seguenti operazioni basandoti su TUTTE le informazioni fornite (abstract e, se disponibile, il testo del PDF):",
                                        "1. Fornisci un riassunto conciso ma completo (circa 150-250 parole) degli obiettivi, metodi principali, risultati chiave e conclusioni dell'articolo.",
                                        "2. Elenca e spiega brevemente i 3-5 punti o scoperte più significativi dell'articolo.",
                                        "3. Se applicabile e desumibile dal testo, discuti brevemente le implicazioni o la rilevanza di questi risultati.",
                                        "4. Se il testo del PDF è stato analizzato, integra le informazioni da esso per arricchire i punti precedenti, specificando se un dettaglio proviene dal PDF.",
                                        "Rispondi in modo chiaro e strutturato, usando Markdown per la formattazione."
                                    ])
                                    prompt_final_ss = "\\n\\n".join(prompt_per_analisi_ss_parts)

                                    chat_history_for_llm_ss = current_chat_data_ss.get("messages", [])[:-1]
                                    
                                    ai_analysis_text_ss = ""
                                    with st.spinner(f"L'AI sta analizzando l'articolo '{item_title_ss}' (Semantic Scholar)..."):
                                        ai_analysis_text_ss, p_tok_ss, c_tok_ss, _ = get_answer_from_llm(
                                            query=prompt_final_ss,
                                            relevant_chunks_data=[],
                                            api_key=st.session_state.gemini_api_key,
                                            model_name=st.session_state.selected_gemini_model,
                                            chat_history=chat_history_for_llm_ss, 
                                            ignore_active_contexts_for_this_query=True, 
                                        )
                                    analysis_output_parts_ss.insert(0, ai_analysis_text_ss)
                                    
                                    current_chat_data_ss['active_article_context'] = {
                                        "title": item_title_ss,
                                        "abstract": item_abstract_ss,
                                        "url": item_url_ss,
                                        "source": "Semantic Scholar (con tentativo analisi PDF)",
                                        "full_text_analyzed": bool(extracted_pdf_text_ss)
                                    }
                                    current_chat_data_ss.pop('active_web_page_context', None)
                                    st.success(f"Articolo '{item_title_ss}' (Semantic Scholar) analizzato e impostato come contesto per la chat corrente!")
                                else:
                                    analysis_output_parts_ss.append("Errore: Nessuna chat attiva per associare l'articolo Semantic Scholar.")
                                    st.warning(analysis_output_parts_ss[0])

                                final_output_for_chat_ss = "\\n\\n".join(analysis_output_parts_ss)
                                add_assistant_response_and_rerun_v2(
                                    assistant_content=final_output_for_chat_ss, 
                                    current_chat_data=current_chat_data_ss,
                                    prompt_tokens=p_tok_ss, 
                                    completion_tokens=c_tok_ss,
                                    mark_action_completed=True
                                )
                            else:
                                analysis_output_parts_ss.append("Errore: Nessuna chat attiva per associare l'articolo Semantic Scholar.")
                                st.warning(analysis_output_parts_ss[0])

                            final_output_for_chat_ss = "\\n\\n".join(analysis_output_parts_ss)
                            add_assistant_response_and_rerun_v2(
                                assistant_content=final_output_for_chat_ss, 
                                current_chat_data=current_chat_data_ss,
                                prompt_tokens=p_tok_ss, 
                                completion_tokens=c_tok_ss,
                                mark_action_completed=True
                            )
                    else:
                            st.caption("Azione precedente in corso...")
                    # Aggiungiamo un piccolo spazio dopo ogni articolo se non è l'ultimo
                    if i < len(st.session_state.last_semantic_scholar_results) - 1:
                        st.markdown("---")

    # Sottosezione PubMed
    with st.expander("🧬 PubMed Search", expanded=False):
        pm_query = st.text_input("Cerca su PubMed:", key="pm_search_in", help="Inserisci termini di ricerca per articoli biomedici (es. autori, parole chiave).")
        if st.button("Cerca PubMed", key="pm_search_btn", use_container_width=True):
            if pm_query.strip():
                with st.spinner("Ricerca PubMed in corso..."):
                    st.session_state.last_pubmed_results = search_pubmed(pm_query)
                st.session_state.last_pubmed_query = pm_query
                st.rerun()
            else:
                st.caption("Inserisci un termine di ricerca.")
        
        if "last_pubmed_results" in st.session_state and st.session_state.last_pubmed_results is not None:
            query_display = st.session_state.get('last_pubmed_query', '')
            st.markdown(f"##### Risultati per: *\"{query_display}\"*")
            if not st.session_state.last_pubmed_results:
                st.caption(f"Nessun risultato da PubMed per '{query_display}'.")
            else:
                for i, item in enumerate(st.session_state.last_pubmed_results):
                    st.markdown(f"**{i+1}. {item.get('title', 'N/A')}** ({item.get('year', 'N/D')})")
                    st.markdown(f"**Autori:** {item.get('authors', 'N/D')}")
                    st.markdown(f"**Journal/Venue:** {item.get('venue', 'N/D')}")
                    if item.get('pmid'): st.markdown(f"**PMID:** {item.get('pmid')}")
                    
                    abstract_text = item.get('abstract', 'N/D')
                    max_abstract_len = 250 # Leggermente ridotto per non occupare troppo spazio
                    display_abstract = (abstract_text[:max_abstract_len] + '...' if len(abstract_text) > max_abstract_len else abstract_text)
                    if display_abstract.strip() and display_abstract != 'Abstract non disponibile.': 
                        st.markdown(f"**Abstract:** _{display_abstract}_")
                    elif display_abstract == 'Abstract non disponibile.':
                        st.caption("Abstract non disponibile.")
                    
                    pm_url = item.get('url', '#')
                    st.markdown(f"[Pagina PubMed]({pm_url})")

                    # Tentativo di trovare un URL PDF diretto per PubMed Central
                    # Questo è un approccio euristico e potrebbe non coprire tutti i casi.
                    # Spesso gli articoli PubMed hanno link a PMC che a sua volta linka al PDF.
                    pdf_url_pmc = None
                    if item.get('pmid'):
                        # Costruisci un possibile link a PMC per l'articolo
                        pmc_article_url = f"https://www.ncbi.nlm.nih.gov/pmc/articles/PMC{item.get('pmid')}/"
                        # A volte il PDF è linkato direttamente da PMC in un formato standard
                        pdf_url_pmc = f"https://www.ncbi.nlm.nih.gov/pmc/articles/PMC{item.get('pmid')}/pdf/"
                        # Potremmo anche aggiungere un bottone per "Tentare di trovare PDF su PMC" che apra pmc_article_url
                    
                    # Diamo priorità a un eventuale link PDF trovato da Semantic Scholar se fosse stato usato prima,
                    # o un campo `pdf_url` se lo aggiungessimo direttamente a PubMed (ma non c'è nativamente)
                    final_pdf_url_to_try = item.get('pdf_url') # Se venisse da altre fonti
                    if not final_pdf_url_to_try and pdf_url_pmc:
                        final_pdf_url_to_try = pdf_url_pmc
                        # st.caption(f"Tenteremo di scaricare il PDF da: {final_pdf_url_to_try}") # Debug

                    button_key_pm = f"analyze_pm_{item.get('pmid', '')}_{i}"
                    if st.button("🔬 Analizza Articolo con AI", key=button_key_pm, use_container_width=True, help="Analizza abstract e tenta di scaricare e analizzare il PDF completo se disponibile su PMC."):
                        if not st.session_state.action_in_progress:
                            st.session_state.action_in_progress = True
                            p_tok_pubmed, c_tok_pubmed = 0, 0
                            analysis_output_parts_pm = []
                            
                            if st.session_state.current_chat_id and st.session_state.current_chat_id in st.session_state.chat_sessions:
                                current_chat_data_pm = st.session_state.chat_sessions[st.session_state.current_chat_id]
                                
                                item_title_pm = item.get('title', 'N/A')
                                item_abstract_pm = item.get('abstract', 'N/A')
                                item_url_pm = item.get('url', '#')
                                item_pmid_pm = item.get('pmid')
                                item_pmcid_pm = item.get('pmcid') # PMCID estratto da search_pubmed

                                user_action_msg_content_pm = f"Richiesta analisi AI per l'articolo PubMed: '{item_title_pm}' (PMID: {item_pmid_pm}, URL: {item_url_pm})"
                                user_action_msg_pm = {"role": "user", "content": user_action_msg_content_pm, "timestamp": datetime.now(timezone.utc).isoformat()}
                                current_chat_data_pm.setdefault("messages",[]).append(user_action_msg_pm)
                                
                                extracted_pdf_text_pm = None
                                attempted_pdf_url_info = "Nessun tentativo di URL PDF specifico."
                                pdf_download_error_detail = ""

                                if item_pmid_pm: # Richiede almeno PMID per tentare
                                    with st.spinner(f"Ricerca PDF full-text per PMID: {item_pmid_pm} (PMCID: {item_pmcid_pm or 'N/A'})..."):
                                        potential_pdf_url = find_pmc_pdf_url(pmid=item_pmid_pm, pmcid=item_pmcid_pm)
                                    
                                    if potential_pdf_url:
                                        # --- INIZIO BLOCCO AGGIUNTO ---
                                        st.write(f"DEBUG PubMed An.: URL iniziale da find_pmc_pdf_url: {potential_pdf_url}")
                                        
                                        headers_check = {
                                            'User-Agent': 'Mozilla/5.0 (compatible; SuperTutorAI/1.0; +http://localhost)',
                                            'Accept': 'application/pdf, text/html;q=0.9, */*;q=0.8' # Tolto xhtml/xml per semplicità
                                        }
                                        
                                        is_direct_pdf = False
                                        refined_url_for_download = potential_pdf_url # Inizia con l'URL che abbiamo

                                        try:
                                            head_resp = requests.head(refined_url_for_download, headers=headers_check, timeout=5, allow_redirects=True, verify=False)
                                            if head_resp.status_code == 200 and 'application/pdf' in head_resp.headers.get('content-type','').lower():
                                                is_direct_pdf = True
                                                refined_url_for_download = head_resp.url 
                                                st.write(f"DEBUG PubMed An.: URL confermato come PDF diretto: {refined_url_for_download}")
                                            else:
                                                st.write(f"DEBUG PubMed An.: HEAD check su {refined_url_for_download} non è PDF diretto (Status: {head_resp.status_code}, CT: {head_resp.headers.get('content-type','')}).")
                                                if 'text/html' in head_resp.headers.get('content-type','').lower():
                                                    refined_url_for_download = head_resp.url 
                                        except requests.exceptions.RequestException as e_head_initial:
                                            st.write(f"DEBUG PubMed An.: HEAD check iniziale fallito per {refined_url_for_download}: {e_head_initial}, procedo con altre strategie.")

                                        if not is_direct_pdf and item_pmcid_pm:
                                            numeric_pmcid = item_pmcid_pm.replace("PMC", "").strip()
                                            if numeric_pmcid.isdigit():
                                                base_article_url_pmc = f"https://www.ncbi.nlm.nih.gov/pmc/articles/PMC{numeric_pmcid}/"
                                                url_is_pmc_non_pdf = False
                                                if 'ncbi.nlm.nih.gov/pmc/articles/' in refined_url_for_download.lower() and \
                                                   f"PMC{numeric_pmcid}".lower() in refined_url_for_download.lower():
                                                    url_is_pmc_non_pdf = True

                                                if url_is_pmc_non_pdf: 
                                                    try:
                                                        parsed_ncbi_base = urlparse(base_article_url_pmc) 
                                                        query_ncbi = urllib.parse.parse_qs(parsed_ncbi_base.query)
                                                        query_ncbi['report'] = ['pdf']
                                                        query_ncbi['format'] = ['raw']
                                                        ncbi_raw_attempt_url = urllib.parse.urlunparse(parsed_ncbi_base._replace(query=urllib.parse.urlencode(query_ncbi, doseq=True)))
                                                        
                                                        st.write(f"DEBUG PubMed An.: Tentativo NCBI raw: {ncbi_raw_attempt_url}")
                                                        head_raw_ncbi = requests.head(ncbi_raw_attempt_url, headers=headers_check, timeout=7, allow_redirects=True, verify=False)
                                                        if head_raw_ncbi.status_code == 200 and 'application/pdf' in head_raw_ncbi.headers.get('content-type','').lower():
                                                            st.success(f"DEBUG PubMed An.: PDF trovato con NCBI raw: {head_raw_ncbi.url}")
                                                            refined_url_for_download = head_raw_ncbi.url
                                                            is_direct_pdf = True 
                                                        else:
                                                            st.write(f"DEBUG PubMed An.: NCBI raw non ha dato PDF ({head_raw_ncbi.status_code}, CT: {head_raw_ncbi.headers.get('content-type','')}).")
                                                    except Exception as e_ncbi_raw:
                                                        st.warning(f"DEBUG PubMed An.: Errore strategia NCBI raw: {e_ncbi_raw}")
                                        
                                        if not is_direct_pdf and item_pmcid_pm: 
                                            numeric_pmcid = item_pmcid_pm.replace("PMC", "").strip() 
                                            if numeric_pmcid.isdigit():
                                                try:
                                                    epmc_article_url = f"https://europepmc.org/articles/PMC{numeric_pmcid}"
                                                    parsed_epmc_base = urlparse(epmc_article_url)
                                                    query_epmc = urllib.parse.parse_qs(parsed_epmc_base.query)
                                                    query_epmc['pdf'] = ['render'] 
                                                    epmc_attempt_url = urllib.parse.urlunparse(parsed_epmc_base._replace(query=urllib.parse.urlencode(query_epmc, doseq=True)))

                                                    st.write(f"DEBUG PubMed An.: Tentativo Europe PMC: {epmc_attempt_url}")
                                                    head_epmc = requests.head(epmc_attempt_url, headers=headers_check, timeout=7, allow_redirects=True, verify=False)
                                                    if head_epmc.status_code == 200 and 'application/pdf' in head_epmc.headers.get('content-type','').lower():
                                                        st.success(f"DEBUG PubMed An.: PDF trovato con Europe PMC: {head_epmc.url}")
                                                        refined_url_for_download = head_epmc.url
                                                    else:
                                                        st.write(f"DEBUG PubMed An.: Europe PMC non ha dato PDF ({head_epmc.status_code}, CT: {head_epmc.headers.get('content-type','')}).")
                                                except Exception as e_epmc:
                                                    st.warning(f"DEBUG PubMed An.: Errore strategia Europe PMC: {e_epmc}")
                                        
                                        potential_pdf_url = refined_url_for_download 
                                        attempted_pdf_url_info = f"URL (post-affinamento): {potential_pdf_url}"
                                        
                                        with st.spinner(f"Download e parsing PDF da: {potential_pdf_url}..."):
                                            # Passa item_title_pm come expected_title
                                            extracted_pdf_text_pm = download_and_parse_pdf_from_url(potential_pdf_url, expected_title=item_title_pm)
                                        
                                        if extracted_pdf_text_pm and not extracted_pdf_text_pm.startswith("Errore") and not extracted_pdf_text_pm.startswith("Timeout"):
                                            st.success(f"PDF scaricato e analizzato con successo da (URL usato): {potential_pdf_url}!")
                                        else:
                                            # Salva il messaggio di errore per un log più dettagliato
                                            pdf_download_error_detail = extracted_pdf_text_pm if extracted_pdf_text_pm else "Nessun testo estratto o errore sconosciuto."
                                            st.warning(f"Non è stato possibile scaricare o analizzare correttamente il PDF da {potential_pdf_url}. L'analisi AI si baserà sull'abstract e sui metadati. (Dettaglio: {pdf_download_error_detail})")
                                            extracted_pdf_text_pm = None 
                                    else:
                                        st.info(f"Nessun URL PDF full-text trovato automaticamente per l'articolo PMID: {item_pmid_pm}. L'analisi si baserà sull'abstract.")
                                        attempted_pdf_url_info = "Nessun URL PDF trovato automaticamente."
                                        pdf_download_error_detail = "Nessun URL PDF trovato."
                                else:
                                    st.info("PMID non disponibile, impossibile cercare PDF. L'analisi si baserà sull'abstract.")
                                    attempted_pdf_url_info = "PMID non disponibile."
                                    pdf_download_error_detail = "PMID non disponibile per ricerca PDF."

                                prompt_per_analisi_pm_parts = [
                                    f"Sto analizzando l'articolo scientifico proveniente da PubMed intitolato: '{item_title_pm}' (PMID: {item_pmid_pm}).",
                                    f"URL Pagina PubMed: {item_url_pm}",
                                ]
                                # Aggiunta sicura dell'abstract
                                abstract_section = "Abstract fornito:\n"
                                if item_abstract_pm:
                                    abstract_section += str(item_abstract_pm) # Assicura che sia una stringa
                                else:
                                    abstract_section += "N/A"
                                prompt_per_analisi_pm_parts.append(abstract_section)

                                if extracted_pdf_text_pm:
                                    prompt_per_analisi_pm_parts.append(f"\nIn aggiunta, è stato possibile estrarre il seguente testo dal PDF completo dell'articolo (potrebbe essere parziale, URL tentato: {attempted_pdf_url_info}):\n---\n{extracted_pdf_text_pm[:10000]}\n---")
                                else:
                                    # Aggiungi il dettaglio dell'errore PDF al prompt per l'AI
                                    prompt_per_analisi_pm_parts.append(f"\nNon è stato possibile analizzare il PDF completo dell'articolo. {attempted_pdf_url_info}. Dettaglio errore PDF: {pdf_download_error_detail}")
                                
                                prompt_per_analisi_pm_parts.append("\nPer favore, esegui le seguenti operazioni basandoti su TUTTE le informazioni fornite (abstract e, se disponibile, il testo del PDF):")
                                prompt_per_analisi_pm_parts.append("1. Fornisci un riassunto conciso ma completo (circa 150-250 parole) degli obiettivi, metodi principali, risultati chiave e conclusioni dell'articolo.")
                                prompt_per_analisi_pm_parts.append("2. Elenca e spiega brevemente i 3-5 punti o scoperte più significativi dell'articolo.")
                                prompt_per_analisi_pm_parts.append("3. Se applicabile e desumibile dal testo, discuti brevemente le implicazioni o la rilevanza di questi risultati.")
                                prompt_per_analisi_pm_parts.append("4. Se il testo del PDF è stato analizzato, integra le informazioni da esso per arricchire i punti precedenti, specificando se un dettaglio proviene dal PDF.")
                                prompt_per_analisi_pm_parts.append("Rispondi in modo chiaro e strutturato, usando Markdown per la formattazione.")
                                prompt_final_pm = "\n\n".join(prompt_per_analisi_pm_parts)

                                chat_history_for_llm_pm = current_chat_data_pm.get("messages", [])[:-1]
                                
                                ai_analysis_text_pm = ""
                                with st.spinner(f"L'AI sta analizzando l'articolo '{item_title_pm}'..."):
                                    ai_analysis_text_pm, p_tok_pubmed, c_tok_pubmed, _ = get_answer_from_llm(
                                        query=prompt_final_pm,
                                        relevant_chunks_data=[],
                                        api_key=st.session_state.gemini_api_key,
                                        model_name=st.session_state.selected_gemini_model,
                                        chat_history=chat_history_for_llm_pm, 
                                        ignore_active_contexts_for_this_query=True, 
                                    )
                                analysis_output_parts_pm.insert(0, ai_analysis_text_pm) # Mette la risposta AI all'inizio
                                
                                current_chat_data_pm['active_article_context'] = {
                                    "title": item_title_pm,
                                    "abstract": item_abstract_pm,
                                    "url": item_url_pm,
                                    "pmid": item_pmid_pm,
                                    "source": "PubMed (con tentativo analisi PDF)",
                                    "full_text_analyzed": bool(extracted_pdf_text_pm) # True se il testo del PDF è stato estratto
                                }
                                current_chat_data_pm.pop('active_web_page_context', None)
                                st.success(f"Articolo '{item_title_pm}' analizzato e impostato come contesto per la chat corrente!")
                            else:
                                analysis_output_parts_pm.append("Errore: Nessuna chat attiva o URL non valido per l'analisi.")
                                st.warning(analysis_output_parts_pm[0])

                            final_output_for_chat_pm = "\n\n".join(analysis_output_parts_pm)
                            add_assistant_response_and_rerun_v2(
                                assistant_content=final_output_for_chat_pm, 
                                current_chat_data=current_chat_data_pm,
                                prompt_tokens=p_tok_pubmed, 
                                completion_tokens=c_tok_pubmed,
                                mark_action_completed=True
                            )
                        else:
                            st.caption("Azione precedente in corso...")

                    if i < len(st.session_state.last_pubmed_results) - 1:
                        st.markdown("---") # Separatore tra i risultati

def extract_google_drive_info(text):
    """Estrae informazioni (tipo e file_id) da un URL di Google Drive, Docs o Sheets."""
    if not text: return None
    # Pattern per Google Docs: https://docs.google.com/document/d/FILE_ID/edit (o view, etc.)
    doc_match = re.search(r"https://docs\.google\.com/document/d/([a-zA-Z0-9_-]+)", text)
    if doc_match:
        return {"type": "doc", "file_id": doc_match.group(1), "original_url": doc_match.group(0)}

    # Pattern per Google Sheets: https://docs.google.com/spreadsheets/d/FILE_ID/edit (o view, etc.)
    sheet_match = re.search(r"https://docs\.google\.com/spreadsheets/d/([a-zA-Z0-9_-]+)", text)
    if sheet_match:
        return {"type": "sheet", "file_id": sheet_match.group(1), "original_url": sheet_match.group(0)}

    # Pattern per Google Drive (file generici): https://drive.google.com/file/d/FILE_ID/view (o edit, etc.)
    # Escludiamo /folders/ per ora per semplicità, ci concentriamo sui file.
    drive_file_match = re.search(r"https://drive\.google\.com/file/d/([a-zA-Z0-9_-]+)", text)
    if drive_file_match:
        return {"type": "drive_file", "file_id": drive_file_match.group(1), "original_url": drive_file_match.group(0)}
    
    return None

def extract_first_url(text):
    """Estrae il primo URL HTTP o HTTPS generico da una stringa di testo."""
    if not text: return None
    # Espressione regolare per trovare URL (semplificata, ma dovrebbe coprire molti casi)
    url_pattern = r'https?://[\w\d:#@%/;$()~_?\+-=\\\.&]*'
    match = re.search(url_pattern, text)
    return match.group(0) if match else None

def add_assistant_response_and_rerun_v2(assistant_content, current_chat_data, prompt_tokens=0, completion_tokens=0, citation_map=None, chunks_info=None, mark_action_completed=False):
    """Aggiunge una risposta dell'assistente, aggiorna i token (se forniti), persiste e fa rerun."""
    ai_msg_data = {
        "role": "assistant", 
        "content": assistant_content,
        "timestamp": datetime.now(timezone.utc).isoformat() # Aggiungi timestamp UTC ISO per AI
    }
    if citation_map: # Attualmente citation_map è sempre vuoto, ma lo manteniamo per coerenza
        ai_msg_data["cited_chunks_map"] = citation_map
    if chunks_info:
        ai_msg_data["chunks_info"] = chunks_info
    
    current_chat_data.setdefault("messages", []).append(ai_msg_data)
    
    # Aggiorna i token solo se sono stati forniti valori positivi (tipicamente da una chiamata LLM)
    if prompt_tokens > 0 or completion_tokens > 0:
        current_chat_data.update({
            "tokens_prompt": current_chat_data.get("tokens_prompt", 0) + prompt_tokens,
            "tokens_completion": current_chat_data.get("tokens_completion", 0) + completion_tokens,
            "tokens_total": current_chat_data.get("tokens_total", 0) + prompt_tokens + completion_tokens
        })
        # Assicurati che current_chat_id e selected_gemini_model siano accessibili o passati se necessario
        # Qui presumiamo che siano in st.session_state come al solito
        if st.session_state.get("current_chat_id") and st.session_state.get("selected_gemini_model"):
            st.session_state.setdefault("token_usage_log", []).append({
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "chat_id": st.session_state.current_chat_id,
                "model": st.session_state.selected_gemini_model,
                "prompt_tokens": prompt_tokens,
                "completion_tokens": completion_tokens,
                "total_tokens": prompt_tokens + completion_tokens
            })

    current_chat_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    
    if mark_action_completed:
        st.session_state.action_in_progress = False # Sblocca per la prossima azione
        
    persist_app_data() # Salva lo stato PRIMA del rerun
    st.rerun()

chat_area_container = st.container()
if not st.session_state.gemini_api_key: 
    chat_area_container.warning("🔑 Inserisci Gemini API Key.")
elif not st.session_state.current_chat_id or st.session_state.current_chat_id not in st.session_state.chat_sessions: 
    chat_area_container.info("➕ Nuova chat o seleziona esistente dalla sidebar per iniziare.")
elif st.session_state.current_chat_id and st.session_state.current_chat_id in st.session_state.chat_sessions:
    cur_chat_data = st.session_state.chat_sessions[st.session_state.current_chat_id]
    chat_name = cur_chat_data.get("name","Chat"); chat_ts_iso = cur_chat_data.get("updated_at")
    try: chat_ts_str = datetime.fromisoformat(chat_ts_iso.replace("Z","+00:00")).strftime('%d/%m %H:%M') if chat_ts_iso else "N/D" # %Y rimosso per brevità
    except: chat_ts_str = "N/D"
    sidebar_update_time_placeholder.caption(f"Chat: '{chat_name}'\nUltimo Agg: {chat_ts_str}")
    
    # ---- Visualizzazione e Rimozione Contesto Articolo Selezionato (dalla chat corrente) ----
    active_article_for_display = None
    active_web_page_for_display = None
    if cur_chat_data: 
        active_article_for_display = cur_chat_data.get("active_article_context")
        active_web_page_for_display = cur_chat_data.get("active_web_page_context")

    if active_article_for_display:
        article_ctx = active_article_for_display
        ctx_title = article_ctx.get("title", "N/A")
        ctx_source = article_ctx.get("source", "N/A")
        ctx_url = article_ctx.get("url", "#")
        
        col1, col2 = st.columns([0.8, 0.2])
        with col1:
            st.info(f"🔬 Contesto AI per QUESTA CHAT: Articolo '{ctx_title}' ({ctx_source}). [Vedi Fonte]({ctx_url})")
        with col2:
            if st.button("🗑️ Rimuovi Contesto Articolo da Chat", key=f"remove_article_ctx_{st.session_state.current_chat_id}", help="Rimuovi il contesto dell'articolo selezionato per le prossime domande AI in questa chat."):
                if cur_chat_data:
                    cur_chat_data['active_article_context'] = None
                    persist_app_data()
                    st.success("Contesto articolo rimosso da questa chat.")
                    st.rerun()
    
    if active_web_page_for_display and not active_article_for_display:
        wp_ctx = active_web_page_for_display
        wp_title = wp_ctx.get("title", "N/A")
        wp_source = wp_ctx.get("source", "Pagina Web")
        wp_url = wp_ctx.get("url", "#")
        
        col1_wp, col2_wp = st.columns([0.8, 0.2])
        with col1_wp:
            st.info(f"🌐 Contesto AI per QUESTA CHAT: Pagina Web '{wp_title}' ({wp_source}). [Vedi Fonte Originale]({wp_url})")
        with col2_wp:
            if st.button("🗑️ Rimuovi Contesto Pagina Web", key=f"remove_wp_ctx_{st.session_state.current_chat_id}", help="Rimuovi il contesto della pagina web analizzata per le prossime domande AI in questa chat."):
                if cur_chat_data:
                    cur_chat_data['active_web_page_context'] = None
                    persist_app_data()
                    st.success("Contesto pagina web rimosso da questa chat.")
                    st.rerun()
    
    # ----- Legenda Comandi e Funzionalità -----
    # st.markdown("---")
    # with st.expander("💡 Guida Rapida e Comandi", expanded=False):
    #     st.markdown(""" ... (contenuto legenda rimosso) ... """)

    # ---- Fine Visualizzazione e Rimozione ----

    new_c_name = st.text_input("Nome Chat:", value=chat_name, key=f"ren_chat_{st.session_state.current_chat_id}")
    if new_c_name != chat_name and new_c_name.strip():
        cur_chat_data["name"] = new_c_name.strip(); cur_chat_data["updated_at"] = datetime.now(timezone.utc).isoformat(); persist_app_data()

    for msg_idx, msg in enumerate(cur_chat_data.get("messages",[])):
        with st.chat_message(msg["role"]):
            if msg["role"] == "assistant" and "cited_chunks_map" in msg and msg.get("cited_chunks_map") and msg.get("content"): 
                display_response_with_citations(msg["content"], msg["cited_chunks_map"])
            elif msg.get("content"):
                st.markdown(msg.get("content"))
                # Aggiungi timestamp per i messaggi utente O AI
                if "timestamp" in msg:
                    try:
                        dt_object = datetime.fromisoformat(msg["timestamp"].replace("Z", "+00:00"))
                        formatted_time = dt_object.astimezone().strftime("%d/%m/%Y %H:%M:%S")
                        # Differenzia leggermente lo stile o la posizione se necessario, ma per ora usiamo solo st.caption
                        st.caption(f"{formatted_time}")
                    except Exception as e:
                        st.caption(f"(Ora messaggio non disponibile)") 
    else:
                st.caption("[Messaggio vuoto]")
            
            # Bottoni sotto i messaggi dell'assistente
                if msg["role"] == "assistant" and msg.get("content"):
                    assistant_message_content = msg.get("content")
                    parsed_table_for_button = parse_markdown_table_from_text(assistant_message_content)
                else: # Assicura che parsed_table_for_button sia definita
                    parsed_table_for_button = None

                cols = st.columns(3) # Sempre 3 colonne

                # Colonna 1: Bottone Sheet (dinamico o statico)
                with cols[0]:
                    unique_key_sheet = f"sheet_action_{st.session_state.current_chat_id}_{msg_idx}"
                    # Determina se parsed_table_for_button dovrebbe essere calcolato qui
                    # Lo ricalcoliamo solo se il messaggio è dell'assistente
                    # per evitare errori se msg non è dell'assistente
                    current_message_is_assistant = msg["role"] == "assistant" and msg.get("content")
                    
                    # Ricalcola parsed_table_for_button solo se necessario
                    if current_message_is_assistant:
                        assistant_content_for_this_button = msg.get("content")
                        parsed_table_for_this_button_action = parse_markdown_table_from_text(assistant_content_for_this_button)
                    else:
                        parsed_table_for_this_button_action = None


                    if parsed_table_for_this_button_action:
                        if st.button("✨ Crea Sheet con Tabella", key=f"dynamic_{unique_key_sheet}", help="Crea e popola un Google Sheet con la tabella rilevata in questo messaggio"):
                            if not st.session_state.action_in_progress:
                                st.session_state.action_in_progress = True
                                if not st.session_state.get("google_creds"):
                                    st.warning("Autorizza prima i Servizi Google nella sidebar.")
                                    st.session_state.action_in_progress = False # Sblocca se l'azione non può procedere
                                else:
                                    sheet_title = f"Sheet da Tabella (msg {msg_idx+1})"
                                    if parsed_table_for_this_button_action[0]: # Usa la versione ricalcolata
                                        header_preview = ' '.join(map(str, parsed_table_for_this_button_action[0]))[:50]
                                        sheet_title = f"Sheet: {header_preview}... (AI msg {msg_idx+1})"
                                    with st.spinner(f"Creazione Google Sheet '{sheet_title}'..."):
                                        sheet_id = create_google_sheet(sheet_title)
                                        action_message = ""
                                        if sheet_id:
                                            with st.spinner(f"Popolamento '{sheet_title}'..."):
                                                write_success = write_data_to_google_sheet(sheet_id, parsed_table_for_this_button_action) # Usa la versione ricalcolata
                                                if write_success:
                                                    action_message = f"✅ Foglio '{sheet_title}' creato e popolato: https://docs.google.com/spreadsheets/d/{sheet_id}/edit"
                                                else:
                                                    action_message = f"⚠️ Foglio '{sheet_title}' creato (https://docs.google.com/spreadsheets/d/{sheet_id}/edit), ma errore nel popolarlo."
                                        else:
                                            action_message = f"❌ Errore creazione Google Sheet '{sheet_title}'."
                                        add_assistant_response_and_rerun_v2(action_message, cur_chat_data, mark_action_completed=True)
                            # else: st.info("Azione precedente in corso...") # Debug opzionale
                    elif current_message_is_assistant: # Mostra il bottone statico solo per messaggi AI
                        if st.button("📊 Crea Google Sheet", key=f"static_{unique_key_sheet}", help="Crea un nuovo Google Sheet (vuoto o da contenuto generico del messaggio)"):
                            if not st.session_state.action_in_progress:
                                st.session_state.action_in_progress = True
                                if not st.session_state.get("google_creds"):
                                    st.warning("Autorizza prima i Servizi Google nella sidebar.")
                                    st.session_state.action_in_progress = False
                                else:
                                    sheet_title = f"Sheet da msg AI {msg_idx+1}"
                                    # Qui, diversamente dal bottone dinamico, potremmo tentare un parsing o creare vuoto
                                    # Si potrebbe estendere per cercare una tabella anche qui, ma il dinamico ha priorità.
                                    data_to_populate_static = parse_markdown_table_from_text(msg.get("content")) # Usa il contenuto del messaggio corrente
                                    if data_to_populate_static and data_to_populate_static[0]:
                                         header_preview = ' '.join(map(str, data_to_populate_static[0]))[:50]
                                         sheet_title = f"Sheet (da contenuto generico): {header_preview}..."
                                    
                                    with st.spinner(f"Creazione Google Sheet '{sheet_title}'..."):
                                        sheet_id = create_google_sheet(sheet_title)
                                        action_message = ""
                                        if sheet_id:
                                            if data_to_populate_static:
                                                with st.spinner(f"Popolamento '{sheet_title}'..."):
                                                    write_success = write_data_to_google_sheet(sheet_id, data_to_populate_static)
                                                    if write_success:
                                                        action_message = f"✅ Foglio '{sheet_title}' creato e popolato (contenuto generico): https://docs.google.com/spreadsheets/d/{sheet_id}/edit"
                                                    else:
                                                        action_message = f"⚠️ Foglio '{sheet_title}' creato (https://docs.google.com/spreadsheets/d/{sheet_id}/edit), ma errore popolarlo (contenuto generico)."
                                            else:
                                                action_message = f"✅ Foglio '{sheet_title}' creato (vuoto): https://docs.google.com/spreadsheets/d/{sheet_id}/edit. Nessuna tabella specifica trovata per auto-popolamento."
                                        else:
                                            action_message = f"❌ Errore creazione Google Sheet '{sheet_title}'."
                                        add_assistant_response_and_rerun_v2(action_message, cur_chat_data, mark_action_completed=True)
                            # else: st.info("Azione precedente in corso...")

                # Colonna 2: Esporta in Docs
                with cols[1]:
                    if current_message_is_assistant: # Mostra bottone solo per messaggi AI
                        if st.button("📄 Esporta in Docs", key=f"export_doc_{st.session_state.current_chat_id}_{msg_idx}", help="Esporta questo messaggio su Google Docs"):
                            if not st.session_state.action_in_progress:
                                st.session_state.action_in_progress = True
                                if not st.session_state.get("google_creds"):
                                    st.warning("Autorizza prima i Servizi Google nella sidebar.")
                                    st.session_state.action_in_progress = False
                                else:
                                    with st.spinner("Creazione documento Google Docs..."):
                                        docs_service = get_google_service('docs', 'v1')
                                        export_action_message = ""
                                        if not docs_service:
                                            export_action_message = "❌ Servizi Google Docs non autorizzati o errore."
                                        else:
                                            doc_title = f"Esportazione SuperTutorAI - Messaggio {msg_idx+1}"
                                            try:
                                                doc = docs_service.documents().create(body={"title": doc_title}).execute()
                                                document_id = doc.get("documentId")
                                                # Usa il contenuto del messaggio AI corrente
                                                requests_body = [{"insertText": {"location": {"index": 1}, "text": msg.get("content")}}]
                                                docs_service.documents().batchUpdate(documentId=document_id, body={"requests": requests_body}).execute()
                                                export_action_message = f"✅ Documento '{doc_title}' creato: https://docs.google.com/document/d/{document_id}/edit"
                                            except Exception as e:
                                                export_action_message = f"❌ Errore creazione Google Doc: {e}"
                                    add_assistant_response_and_rerun_v2(export_action_message, cur_chat_data, mark_action_completed=True)
                            # else: st.info("Azione precedente in corso...")

                # Colonna 3: Carica su Drive (logica aggiornata con priorità)
                with cols[2]:
                    upload_key_suffix = f"upload_drive_action_{st.session_state.current_chat_id}_{msg_idx}"
                    
                    # Mostra bottone solo per messaggi AI, dato che prende il contenuto del messaggio AI come fallback
                    if current_message_is_assistant:
                        if st.button("⬆️ Carica su Drive", key=upload_key_suffix, help="Carica su Drive (file allegato > file da link nel msg > testo msg)"):
                            if not st.session_state.action_in_progress:
                                st.session_state.action_in_progress = True
                                drive_action_message = ""

                                if not st.session_state.get("google_creds"):
                                    drive_action_message = "❌ Servizi Google non autorizzati. Autorizza prima dalla sidebar."
                                    st.session_state.action_in_progress = False 
                                    st.warning(drive_action_message) 
                                    st.rerun()
                                else:
                                    # Priorità 1: File allegato all'applicazione
                                    if st.session_state.get("uploaded_file_bytes") and st.session_state.get("uploaded_file_name"):
                                        file_bytes = st.session_state.uploaded_file_bytes
                                        file_name = st.session_state.uploaded_file_name
                                        mime_type = st.session_state.uploaded_file_mime
                                        source_description = f"""file allegato all'app ('{file_name}')"""
                                        with st.spinner(f"""Caricamento di {source_description} su Google Drive..."""):
                                            try:
                                                file_id = upload_file_to_drive(file_bytes, file_name, mime_type)
                                                if file_id:
                                                    drive_action_message = f"""✅ {source_description.capitalize()} caricato: https://drive.google.com/file/d/{file_id}/view"""
                                                else:
                                                    drive_action_message = f"""❌ Errore caricamento {source_description} su Drive."""
                                            except Exception as e:
                                                drive_action_message = f"""❌ Eccezione caricamento {source_description}: {repr(e)}""" 
                                        add_assistant_response_and_rerun_v2(drive_action_message, cur_chat_data, mark_action_completed=True)
                                    
                                    # Priorità 2: File da URL nel messaggio AI (usa il contenuto del messaggio corrente)
                                    else:
                                        assistant_content_for_drive_url = msg.get("content")
                                        google_drive_info = extract_google_drive_info(assistant_content_for_drive_url)
                                        url_to_process_generically = None

                                        if google_drive_info:
                                            file_id_from_url = google_drive_info["file_id"] # Rinominato per chiarezza
                                            file_type_from_url = google_drive_info["type"] # Rinominato
                                            original_url_from_msg = google_drive_info["original_url"] # Rinominato
                                            source_description = f"file Google ({file_type_from_url}) dall\'URL: {original_url_from_msg}"
                                            
                                            with st.spinner(f"Accesso e download di {source_description}..."):
                                                try:
                                                    drive_service = get_google_service('drive', 'v3')
                                                    if not drive_service:
                                                        drive_action_message = "❌ Servizio Google Drive non disponibile o non autorizzato correttamente."
                                                    else:
                                                        downloaded_bytes = None
                                                        downloaded_file_name = f"{file_id_from_url}_export"
                                                        downloaded_mime_type = 'application/octet-stream' 

                                                        if file_type_from_url == 'doc':
                                                            export_mime_type = 'application/pdf'
                                                            request_export = drive_service.files().export_media(fileId=file_id_from_url, mimeType=export_mime_type)
                                                            downloader_export = io.BytesIO(request_export.execute())
                                                            downloaded_bytes = downloader_export.getvalue()
                                                            file_metadata = drive_service.files().get(fileId=file_id_from_url, fields='name').execute()
                                                            original_name = file_metadata.get('name', file_id_from_url)
                                                            downloaded_file_name = f"{original_name}.pdf"
                                                            downloaded_mime_type = export_mime_type
                                                        
                                                        elif file_type_from_url == 'sheet':
                                                            export_mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                                                            request_export = drive_service.files().export_media(fileId=file_id_from_url, mimeType=export_mime_type)
                                                            fh_export = io.BytesIO(request_export.execute())
                                                            downloaded_bytes = fh_export.getvalue()
                                                            file_metadata = drive_service.files().get(fileId=file_id_from_url, fields='name').execute()
                                                            original_name = file_metadata.get('name', file_id_from_url)
                                                            downloaded_file_name = f"{original_name}.xlsx"
                                                            downloaded_mime_type = export_mime_type

                                                        elif file_type_from_url == 'drive_file':
                                                            file_metadata = drive_service.files().get(fileId=file_id_from_url, fields='name, mimeType').execute()
                                                            downloaded_file_name = file_metadata.get('name', f"{file_id_from_url}_drive_file")
                                                            downloaded_mime_type = file_metadata.get('mimeType', 'application/octet-stream')
                                                            request_get_media = drive_service.files().get_media(fileId=file_id_from_url)
                                                            fh_get_media = io.BytesIO()
                                                            from googleapiclient.http import MediaIoBaseDownload
                                                            downloader_get_media = MediaIoBaseDownload(fh_get_media, request_get_media)
                                                            done_get_media = False
                                                            while done_get_media is False: status_get_media, done_get_media = downloader_get_media.next_chunk()
                                                            downloaded_bytes = fh_get_media.getvalue()

                                                        if downloaded_bytes:
                                                            uploaded_file_id_drive_new = upload_file_to_drive(downloaded_bytes, downloaded_file_name, downloaded_mime_type) # Rinominato
                                                            if uploaded_file_id_drive_new:
                                                                drive_action_message = f"""✅ File '{downloaded_file_name}' (da Google Drive URL) caricato su Drive: https://drive.google.com/file/d/{uploaded_file_id_drive_new}/view"""
                                                            else:
                                                                drive_action_message = f"""❌ Errore durante il caricamento del file '{downloaded_file_name}' (da Google Drive URL) su Drive."""
                                                        else:
                                                            drive_action_message = f"""❌ Impossibile scaricare/processare il file Google: {file_id_from_url} (tipo: {file_type_from_url})."""
                                                
                                                except HTTPError as e: # Riferimento corretto a e.resp
                                                    if e.resp.status == 401: drive_action_message = f"""❌ Accesso negato (401) per Google file {original_url_from_msg}."""
                                                    elif e.resp.status == 403: drive_action_message = f"""❌ Accesso proibito (403) per Google file {original_url_from_msg}."""
                                                    elif e.resp.status == 404: drive_action_message = f"""❌ File Google non trovato (404) all'URL: {original_url_from_msg}."""
                                                    else: drive_action_message = f"""❌ Errore HTTP {e.resp.status} ({e.reason if hasattr(e, 'reason') else e.resp.reason}) con Google file: {original_url_from_msg}"""
                                                except Exception as e:
                                                    drive_action_message = f"""❌ Errore imprevisto ({type(e).__name__}) processando Google file {original_url_from_msg}: {repr(e)}"""
                                            add_assistant_response_and_rerun_v2(drive_action_message, cur_chat_data, mark_action_completed=True)
                                        
                                        else: # Non è un link Google Drive, prova come URL generico
                                            url_to_process_generically = extract_first_url(assistant_content_for_drive_url)

                                        if url_to_process_generically:
                                            source_description = f"""file da URL ({url_to_process_generically})"""
                                            with st.spinner(f"""Tentativo di download e caricamento di {source_description}..."""):
                                                try:
                                                    req = Request(url_to_process_generically, headers={'User-Agent': 'Mozilla/5.0'})
                                                    with urlopen(req, timeout=15) as response:
                                                        downloaded_bytes_generic = response.read() # Rinominato
                                                        content_type_header_generic = response.info().get('Content-Type') # Rinominato
                                                        content_disposition_header_generic = response.info().get('Content-Disposition') # Rinominato
                                                        
                                                        fname_from_header_generic = None # Rinominato
                                                        if content_disposition_header_generic:
                                                            disp_parts_generic = content_disposition_header_generic.split(';') # Rinominato
                                                            for part_generic in disp_parts_generic: # Rinominato
                                                                part_generic = part_generic.strip()
                                                                if part_generic.lower().startswith('filename='):
                                                                    fname_from_header_generic = part_generic.split('=', 1)[1].strip('" \'')
                                                                    break
                                                        
                                                        downloaded_file_name_generic = fname_from_header_generic # Rinominato
                                                        if not downloaded_file_name_generic:
                                                            path_from_url_generic = urlparse(url_to_process_generically).path # Rinominato
                                                            downloaded_file_name_generic = os.path.basename(path_from_url_generic) if path_from_url_generic else "downloaded_file"
                                                            if not downloaded_file_name_generic: downloaded_file_name_generic = "downloaded_content"

                                                        downloaded_mime_type_generic = content_type_header_generic.split(';')[0].strip() if content_type_header_generic else None # Rinominato
                                                        if not downloaded_mime_type_generic:
                                                            guessed_type_generic, _ = mimetypes.guess_type(downloaded_file_name_generic) # Rinominato
                                                            downloaded_mime_type_generic = guessed_type_generic if guessed_type_generic else 'application/octet-stream'
                                                        
                                                        if '.' not in downloaded_file_name_generic and downloaded_mime_type_generic != 'application/octet-stream':
                                                            guessed_ext_generic = mimetypes.guess_extension(downloaded_mime_type_generic) # Rinominato
                                                            if guessed_ext_generic: downloaded_file_name_generic += guessed_ext_generic

                                                        file_id_drive_generic = upload_file_to_drive(downloaded_bytes_generic, downloaded_file_name_generic, downloaded_mime_type_generic) # Rinominato
                                                        if file_id_drive_generic:
                                                            drive_action_message = f"""✅ File '{downloaded_file_name_generic}' da URL caricato: https://drive.google.com/file/d/{file_id_drive_generic}/view"""
                                                        else: 
                                                            drive_action_message = f"""❌ Errore caricamento file '{downloaded_file_name_generic}' (da URL) su Drive."""

                                                except HTTPError as e: # Riferimento corretto a e.code, e.reason
                                                    if e.code == 401: drive_action_message = f"""❌ Accesso negato (401) scaricando da URL: {url_to_process_generically}."""
                                                    elif e.code == 403: drive_action_message = f"""❌ Accesso proibito (403) scaricando da URL: {url_to_process_generically}."""
                                                    else: drive_action_message = f"""❌ Errore HTTP {e.code} ({e.reason}) scaricando da URL: {url_to_process_generically}"""
                                                except URLError as e: drive_action_message = f"""❌ Errore URL (es. rete): {e.reason} - URL: {url_to_process_generically}"""
                                                except Exception as e: drive_action_message = f"""❌ Errore imprevisto ({type(e).__name__}) download/upload da URL: {repr(e)}"""
                                            add_assistant_response_and_rerun_v2(drive_action_message, cur_chat_data, mark_action_completed=True)
                                        
                                        # Priorità 3: Testo del messaggio AI corrente come .txt
                                        else: 
                                            source_description = "risposta AI corrente"
                                            # Usa il contenuto del messaggio AI corrente (msg.get("content"))
                                            file_bytes_from_ai_msg = msg.get("content").encode('utf-8') # Rinominato
                                            file_name_for_ai_msg = f"Risposta_AI_msg{msg_idx+1}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt" # Rinominato
                                            mime_type_for_ai_msg = "text/plain" # Rinominato
                                            with st.spinner(f"Caricamento della {source_description} come '{file_name_for_ai_msg}' su Google Drive..."):
                                                try:
                                                    file_id_ai_msg_drive = upload_file_to_drive(file_bytes_from_ai_msg, file_name_for_ai_msg, mime_type_for_ai_msg) # Rinominato
                                                    if file_id_ai_msg_drive:
                                                        drive_action_message = f"""✅ {source_description.capitalize()} caricata come '{file_name_for_ai_msg}': https://drive.google.com/file/d/{file_id_ai_msg_drive}/view"""
                                                    else:
                                                        drive_action_message = f"""❌ Errore caricamento {source_description} come '{file_name_for_ai_msg}' su Drive."""
                                                except Exception as e:
                                                    drive_action_message = f"""❌ Eccezione caricamento {source_description} come '{file_name_for_ai_msg}': {repr(e)}"""
                                                add_assistant_response_and_rerun_v2(drive_action_message, cur_chat_data, mark_action_completed=True)
                            # else: st.info("Azione precedente in corso...")
            
            # Espansore per i chunk considerati (logica invariata)
                if "chunks_info" in msg and msg["role"] == "assistant":
                    with st.expander("Passaggi considerati per questa risposta", expanded=False):
                        if not msg.get("chunks_info"): st.caption("Nessun passaggio specifico identificato come altamente rilevante per questa domanda.")
                        else:
                            for i, ci in enumerate(msg.get("chunks_info",[])): # msg.get("chunks_info",[]) è una lista di dizionari
                                sim = ci.get('similarity'); sim_str = f"{float(sim):.2f}" if isinstance(sim,(int,float)) else "N/A"
                                st.caption(f"Passaggio Candidato {i+1} (Similarità: {sim_str}):\n```{ci.get('chunk','')[:150]}...```")
    
                # ---- NUOVI SUGGERIMENTI SOTTO OGNI MESSAGGIO AI (COME EXPANDER) ----
                if msg["role"] == "assistant" and msg.get("content"):
                    with st.expander("💡 Suggerimenti Rapidi e Comandi Utili", expanded=False):
                        st.markdown("""
                        **Gestione Contesti Specifici:**
                        *   **Articoli Scientifici/Pagine Web Analizzate:** 
                            *   Usa le funzioni di ricerca nella sidebar (Google, Semantic Scholar, PubMed) per trovare fonti.
                            *   Clicca "🔬 **Analizza con AI**" sotto un risultato per impostarlo come contesto specifico per la *chat corrente*.
                            *   Un avviso blu sopra l'area di chat ti ricorderà il contesto attivo. Rimuovilo con "🗑️ **Rimuovi Contesto ...**".
                        *   **Documenti PDF/Testo/Immagini Caricati:**
                            *   Carica un file tramite "📄 Documento Allegato" nella sidebar.
                            *   Se attivo, verrà usato come contesto (priorità inferiore rispetto a un articolo/pagina web specifici per la chat).
                            *   Per scollegarlo: clicca "**Rimuovi documento**".
                        *   **Ignorare Contesti Esterni (per una singola domanda):**
                            *   Includi nella tua domanda: `staccandoci da questo contesto` per forzare l'AI a usare solo la sua conoscenza generale.

                        **Istruzioni Temporanee per Singole Risposte:**
                        *   Dai istruzioni valide solo per la prossima risposta AI, senza che vengano memorizzate.
                        *   Inizia il tuo messaggio con:
                            *   `solo per questa risposta, segui queste istruzioni: [le tue direttive specifiche qui]`
                            *   `istruzioni temporanee: [le tue direttive specifiche qui]`
                            *   L'AI tenterà di seguire queste direttive per elaborare la parte restante del tuo messaggio (se presente prima della keyword) o il contesto implicito.
                        
                        **Ricerca nell'Intero Storico Chat (RAG Globale):**
                        *   Fai domande su argomenti trattati in qualsiasi tua chat passata.
                        *   Inizia il tuo messaggio con:
                            *   `cerca in tutto lo storico chat: [la tua domanda/argomento]`
                            *   `cosa abbiamo detto in passato riguardo a: [la tua domanda/argomento]`
                            *   `riassumi da tutte le chat l'argomento: [la tua domanda/argomento]`
                        *   L'AI cercherà nei messaggi di tutte le tue chat e sintetizzerà una risposta.

                        **Ricerca Web Esplicita:**
                        *   Per informazioni aggiornate o ricerche web dirette.
                        *   Usa frasi come: `cerca su internet [tua query]`, `ultime notizie su [argomento]`, ecc.
                        
                        **Azioni Rapide (Bottoni sotto questa sezione, sotto la risposta AI):**
                        *   📊 **Crea Google Sheet**: Crea e popola un Google Sheet (anche da tabelle nella risposta AI).
                        *   📄 **Esporta in Docs**: Salva la risposta AI in un Google Doc.
                        *   ⬆️ **Carica su Drive**: Carica file su Google Drive (allegato all'app, da URL nella risposta AI, o il testo stesso della risposta AI).

                        **Altre Funzionalità Utili (nella Sidebar):**
                        *   🔍 **Ricerca Google**: Per ricerche web generiche (i risultati appaiono nella sidebar).
                        *   📜 **Storico Chat**: Gestisci le tue chat, creane di nuove, o cerca tra le vecchie (la ricerca qui è testuale semplice o dinamica sui messaggi se digiti).
                        *   ✍️ **Rinomina Chat**: Puoi modificare il nome della chat corrente direttamente sopra l'area di conversazione.
                        """)
                            
    user_q = st.chat_input("Scrivi il tuo messaggio...", key=f"chat_in_main_{st.session_state.current_chat_id}")
    if user_q:
        if not st.session_state.action_in_progress:
            st.session_state.action_in_progress = True
            
            user_input_original = user_q # Salva l'input originale completo dell'utente
            actual_query_for_llm = user_input_original
            temp_instructions_for_llm = None
            retrieved_history_context_for_llm = None # Inizializza per RAG storico
            retrieved_history_search_active = False # Flag per sapere se la ricerca RAG storico è attiva

            # --- Keyword per RAG su intero storico chat ---
            history_rag_keywords = {
                "cerca in tutto lo storico chat:": 
                    "Sto cercando informazioni in tutte le tue conversazioni passate.",
                "cosa abbiamo detto in passato riguardo a:":
                    "Consulto lo storico completo delle nostre chat per rispondere.",
                "riassumi da tutte le chat l'argomento:":
                    "Analizzo tutte le chat per riassumere l'argomento richiesto."
            }

            for keyword, response_intro in history_rag_keywords.items():
                if user_input_original.lower().startswith(keyword.lower()):
                    query_for_history_search = user_input_original[len(keyword):].strip()
                    if query_for_history_search:
                        actual_query_for_llm = query_for_history_search # La query per l'LLM è ciò che l'utente vuole sapere dallo storico
                        # st.info(f"{response_intro} Query: '{query_for_history_search}'") # Feedback opzionale all'utente
                        
                        embedding_model_instance = load_embedding_model_cached()
                        if embedding_model_instance:
                            # Creiamo un hash semplice dello stato delle chat per invalidare la cache se cambiano
                            # Questo potrebbe essere ottimizzato, ma per ora è un segnale di cambiamento.
                            chat_sessions_json_for_hash = json.dumps(st.session_state.chat_sessions, sort_keys=True)
                            sessions_hash = hash(chat_sessions_json_for_hash)

                            retrieved_history_context_for_llm = search_entire_chat_history_semantically(
                                sessions_hash, 
                                query_for_history_search, 
                                embedding_model_instance,
                                top_n=7, # Configurabile
                                similarity_threshold=0.25 # Configurabile
                            )
                            retrieved_history_search_active = True
                        else:
                            st.error("Modello di embedding non disponibile per la ricerca nello storico.")
                        # Non cercare altre keyword se una per RAG storico è stata trovata
                        temp_instructions_for_llm = None # Assicurati che le istruzioni temporanee non interferiscano
                        break 
                    else:
                        st.warning("Per favore, specifica cosa cercare nello storico chat dopo la frase chiave.")
                        # Impedisci di procedere se la query di ricerca è vuota
                        st.session_state.action_in_progress = False
                        st.rerun()
                        # Aggiungere `return` o un modo per uscire qui se dentro una funzione, 
                        # altrimenti lo script continua. Dato che siamo nel flusso principale, st.rerun() dovrebbe bastare.
            
            # --- Keyword per istruzioni temporanee (solo se RAG storico non è attivo) ---
            if not retrieved_history_search_active:
                keyword_phrases_map = {
                    "solo per questa risposta, segui queste istruzioni:": "...", # Contenuto effettivo non cruciale qui perché temp_instructions_for_llm prende il raw
                    "istruzioni temporanee:": "..."
                }
                for keyword_text, _ in keyword_phrases_map.items():
                    if keyword_text.lower() in user_input_original.lower():
                        keyword_start_index = user_input_original.lower().find(keyword_text.lower())
                        actual_query_for_llm = user_input_original[:keyword_start_index].strip()
                        temp_instructions_raw = user_input_original[keyword_start_index + len(keyword_text):].strip()
                        if not actual_query_for_llm: pass 
                        if temp_instructions_raw:
                            temp_instructions_for_llm = temp_instructions_raw
                            break 
            
            user_msg_data = {
                "role": "user", 
                "content": user_input_original, # Salva SEMPRE l'input utente originale completo
                "timestamp": datetime.now(timezone.utc).isoformat()
            }
            cur_chat_data.setdefault("messages",[]).append(user_msg_data)
            cur_chat_data["updated_at"] = datetime.now(timezone.utc).isoformat()

            lower_q = user_q.lower()
            
            # ---- NUOVA LOGICA PER RICERCA WEB ESPLICITA ----
            processed_web_search_results = None
            keywords_for_web_search = [
                "cerca su internet", "cerca online", "trova online", 
                "ultime notizie su", "risultato di", "meteo a", "meteo oggi",
                "cosa sta facendo", "come sta andando la partita", "risultati calcio",
                "quotazione azioni", "prezzo azioni"
            ]

            should_perform_search = any(keyword in lower_q for keyword in keywords_for_web_search)

            if should_perform_search:
                if st.session_state.google_api_key and st.session_state.google_cse_id:
                    with st.spinner("Ricerca Google in corso per te..."):
                        # Usiamo la funzione perform_google_custom_search che restituisce una lista o una stringa di errore
                        search_output = perform_google_custom_search(user_q, num_results=3)
                        if isinstance(search_output, list): # Successo, search_output è una lista di dizionari
                            # Formattiamo i risultati in una singola stringa per il prompt
                            results_str_parts = []
                            if not search_output: # Lista vuota di risultati
                                results_str_parts.append("La ricerca Google non ha prodotto risultati.")
                            else:
                                results_str_parts.append("Ecco i risultati principali dalla ricerca Google:")
                                for i, res in enumerate(search_output):
                                    results_str_parts.append(f"Risultato [{i+1}]:")
                                    results_str_parts.append(f"  Titolo: {res.get('title', 'N/A')}")
                                    results_str_parts.append(f"  Snippet: {res.get('snippet', 'N/A')}")
                                    results_str_parts.append(f"  URL: {res.get('url', '#')}")
                                    results_str_parts.append("---")
                            processed_web_search_results = "\n".join(results_str_parts)
                        elif isinstance(search_output, str): # Errore, search_output è una stringa di errore
                            processed_web_search_results = f"Tentativo di ricerca Google fallito: {search_output}"
                            st.warning(processed_web_search_results) # Mostra avviso all'utente
                        # Se non è né lista né stringa (improbabile), processed_web_search_results rimane None
                else:
                    st.warning("La ricerca web è stata richiesta, ma le API Key di Google non sono configurate.")
                    processed_web_search_results = "L'utente ha richiesto una ricerca web, ma le API non sono configurate."
            # ---- FINE NUOVA LOGICA PER RICERCA WEB ESPLICITA ----

            is_complex_generation_request = False # Logica esistente per comandi speciali
            generation_keywords = ["struttura", "analizza", "fammi", "costruiscimi", "scrivi", "genera", "crea una serie", "fammi un report"]
            tool_keywords_present = ("drive" in lower_q and "carica" in lower_q) or \
                                    (("sheet" in lower_q or "spreadsheet" in lower_q) and "crea" in lower_q)
            if len(user_q.split()) > 12 and any(g_kw in lower_q for g_kw in generation_keywords) and tool_keywords_present:
                is_complex_generation_request = True

            command_executed = False
            # ... (la logica dei comandi testuali rimane invariata qui, la ometto per brevità)
            if not is_complex_generation_request and "drive" in lower_q and "carica" in lower_q:
                command_executed = True
                assistant_content = ""
                if not st.session_state.get("google_creds"):
                    assistant_content = "Servizi Google non autorizzati."
                    st.session_state.action_in_progress = False 
                else:
                    file_to_upload_bytes = st.session_state.get("uploaded_file_bytes")
                    file_to_upload_name = st.session_state.get("uploaded_file_name")
                    file_to_upload_mime = st.session_state.get("uploaded_file_mime")
                    if file_to_upload_bytes and file_to_upload_name:
                        with st.spinner(f"Caricamento '{file_to_upload_name}'..."):
                            file_id = upload_file_to_drive(file_to_upload_bytes, file_to_upload_name, file_to_upload_mime)
                            if file_id: assistant_content = f"File allegato '{file_to_upload_name}' caricato: https://drive.google.com/file/d/{file_id}/view"
                            else: assistant_content = f"Errore caricamento '{file_to_upload_name}'."
                    elif cur_chat_data.get("messages") and len(cur_chat_data.get("messages")) >= 2 and cur_chat_data["messages"][-2].get("role") == "assistant": 
                        last_ai_content = cur_chat_data["messages"][-2]["content"]
                        text_bytes = last_ai_content.encode('utf-8'); text_name = f"risposta_ai_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
                        with st.spinner(f"Caricamento ultima risposta AI '{text_name}'..."):
                            file_id = upload_file_to_drive(text_bytes, text_name, "text/plain")
                            if file_id: assistant_content = f"Ultima risposta AI caricata come '{text_name}': https://drive.google.com/file/d/{file_id}/view"
                            else: assistant_content = f"Errore caricamento ultima risposta AI '{text_name}'."
                    else: assistant_content = "Nessun file da caricare."
                add_assistant_response_and_rerun_v2(assistant_content, cur_chat_data, mark_action_completed=True)
            
            elif not is_complex_generation_request and (("crea google sheet" in lower_q) or ("crea" in lower_q and ("sheet" in lower_q or "spreadsheet" in lower_q)) or ("converti in file google sheets" in lower_q)):
                command_executed = True
                assistant_content = ""
                if not st.session_state.get("google_creds"):
                    assistant_content = "Servizi Google non autorizzati."
                    st.session_state.action_in_progress = False
                else:
                    sheet_title_from_query = f"Foglio per: {user_q} (da SuperTutorAI)"
                    data_to_populate = None
                    populate_keywords = ["messaggio precedente", "tabella precedente", "sopra", "questi dati", "con la tabella", "usando la tabella", "dal messaggio prima", "procedi con la creazione dello sheet", "crea lo sheet con questi dati"]
                    should_populate_from_last_message = any(keyword in lower_q for keyword in populate_keywords)
                    if should_populate_from_last_message:
                        if cur_chat_data.get("messages") and len(cur_chat_data.get("messages")) >= 2:
                            reversed_messages = cur_chat_data["messages"][:-1][::-1]
                            last_ai_content = next((m.get("content") for m in reversed_messages if m.get("role") == "assistant"), None)
                            if last_ai_content:
                                parsed_table = parse_markdown_table_from_text(last_ai_content)
                                if parsed_table: data_to_populate = parsed_table; header_preview = ' '.join(map(str, parsed_table[0]))[:50]; sheet_title_from_query = f"Sheet: {header_preview}..."
                                else: st.warning("Provato a usare messaggio precedente, ma non è stato possibile estrarre tabella.")
                            else: st.warning("Nessun messaggio AI precedente trovato per popolare lo sheet.")
                        else: st.warning("Non abbastanza messaggi per usare il precedente.")
                    with st.spinner(f"Creazione Google Sheet '{sheet_title_from_query}'..."):
                        sheet_id = create_google_sheet(sheet_title_from_query)
                    if sheet_id:
                        if data_to_populate:
                            with st.spinner(f"Popolamento '{sheet_title_from_query}'..."):
                                ws = write_data_to_google_sheet(sheet_id, data_to_populate)
                                if ws: assistant_content = f"Foglio '{sheet_title_from_query}' creato e popolato: https://docs.google.com/spreadsheets/d/{sheet_id}/edit"
                                else: assistant_content = f"Foglio '{sheet_title_from_query}' creato (https://docs.google.com/spreadsheets/d/{sheet_id}/edit), ma errore nel popolarlo."
                        else: assistant_content = f"Foglio '{sheet_title_from_query}' creato: https://docs.google.com/spreadsheets/d/{sheet_id}/edit"
                        if should_populate_from_last_message and not data_to_populate: assistant_content += " (Non sono riuscito a estrarre dati tabellari dal messaggio precedente per popolarlo.)"
                    else: assistant_content = "Errore creazione Google Sheet."
                    add_assistant_response_and_rerun_v2(assistant_content, cur_chat_data, mark_action_completed=True)
            
            if not command_executed:
                if len(cur_chat_data["messages"]) == 1 and cur_chat_data.get("name","") == "Nuova Chat...":
                    with st.spinner("Generazione titolo..."):
                        cur_chat_data["name"] = generate_chat_title_from_query(user_q, st.session_state.gemini_api_key, st.session_state.selected_gemini_model)
                
                scraped_content_to_pass = None
                should_ignore_document_context_for_query = "staccandoci da questo contesto".lower() in user_q.lower()
                
                if cur_chat_data.get('active_web_page_context') and not should_ignore_document_context_for_query:
                    active_wp_info = cur_chat_data['active_web_page_context']
                    if active_wp_info.get('url'):
                        with st.spinner(f"Recupero contenuto da pagina precedentemente analizzata: {active_wp_info.get('title', '')[:50]}..."):
                            # Ri-esegui lo scraping per passare il contenuto completo alla chiamata LLM
                            scraped_content_to_pass = fetch_and_extract_text(active_wp_info['url'])
                            if not scraped_content_to_pass:
                                st.warning(f"Non è stato possibile ri-estrarre il contenuto da {active_wp_info['url']} per questa domanda. L'AI si baserà sulla cronologia e sul titolo/URL.")

                with st.spinner("Il tutor (Gemini) sta pensando..."):
                    emb_model = load_embedding_model_cached() # Caricato qui ora
                    rel_chunks_info_for_llm = []
                    direct_file_text_context_for_llm = None
                    image_data_for_llm = None
                    image_mime_for_llm = None
                    
                    phrase_to_ignore_context = "staccandoci da questo contesto"
                    should_ignore_document_context_for_query = phrase_to_ignore_context.lower() in user_q.lower()

                    if st.session_state.get("document_processed", False) and not should_ignore_document_context_for_query:
                        uploaded_mime = st.session_state.get("uploaded_file_mime")
                        uploaded_bytes = st.session_state.get("uploaded_file_bytes")
                        
                        if uploaded_mime and uploaded_bytes:
                            if uploaded_mime.startswith("image/"):
                                image_data_for_llm = uploaded_bytes
                                image_mime_for_llm = uploaded_mime
                                rel_chunks_info_for_llm = [] 
                                st.session_state.text_chunks = [] 
                                st.session_state.chunk_embeddings = np.array([])
                            else: # Non è un'immagine, potrebbe essere testo/PDF
                                text_code_mime_prefixes = ["text/"]
                                text_code_specific_mimes = [
                                    "application/json", "application/xml", "application/csv", 
                                    "application/rtf", "application/xhtml+xml", "application/javascript",
                                    "application/x-python-code", "application/x-csharp"
                                ]
                                text_code_extensions = [
                                    '.md', '.csv', '.json', '.xml', '.html', '.htm', '.rtf',
                                    '.py', '.js', '.java', '.c', '.cpp', '.cs', '.go', '.rb', 
                                    '.php', '.swift', '.kt', '.ts', '.sql', '.sh', '.bat', '.ps1'
                                ]
                                uploaded_file_name_lower = st.session_state.get("uploaded_file_name", "").lower()

                                is_potentially_direct_text_type = (
                                    any(uploaded_mime.startswith(p) for p in text_code_mime_prefixes) or
                                    uploaded_mime in text_code_specific_mimes or
                                    any(uploaded_file_name_lower.endswith(ext) for ext in text_code_extensions)
                                )

                                if is_potentially_direct_text_type and uploaded_mime != "application/pdf":
                                    try:
                                        decoded_text = uploaded_bytes.decode('utf-8', errors='replace')
                                        direct_file_text_context_for_llm = decoded_text
                                        st.session_state.text_chunks = [] 
                                        st.session_state.chunk_embeddings = np.array([])
                                        rel_chunks_info_for_llm = []
                                    except UnicodeDecodeError:
                                        st.warning(f"Impossibile decodificare il file {st.session_state.get('uploaded_file_name')} come UTF-8 per il contesto diretto. Verrà tentato il RAG se applicabile.")
                                
                                # Logica RAG (si applica a PDF, o se il file testuale non è stato messo in direct_file_text_context_for_llm causa errore decodifica)
                                if (
                                    not image_data_for_llm and 
                                    not direct_file_text_context_for_llm and
                                    st.session_state.get("text_chunks") and 
                                    isinstance(st.session_state.get("chunk_embeddings"), np.ndarray) and
                                    st.session_state.chunk_embeddings.size > 0 and 
                                    emb_model
                                ):
                                    q_embed = emb_model.encode([user_q])
                                    rel_chunks_info_for_llm = find_relevant_chunks(
                                        q_embed, 
                                        st.session_state.chunk_embeddings, 
                                        st.session_state.text_chunks, 
                                        st.session_state.top_n, 
                                        st.session_state.similarity_threshold
                                    )
                    
                    chat_hist_llm = [{"role":m["role"], "content":m["content"]} for m in cur_chat_data.get("messages",[])[:-1]]
                    
                    ai_res_text, p_tok, c_tok, citation_map = get_answer_from_llm(
                        query=user_q, 
                        relevant_chunks_data=rel_chunks_info_for_llm, 
                        api_key=st.session_state.gemini_api_key, 
                        model_name=st.session_state.selected_gemini_model, 
                        chat_history=chat_hist_llm,
                        ignore_active_contexts_for_this_query=should_ignore_document_context_for_query, 
                        direct_file_text_context=direct_file_text_context_for_llm, 
                        image_data=image_data_for_llm,            
                        image_mime_type=image_mime_for_llm,
                        external_web_search_results=processed_web_search_results,
                        scraped_page_content_for_query=scraped_content_to_pass,
                        temporary_instructions_for_this_turn=temp_instructions_for_llm,
                        retrieved_history_context=retrieved_history_context_for_llm 
                    )
                    
                    if should_ignore_document_context_for_query and not temp_instructions_for_llm and not retrieved_history_search_active:
                        ai_res_text += "\n\n*(Nota: Per questa risposta, il contesto da eventuali documenti o articoli attivi è stato temporaneamente ignorato come da tua richiesta.)*"

                    add_assistant_response_and_rerun_v2(
                        assistant_content=ai_res_text, current_chat_data=cur_chat_data,
                        prompt_tokens=p_tok, completion_tokens=c_tok,
                        citation_map=citation_map, chunks_info=rel_chunks_info_for_llm,
                        mark_action_completed=True
                    )
        # else: st.caption("Azione precedente in corso o input utente già processato...")

